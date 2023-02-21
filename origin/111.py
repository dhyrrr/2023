from pptx import Presentation
from pptx.shapes.picture import Picture
import os
import re
import shutil
import collections.abc
from pptx import Presentation, util
from pptx.util import Cm,Pt,Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

######将ppt的图片保存，并按照id重命名

image_name_list=[]
image_path=r'C:\Users\Administrator.DESKTOP-8279PGK\Desktop\小试\1'
file_path = r'C:\Users\Administrator.DESKTOP-8279PGK\Desktop\小试\3.pptx'
prs = Presentation(file_path)
index=0
index_list=[]
#读取幻灯片的每一页
for i,slide in enumerate(prs.slides):
    # 读取每一板块
    sortedShapes = sorted(slide.shapes, key=lambda x: (x.left))
    for shape in sortedShapes:
        # print(dir(shape))
        ######对每一页每一块进行分析
        #######如果含有PPTshape含有表格
        if isinstance(shape, Picture):
            if i==4:
                index+=1
                print(index)
                index_list.append(index)
        print(index_list)
        # print(max(index_list))

        if shape.has_table:
            one_table_data = []
            for row in shape.table.rows:  # 读每行
                row_data = []
                for cell in row.cells:  # 读一行中的所有单元格
                    c = cell.text
                    row_data.append(c)
                one_table_data.append(row_data)
            if one_table_data[0][0]=='Biortus code':
                # print(one_table_data[1][0])
                try:
                    plasmidno=one_table_data[1][0].split("#")[1]

                except:
                    plasmidno=one_table_data[1][0]
                # print(plasmidno)
            try:
                if one_table_data[6][0]=='Score':
                    # print(one_table_data)
                    plasmidno_score_list=[]
                    score=one_table_data[6][1]
                    # plasmidno=one_table_data[0][1].split('#')[1]
                    plasmidno_score=plasmidno+'_'+score
                    # print(plasmidno_score)
                    # image_name_list=[]
                    for i in range(2):
                        image_name=plasmidno_score+"_"+str(i)
                        image_name_list.append(image_name)
            except:
                pass
        # if isinstance(shape, Picture):
        #     with open(os.path.join(image_path,f'{index}.jpg'), 'wb') as f:
        #         f.write(shape.image.blob)
        #         index+=1

print(image_name_list)

# ####重命名
# for roots,dirs,files in os.walk(image_path):
#     for file in files:
#         if re.search("jpg$",file):
#             # print(file)
#             file=file.split(".")[0]
#             # print(image_name_list)
#             try:
#                 file_name=image_name_list[int(file)]
#                 print(file,file_name)
#                 oldname = image_path + os.sep + file + '.jpg'
#                 newname = image_path + os.sep + file_name+ '.jpg'
#                 os.rename(oldname, newname)
#             except:
#                 pass


######加载图片


########向ppt写入图片和文本框
'''
template_ppt_path = r'C:\Administrator.DESKTOP-8279PGK\Desktop/3\ppt\FSEC_Analysis_2023-02-06 16-37-39.pptx'
ppt_file = Presentation(template_ppt_path)
slice = ppt_file.slides.add_slide(ppt_file.slide_layouts[1])
img_path = 'F:\image/27244_7_0.jpg' # 图片路径
#插入图片
left = Cm(17.15)
top = Cm(6.88)
width = Cm(2.88)
height = Cm(5.1)

slice.shapes.add_picture(img_path,left,top,width=width,height=height)
# 添加文本框

textbox= slice.shapes.add_textbox(left=Cm(16.93),
                                    top=Cm(5),
                                    width=Cm(4.02),
                                    height=Cm(0.94)
                                   )
## 向文本框加入文字
tf = textbox.text_frame
para = tf.add_paragraph()    # 添加段落
para.text = 'SDS-PAGE'
para.alignment = PP_ALIGN.CENTER  # 居中

font = para.font
font.size = Pt(12)    # 大小
font.name = 'Times New Roman'    # 字体
font.bold = False    # 加粗
font.italic = False  # 倾斜
font.color.rgb = RGBColor(0, 0, 0)
ppt_file.save('2.pptx')
'''
#######根据plasmidno匹配图片，只匹配原ppt右边的，如果含有多个score，则取最高score的那个
# plasmidno_list = ['17792','17793']
# file_name_list=[]
# name_list=[]
# image_path = r'C:\Users\Administrator.DESKTOP-8279PGK\Desktop\小试\1'
# for roots,dirs,files in os.walk(image_path):
#     for file in files:
#         if re.search("jpg$",file):
#             file_name = file.split(".")[0]
#             # print(file_name)
#             file_name_list.append(file_name)
# ###########
# ########选择score最高的image
# for plasmidno in plasmidno_list:
#     print(plasmidno)
#     code_list = []
#     max_list = []
#     for file_name in file_name_list:
#         if plasmidno == file_name.split("_")[0]:
#             # print(i)
#             code_list.append(file_name)
#             try:
#                 s = int(file_name.split("_")[1])
#                 # print(s)
#                 max_list.append(s)
#             except:
#                 s = ''
#     for i in code_list:
#         try:
#             if i.split("_")[1] == str(max(max_list)):
#                 if i.split("_")[2] == str(1):
#                     # print(i)
#                     name = image_path + os.sep + i + '.jpg'
#                     print(name)
#         except:
#             if i.split("_")[2] == str(1):
#                 name = image_path + os.sep + i + '.jpg'
#                 print(name)
#         # print(name)
#
#     name_list.append(name)
#     print(name_list)
    # print(name)
    # print(name_list)
    # print(code_list)
    # print(code_list)
    # for i in code_list:
    #     if i.split("_")[1] == str(max(max_list)):
    #         if i.split("_")[2] == str(0):
    #             # print(i)
    #             name = image_path + os.sep + i + '.jpg'
    # print(name)
    # name_list.append(name)
# print(name_list)





#######获取pptx占位符
####邵博
# 查看占位符的序号
# for placeholder in slice.placeholders:
#     info = placeholder.placeholder_format
#     print("索引{0},名称{1},类型{2},文本{3}".format(info.idx,placeholder.name,info.type,placeholder.text))
# from pptx import Presentation
#
# prs = Presentation("F:\PycharmProjects\origin\For DR\Template.pptx")
# slide = prs.slides.add_slide(prs.slide_layouts[2])  # 用母版生成一页ppt
# for shape in slide.placeholders:         # 获取这一页所有的占位符
#     phf = shape.placeholder_format
#     print(f'{phf.idx}--{shape.name}--{phf.type}')  # id号--占位符形状名称-占位符的类型











