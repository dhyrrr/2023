import os
import re
import sys
import psutil
import pythoncom
from pandas import DataFrame
import originpro as op
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from win32com.client import Dispatch
from scipy.signal import find_peaks
import time
import pptx
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.oxml.xmlchemy import OxmlElement
from szpscript.Public_Script import PublicScript
from szpscript.Xlsx_script import Xlxs
from gooey import Gooey, GooeyParser
import codecs
import json
from pptx import Presentation
from pptx.shapes.picture import Picture
import shutil
import collections.abc
from pptx import Presentation, util
from pptx.util import Cm,Pt,Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

from pptx.enum.text import PP_ALIGN


class Fsec:
    def __init__(self, data_file_path, pptx_template_file_path, complex_protein_info_xlsx, writer, standard_curve_dict, peak_x_error, peak_prominence, Ex_Wavelength, save_dir,Free_GFP,ppt_file
                 ):
        self.data_file_path = data_file_path
        self.pptx_template_file_path = pptx_template_file_path
        self.complex_protein_info_xlsx = complex_protein_info_xlsx
        self.writer = writer
        self.standard_curve_dict = standard_curve_dict
        self.peak_x_error = peak_x_error
        self.peak_prominence = peak_prominence
        self.Ex_Wavelength = float(Ex_Wavelength)
        self.save_dir = save_dir
        self.Free_GFP = Free_GFP
        # self.image_file_path=image_file_path
        self.ppt_file=ppt_file
        self.data_dict = {}
        self.ppt_opju_name_list = []
        self.font_dict = {"Times New Roman": 410,
                          "Courier New": 181}
        self.complex_protein_dict = {}

    @staticmethod
    def load_configuration():
        if os.path.exists("configuration.json"):
            with open("configuration.json", "r", encoding="utf-8") as f:
                configuration_dict = json.load(f)

            try:
                data_file_path = configuration_dict["data_file_path"]
            except:
                data_file_path = ""

            try:
                pptx_template_file_path = configuration_dict["pptx_template_file_path"]
            except:
                pptx_template_file_path = ""

            try:
                complex_protein_info_xlsx = configuration_dict["complex_protein_info_xlsx"]
            except:
                complex_protein_info_xlsx = ""

            try:
                writer = configuration_dict["writer"]
            except:
                writer = ""

            try:
                a, b = configuration_dict["standard_curve_dict"]["a"], configuration_dict["standard_curve_dict"]["b"]
            except:
                a, b = "", ""

            try:
                peak_x_error = configuration_dict["peak_x_error"]
            except:
                peak_x_error = ""

            try:
                peak_prominence = configuration_dict["peak_prominence"]
            except:
                peak_prominence = ""

            try:
                Ex_Wavelength = configuration_dict["Ex_Wavelength"]
            except:
                Ex_Wavelength = ""

            try:
                save_dir = configuration_dict["save_dir"]
            except:
                save_dir = ""
            try:
                Free_GFP = configuration_dict["Free_GFP "]
            except:
                Free_GFP  = ""
            try:
                ppt_file = configuration_dict["ppt_file"]
            except:
                ppt_file = ''
            # try:
            #     image_file_path = configuration_dict["image_file_path"]
            # except:
            #     image_file_path = ''


        else:
            data_file_path = ""
            pptx_template_file_path = ""
            complex_protein_info_xlsx = ""
            writer = ""
            a, b = "", ""
            peak_x_error = ""
            peak_prominence = ""
            Ex_Wavelength = ""
            save_dir = ""
            Free_GFP=''
            ppt_file=''
            # image_file_path=''

        return data_file_path, pptx_template_file_path, complex_protein_info_xlsx, writer, a, b, peak_x_error, peak_prominence, Ex_Wavelength, save_dir,Free_GFP,ppt_file

    def process_maximum_value(self, input_maximum_value_list):
        output_maximum_value_list = []

        # 按x轴大小从小到大排序
        input_maximum_value_list.sort(key=lambda x:x[0])

        # 合并x轴大小相差小于10%的数据
        temp_x, temp_y = None, None
        for input_maximum_value in input_maximum_value_list:
            if temp_x != None and temp_y != None:
                # print(input_maximum_value[0], input_maximum_value[1], type(input_maximum_value[0]), type(input_maximum_value[1]), flush=True)
                if abs((input_maximum_value[0] - temp_x) / temp_x) > self.peak_x_error:
                    output_maximum_value_list.append([temp_x, temp_y])  # 数据相差大于10%, 确定前一数据的值
                    temp_x, temp_y = input_maximum_value[0], input_maximum_value[1]  # 更新数据
                else:
                    if input_maximum_value[1] > temp_y:
                        temp_x, temp_y = input_maximum_value[0], input_maximum_value[1]  # 取大值

            else:
                temp_x, temp_y = input_maximum_value[0], input_maximum_value[1]

        output_maximum_value_list.append([temp_x, temp_y])  # 更新最后的数据

        return output_maximum_value_list

    def get_complex_protein_info(self):
        temp = Xlxs.get_info_with_head(self.complex_protein_info_xlsx)
        for index, info in temp.items():
            plasmidno = info["Plasmid No"]
            self.complex_protein_dict[plasmidno] = {}
            for key, value in info.items():
                self.complex_protein_dict[plasmidno][key] = value

    def get_fsec_data(self):
        for roots, dirs, files in os.walk(self.data_file_path):
            txt_path_dict = {}

            for file in files:
                if re.search("\.txt$", file):
                    txt_path_dict[file.replace(".txt", "")] = os.path.join(roots, file)

            # 数据存在
            if txt_path_dict:
                txt_path_dict_sorted = sorted(txt_path_dict.items(), key=lambda x: str(x[0].split("_")[0]))

                # 写时间轴
                for filename_filepath in txt_path_dict_sorted:
                    fileid = filename_filepath[0].split("_")[0]
                    if self.data_dict.get(fileid, None):
                        print("*" * 20 + "提示" + "*" * 20)
                        print("!" * 45)
                        print("文件{}的命名其其他文件命名重合, 该文件数据不会被处理".format(filename_filepath[0]))
                        print("!" * 45)
                        print("!" * 45)

                        continue

                    else:
                        self.data_dict[fileid] = {}

                    temp_fsec_data_dict = {}
                    with open(filename_filepath[1], "r", encoding="GBK") as f:
                        for line in f:
                            if "Ex. Wavelength(nm)" in line:
                                # 载入当前波长的数据
                                data_ex_wavelength = float(re.search("\d+", line).group())
                                temp_fsec_data_dict[data_ex_wavelength] = {}

                                for line in f:
                                    if not line.strip():
                                        break

                                    if re.search("^\d", line):
                                        time = re.split("\s+", line.strip())[0]
                                        data = float(re.split("\s+", line.strip())[1]) / 1000
                                        temp_fsec_data_dict[data_ex_wavelength][time] = data

                    # 分析符合激光波长的数据
                    right_data_ex_wavelength = None
                    for data_ex_wavelength in temp_fsec_data_dict.keys():
                        if right_data_ex_wavelength == None:
                            right_data_ex_wavelength = data_ex_wavelength
                        else:
                            if abs(data_ex_wavelength - self.Ex_Wavelength) < abs(right_data_ex_wavelength - self.Ex_Wavelength):
                                right_data_ex_wavelength = data_ex_wavelength

                    # 保存符合激光波长的数据
                    for time, data in temp_fsec_data_dict[right_data_ex_wavelength].items():
                        self.data_dict[fileid][time] = data

    def analysis_fsec_data(self):
        # print(self.data_dict, flush=True)
        data_df = DataFrame(self.data_dict)

        column_name_dict = {}

        for column_name in data_df.columns.values:
            try:
                plasmid_id = re.search("^\d+", column_name.strip()).group()
            except:
                print("*" * 20 + "提示" + "*" * 20)
                print("!" * 45)
                print("文件{}命名不是以质粒编号开头, 该文件数据不会被处理".format(column_name))
                print("!" * 45)
                print("!" * 45)

                continue

            column_df = data_df[column_name]

            max_value_y_coordinate = float(column_df.max())

            # 获取数据极值点Index
            coordinate_x, coordinate_y = column_df.index, column_df.values
            maximun_value_index_list, _ = find_peaks(column_df, prominence=self.peak_prominence)
            # print(_)

            if column_name_dict.get(plasmid_id, None):
                temp_maximun_value_list = []
                for maximun_value_index in maximun_value_index_list:
                    x, y = float(coordinate_x[maximun_value_index]), float(coordinate_y[maximun_value_index])
                    temp_maximun_value_list.append([x, y])

                column_name_dict[plasmid_id]["column_name_list"].append(column_name)
                column_name_dict[plasmid_id]["maximun_value_list"].extend(temp_maximun_value_list)

                if max_value_y_coordinate > column_name_dict[plasmid_id]["maximun_y_value"]:
                    # print(max_value_y_coordinate, column_name)
                    column_name_dict[plasmid_id]["maximun_y_value"] = max_value_y_coordinate
                    column_name_dict[plasmid_id]["maximum_peak_column_name"] = column_name

            else:
                temp_maximun_value_list = []
                for maximun_value_index in maximun_value_index_list:
                    x, y = float(coordinate_x[maximun_value_index]), float(coordinate_y[maximun_value_index])
                    temp_maximun_value_list.append([x, y])

                column_name_dict[plasmid_id] = {"column_name_list": [column_name],
                                                "maximun_y_value": max_value_y_coordinate,
                                                "maximum_peak_column_name": column_name,
                                                "maximun_value_list": temp_maximun_value_list
                                                }


        for plasmid_id, info_list in column_name_dict.items():
            print("    处理质粒编号: '{}'数据".format(plasmid_id), flush=True)
            column_name_list = info_list["column_name_list"]
            column_df = data_df[column_name_list]

            maximum_peak_column_name = info_list["maximum_peak_column_name"]
            # maxinum_peak_column_df = data_df[maximum_peak_column_name]

            # 获取数据极值点Index
            # maximun_value_index_list, _ = find_peaks(maxinum_peak_column_df, prominence=1)
            maximun_value_list = info_list["maximun_value_list"]
            # print(maximun_value_list)
            merged_maximun_value_list = self.process_maximum_value(maximun_value_list)
            # print(merged_maximun_value_list)

            # self.origin_plot(column_df, column_name, maximun_value_index_list)
            self.ppt_opju_name_list.append(maximum_peak_column_name)
            # self.origin_plot(column_df, maxinum_peak_column_df, maximum_peak_column_name, maximun_value_index_list)
            self.origin_plot(column_df, maximum_peak_column_name, merged_maximun_value_list)
        
    def origin_plot(self, column_df, maximum_peak_column_name, merged_maximun_value_list):
        pythoncom.CoInitialize()

        def origin_shutdown_exception_hook(exctype, value, traceback):
            op.exit()
            sys.__excepthook__(exctype, value, traceback)

        if op and op.oext:
            sys.excepthook = origin_shutdown_exception_hook

        # Set Origin instance visibility.
        if op.oext:
            op.set_show(False)

        # import data
        wks = op.new_sheet('w')
        wks.from_list(0, column_df.index)   # 自变量
        wks.from_df(column_df, 1)   # 因变量
        # wks.set_int("font", 410)
        # wks.set_int("fsize", 35)

        # create layer
        gp = op.new_graph()
        gl = gp[0]

        # plot
        legend_list = []
        for data_index in range(len(column_df.columns)):
            plt = gl.add_plot(wks, data_index + 1, 0, -1, 'l')
            plt.set_int('line.width', 3)
            plt.color = data_index + 1
            legend_list.append("\l({0}) %({0})".format(data_index + 1))

        # Legend
        lgnd = gl.label('Legend')
        lgnd.text = column_df.columns.values[data_index]
        lgnd.text = "\n".join(legend_list)
        # "page -FLS -m 0.01"

        # axis
        gl.rescale()

        axisx = gl.axis('x')
        axisy = gl.axis('y')

        axisx.set_limits(0, 15)

        axisx.title = "Time (min)"
        axisy.title = "Fluorescence (mV)"

        # generate ogs
        y_max_value = gl.ylim[1]

        # 生成ogs
        self.generate_ogs(maximum_peak_column_name, merged_maximun_value_list, y_max_value)

        latalk_string = 'run.SECTION({0}, Main)'.format(os.path.join(self.save_dir, "ogs", maximum_peak_column_name))
        gl.lt_exec(latalk_string)

        # 循环多次, 源程序Origin单次执行存在Bug
        for i in range(10):
            latalk_string = "page -FLS -m 0.01"
            gl.lt_exec(latalk_string)

        fpath = os.path.join(self.save_dir, "opju", "{}.opju".format(maximum_peak_column_name))
        op.save(fpath)
        # print(f'{gl} is exported as {fpath}')

        # if op.oext:
        op.exit()

        pythoncom.CoUninitialize()

    def generate_ogs(self, column_name, maximun_value_index_list, y_max_value):
        with open(os.path.join(self.save_dir, "ogs", "{}.ogs".format(column_name)), "w", encoding="utf-8") as f:
            f.write('''[Main]\n
//uncomment following line to define functions using outside of this scope\n
//@global=1;\n
//type -b "Hello LabTalk!";\n''')


            for num, xy_value_list in enumerate(maximun_value_index_list):
                if xy_value_list[0] and xy_value_list[1]: # 值可能为None
                    x = float(xy_value_list[0])
                    y = float(xy_value_list[1])
                    if x >= 1:
                        f.write("GObject myline{0} = Arrow{0};\n".format(num))
                        # print(max_value_x_coordinate, max_value_y_coordinate, max_value_y_coordinate * 1.2)
                        f.write("draw -n myline{0} ".format(num) + "-l {" + "{0},{1},{0},{2}".format(x, y + y_max_value * 0.01, y + y_max_value * 0.1) + "};\n")
                        f.write("myline{0}.ARROWBEGINSHAPE = 2;\n".format(num))
                        f.write("myline{0}.COLOR =2;\n".format(num))

                    # protein label
                    if x < 1:
                        pass
                    elif 1 <= x <= 5:
                        protein_label = "void"
                        f.write("label -a {0} {1} \c2({2});\n".format(x - 0.5, y + y_max_value * 0.17, protein_label))
                    elif 8 <= x <= 9:
                        protein_label = self.Free_GFP
                        f.write("label -a {0} {1} \c2({2});\n".format(x - 1, y + y_max_value * 0.17, protein_label))
                    elif x >= 10:
                        pass
                    else:
                        protein_label = "~ {} KD".format(int(10 ** ((x - 0.7) * self.standard_curve_dict["a"] + self.standard_curve_dict["b"])))
                        f.write("label -a {0} {1} \c2({2});\n".format(x - 1, y + y_max_value * 0.17, protein_label))
                
    def generate_image(self):
        ######将小试ppt图片保存
        # image_file_path = self.image_file_path
        image_file_path=os.path.join(self.save_dir, "image")
        for roots, dirs, files in os.walk(image_file_path):
            for file in files:
                    os.remove(os.path.join(image_file_path, file))
        ppt_file = self.ppt_file
        ppt_dic={}
        for roots, dirs, files in os.walk(ppt_file):
            for file in files:
                if re.search("pptx$", file):
                    ppt_path = os.path.join(ppt_file, file)
                    # print(ppt_path)
                    prs = Presentation(ppt_path)
                    index = 0
                    # 读取幻灯片的每一页
                    slides = prs.slides
                    plasmidno_score_dic = {}
                    plasmidno_dic = {}
                    info_dic = {}
                    list = []
                    list1 = []
                    for i, slide in enumerate(prs.slides):
                        sortedShapes = sorted(slide.shapes, key=lambda x: (x.left))
                        for shape in sortedShapes:
                            if shape.has_table:
                                one_table_data = []
                                for row in shape.table.rows:  # 读每行
                                    row_data = []
                                    for cell in row.cells:  # 读一行中的所有单元格
                                        c = cell.text
                                        row_data.append(c)
                                    one_table_data.append(row_data)
                                if one_table_data[0][0] == 'Biortus code':
                                    # print(one_table_data[1][0])
                                    try:
                                        plasmidno = one_table_data[1][0].split("#")[1].strip()

                                    except:
                                        plasmidno = one_table_data[1][0].strip()
                                    # print(plasmidno)
                                try:
                                    info_list = []
                                    info_plasmidno_list = []
                                    Competent_cell_list = []
                                    plasmidno_score_list = []
                                    IPTG_list = []
                                    Temperature_list = []
                                    Induction_time_list = []
                                    if one_table_data[6][0] == 'Score':
                                        # print(one_table_data)
                                        score = one_table_data[6][1]
                                        if score == '\\':
                                            score = ''
                                        else:
                                            score = score
                                        plasmidno_score = plasmidno + '_' + score
                                        plasmidno_score_list.append(plasmidno_score)
                                    if one_table_data[1][0] == 'Competent cell':
                                        Competent_cell = one_table_data[1][1]
                                        Competent_cell_list.append(Competent_cell)
                                    if one_table_data[3][0] == 'IPTG':
                                        IPTG = one_table_data[3][1]
                                        IPTG_list.append(IPTG)
                                    if one_table_data[4][0] == 'Temperature':
                                        Temperature = one_table_data[4][1]
                                        Temperature_list.append(Temperature)
                                    if one_table_data[5][0] == 'Induction time':
                                        Induction_time = one_table_data[5][1]
                                        Induction_time_list.append(Induction_time)
                                    info_plasmidno_list = plasmidno_score_list + Competent_cell_list + IPTG_list + Temperature_list + Induction_time_list
                                    info_list = Competent_cell_list + IPTG_list + Temperature_list + Induction_time_list
                                    plasmidno_score_dic[i] = info_plasmidno_list
                                    list = plasmidno_score_list
                                    list.append(info_list)
                                    list1.append(list)
                                    info_dic[plasmidno_score] = info_list
                                except:
                                    pass
                        ######获取每一页对应的质粒编号
                    # print(plasmidno_score_dic)
                    # print(info_dic)
                    ########根据质粒编号遍历图片并保存
                    for i, slide in enumerate(prs.slides):
                        sortedShapes = sorted(slide.shapes, key=lambda x: (x.left))
                        shape_list = []
                        ########找到最右边的那个图
                        for shape in sortedShapes:
                            if isinstance(shape, Picture):
                                right = shape.left
                                # print(shape, right)
                                shape_list.append(shape)
                        try:
                            rightmost_shape = max(shape_list, key=lambda shape: shape.left + shape.width)
                            # print(rightmost_shape)
                            try:
                                index += 1
                                name = file + "_" + plasmidno_score_dic[i][0] + "_" + str(index)+".jpg"
                                info = plasmidno_score_dic[i]
                                #######判断文件后缀类型
                                with open(os.path.join(image_file_path, f'{name}'), 'wb') as f:
                                    f.write(rightmost_shape.image.blob)
                            except:
                                pass
                            ppt_dic[name] = info
                        except:
                            pass
            # print(ppt_dic)
            with open("ppt_dic.json", "w", encoding="utf-8") as f:
                json.dump(ppt_dic, f, indent=4, ensure_ascii=False)
            return ppt_dic

    def choose_image(self):
        image_file_path = os.path.join(self.save_dir, "image")
        file_name_list = []
        image_path_list=[]
        plasmidno_list = []
        dic={}
        # print(self.ppt_opju_name_list)
        for column_name in self.ppt_opju_name_list:
            plasmidno = re.search("^\d+", column_name).group()
            plasmidno_list.append(plasmidno)
        for roots, dirs, files in os.walk(image_file_path):
            for file in files:
                if re.search("jpg$", file):
                    # file_name = file.rpartition(".")[0]
                    file_name_list.append(file)
        for plasmidno in plasmidno_list:
            # print(plasmidno)
            code_list = []
            max_list = []
            max_u_list = []
            # print(plasmidno)
            for file_name in file_name_list:
                if plasmidno == file_name.split("_")[1]:
                    # print(i)
                    code_list.append(file_name)
                # print(code_list)
            dic[plasmidno] = code_list
        # print(dic)
        return dic

    def generate_ppt(self):

        file = Presentation(self.pptx_template_file_path)
        ## ppt cover
        eln_no_list = []
        for row_index, column_name in enumerate(self.ppt_opju_name_list):
            plasmidno = re.search("^\d+", column_name).group()
            try:
                eln_no = self.complex_protein_dict[plasmidno]["ELN NO."]
            except:
                continue

            if not eln_no in eln_no_list:
                eln_no_list.append(eln_no)

        slice_cover = file.slides.add_slide(file.slide_layouts[0])

        cover_title1 = slice_cover.placeholders[17]
        cover_title2 = slice_cover.placeholders[14]
        cover_title3 = slice_cover.placeholders[18]

        cover_title1.text = "项目编号: {}".format(", ".join(eln_no_list))
        cover_title2.text = "FSEC Test Report"
        cover_title3.text = "{}\n".format(self.writer) + time.strftime("%Y%m%d", time.localtime())

        ## info
        slice_info = file.slides.add_slide(file.slide_layouts[1])

        info_title1 = slice_info.placeholders[17]
        info_title2 = slice_info.placeholders[19]

        info_title1.text = "FSEC Test Scheme"
        # info_title2.text = "Objective:"

        info_text_frame = info_title2.text_frame
        info_paragraph = info_text_frame.add_paragraph()

        # objective
        plasmidno_list = []
        for row_index, column_name in enumerate(self.ppt_opju_name_list):
            plasmidno = re.search("^\d+", column_name).group()
            plasmidno_list.append(plasmidno)

        objective_name = info_text_frame.paragraphs[0].add_run()
        objective_name.text = "Objective: "
        objective_name.font.bold = True

        objective_info = info_text_frame.paragraphs[0].add_run()
        objective_info.text = "Test {}".format(", ".join(plasmidno_list))

        # Sample information
        sample_name = info_paragraph.add_run()
        sample_name.text = "Sample information: "
        sample_name.font.bold = True

        sample_info = info_paragraph.add_run()
        sample_info.text = "\n"

        # Table
        info_table_row_num = len(self.ppt_opju_name_list) + 1
        x, y, cx, cy = Cm(0.71), Cm(4.27), Cm(23.98), Cm(5.05 * info_table_row_num / 7)
        info_table = slice_info.shapes.add_table(info_table_row_num, 6, x, y, cx, cy).table

        # 设置宽度
        info_table.columns[0].width = Cm(2.5)
        info_table.columns[1].width = Cm(2.5)
        info_table.columns[2].width = Cm(9.48)
        info_table.columns[3].width = Cm(3)
        info_table.columns[4].width = Cm(3.5)
        info_table.columns[5].width = Cm(3)

        #table text
        title_list = ["From", "No.", "Protein", "Cell", "Plasmid/PEI", "HPI"]
        for title_index, title in enumerate(title_list):
            info_table.cell(0, title_index).text = title

        for row_index, column_name in enumerate(self.ppt_opju_name_list):
            plasmidno = re.search("^\d+", column_name).group()


            try:
                protein = self.complex_protein_dict[plasmidno]["Plasmid name"]
            except:
                protein = ""

            info_table.cell(row_index + 1, 1).text = plasmidno
            info_table.cell(row_index + 1, 2).text = protein
            info_table.cell(row_index + 1, 2).margin_bottom = Cm(0)
            info_table.cell(row_index + 1, 2).margin_top = Cm(0)

        # table format
        for cell in info_table.iter_cells():
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                # paragraph.line_spacing = 1
                paragraph.font.name = "Times New Roman"
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                # paragraph.space_after = Cm(0)
                # paragraph.space_before = Cm(0)
                # paragraph.margin_bottom = Cm(0)
                # paragraph.margin_top = Cm(0)

        for column_index in range(6):
            info_table.cell(0, column_index).text_frame.paragraphs[0].font.size = Pt(11)

        # blank
        blank_info = info_paragraph.add_run()
        if info_table_row_num <= 5:
            blank_info.text = "\n" * 7
        else:
            blank_info.text = "\n" * int(8 * info_table_row_num / 7)

        # Procedures
        procedure_name = info_paragraph.add_run()
        procedure_name.text = "Procedures: "
        procedure_name.font.bold = True

        procedure_info = info_paragraph.add_run()
        procedure_info.text = ""

        ## opju
        data_df = DataFrame(self.data_dict)
        # print(len(self.data_dict.keys()))
        # print(len(data_df))

        for page_index, column_name in enumerate(self.ppt_opju_name_list):

            plasmidno = re.search("^\d+", column_name).group()

            slice_opju = file.slides.add_slide(file.slide_layouts[2])

            opju_title1 = slice_opju.placeholders[17]
            opju_title2 = slice_opju.placeholders[19]
            opju_title3 = slice_opju.placeholders[20]
            opju_table1 = slice_opju.placeholders[18]
            # opju_pic1 = slice_opju.placeholders[21]
            # opju_title4 = slice_opju.placeholders[22]

            # Fsec results
            opju_title1.text = "FSEC Results"

            # info
            info_text_frame = opju_title2.text_frame
            info_paragraph = info_text_frame.add_paragraph()

            # Sample
            buffer_name = info_text_frame.paragraphs[0].add_run()
            buffer_name.text = "Sample: "
            buffer_name.font.bold = True

            buffer_info = info_text_frame.paragraphs[0].add_run()
            buffer_info.text = "lysate"

            # Running buffer
            buffer_name = info_paragraph.add_run()
            buffer_name.text = "Running buffer: "
            buffer_name.font.bold = True

            ######buffer条件修改
            buffer_info = info_paragraph.add_run()
            buffer_info.text = "20 mM HEPES pH 7.5, 150 mM NaCl.\n"

            # Injection volume
            volume_name = info_paragraph.add_run()
            volume_name.text = "Injection volume: "
            volume_name.font.bold = True

            volume_info = info_paragraph.add_run()
            volume_info.text = "25 μL\n"

            # Column
            column_name = info_paragraph.add_run()
            column_name.text = "Column: "
            column_name.font.bold = True

            column_info = info_paragraph.add_run()
            column_info.text = "Superdex 200 Increase 5/150, 3 mL (ID0028); \n"

            # Ex
            ex_name = info_paragraph.add_run()
            ex_name.text = "Ex: "
            ex_name.font.bold = True

            ex_info = info_paragraph.add_run()
            ex_info.text = "498 nm  "

            em_name = info_paragraph.add_run()
            em_name.text = "Em: "
            em_name.font.bold = True

            em_info = info_paragraph.add_run()
            em_info.text = "526 nm"

            # table
            opju_table1_opj = opju_table1.insert_table(2, 3).table
            for r in range(2):
                if r == 0:
                    position_list = ['a:lnT']
                elif r == 1:
                    position_list = ['a:lnB']
                for c in range(3):
                    cell = opju_table1_opj.cell(r, c)
                    cell = self._set_cell_border(cell, position_list=position_list)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = pptx.dml.color.RGBColor(255, 255, 255)

            opju_table1_opj.columns[0].width = Cm(2)
            opju_table1_opj.columns[1].width = Cm(5.16)
            opju_table1_opj.columns[2].width = Cm(3)

            # table text
            title_list = ["No.", "Protein", "Theoretical MW (Da)"]
            for title_index, title in enumerate(title_list):
                opju_table1_opj.cell(0, title_index).text = title

            opju_table1_opj.cell(1, 0).text = plasmidno

            try:
                protein = self.complex_protein_dict[plasmidno]["Plasmid name"]

            except:
                protein = ""
            opju_table1_opj.cell(1, 1).text = protein

            try:
                molecularweight = str(self.complex_protein_dict[plasmidno]["Theoretical molecular weight (Da)"])
            except:
                molecularweight = ""
            opju_table1_opj.cell(1, 2).text = molecularweight

            # table format
            for cell in opju_table1_opj.iter_cells():
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    # paragraph.line_spacing = 1
                    paragraph.font.name = "Times New Roman"
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.font.color.rgb = RGBColor(0, 0, 0)

            #####添加图片和图标
            # opju_title4.text = "SDS-PAGE"
            # pic
            # left = Cm(1.9)
            # top = Cm(3.9)
            # # height = Cm(15)
            # width = Cm(23)
            dic=self.choose_image()
            info=self.generate_image()
            # print(info)
            image_path_list=dic[plasmidno]
            print(plasmidno,image_path_list)
            # for i in image_path:
            #     info=pic1[i]
            #     print(info)
            #######打印图片路径
            path = os.path.join(self.save_dir, "image")
            if len(image_path_list)==0:
                print(plasmidno + '_小式ppt中不含该质粒胶图')
            # print(plasmidno+"_胶图路径："+image_path)
            if len(image_path_list) == 1:
                # 循环添加图片
                pic_width = Cm(2.88)
                pic_aspect_ratio = 2.88 / 5.1
                pic_height = pic_width / pic_aspect_ratio
                image_path = os.path.join(path, image_path_list[0])
                    # 计算当前图片的左侧和顶部位置
                pic_left = Cm(17.15)
                pic_top=Cm(6.88)
                # 添加图片到幻灯片
                slice_opju.shapes.add_picture(image_path, pic_left, pic_top, pic_width, pic_height)
                txBox_left = pic_left-Cm(0.5)
                txBox_top = pic_top + pic_height + Cm(0.2)
                txBox_width = pic_width
                txBox_height = Cm(3)
                txBox = slice_opju.shapes.add_textbox(txBox_left, txBox_top, txBox_width, txBox_height)
                tf = txBox.text_frame
                # tf.vertical_anchor = MSO_ANCHOR.TOP_LEFT
                tf.auto_size = True
                tf.text = "Competent cell:" + info[image_path_list[0]][1] + "\n" + 'IPTG:' + info[image_path_list[0]][
                    2] + "\n" + "Temperature:" + info[image_path_list[0]][3] + "\n" + "Induction time:" + info[image_path_list[0]][4]
                #####左对齐
                for paragraph in tf.paragraphs:
                    paragraph.font.size = Pt(10)
                    # paragraph.line_spacing = 1
                    paragraph.font.name = "Times New Roman"
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.font.color.rgb = RGBColor(0, 0, 0)
                    paragraph.font.name = "Times New Roman"

            if len(image_path_list)==2:
                pic_area_left = Cm(13.3)
                pic_area_top = Cm(6.88)

                spacing = Cm(2.4)
                # 循环添加图片
                pic_width = Cm(2.88)
                pic_aspect_ratio = 2.88 / 5.1
                pic_height = pic_width / pic_aspect_ratio

                for i, object, in enumerate(image_path_list):
                    image_path = os.path.join(path, object)
                    # 计算当前图片的左侧和顶部位置
                    pic_left = pic_area_left + i * spacing + i * pic_width
                    pic_top = pic_area_top
                    # 添加图片到幻灯片
                    slice_opju.shapes.add_picture(image_path, pic_left, pic_top, pic_width, pic_height)
                    txBox_left = pic_left-Cm(0.5)
                    txBox_top = pic_top + pic_height + Cm(0.2)
                    txBox_width = pic_width
                    txBox_height = Cm(3)
                    txBox = slice_opju.shapes.add_textbox(txBox_left, txBox_top, txBox_width, txBox_height)
                    tf = txBox.text_frame
                    # tf.vertical_anchor = MSO_ANCHOR.TOP_LEFT
                    tf.auto_size = True

                    tf.text = "Competent cell:" + info[object][1] + "\n" + 'IPTG:' + info[object][
                        2] + "\n" + "Temperature:" + info[object][3] + "\n" + "Induction time:" + info[object][4]
                    #####左对齐
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(10)
                        # paragraph.line_spacing = 1
                        paragraph.font.name = "Times New Roman"
                        paragraph.alignment = PP_ALIGN.LEFT
                        paragraph.font.color.rgb = RGBColor(0, 0, 0)
                        paragraph.font.name = "Times New Roman"

            if len(image_path_list)==3:
                pic_area_left = Cm(11.85)
                pic_area_top = Cm(6.05)

                spacing = Cm(5)
                # 循环添加图片
                pic_width = Cm(2.88)
                pic_aspect_ratio = 2.88 / 5.1
                pic_height = pic_width / pic_aspect_ratio

                for i, object, in enumerate(image_path_list):
                    image_path = os.path.join(path, object)
                    # 计算当前图片的左侧和顶部位置
                    pic_left = pic_area_left + i * spacing
                    pic_top = pic_area_top
                    # 添加图片到幻灯片
                    slice_opju.shapes.add_picture(image_path, pic_left, pic_top, pic_width, pic_height)
                    txBox_left = pic_left-Cm(0.5)
                    txBox_top = pic_top + pic_height + Cm(0.2)
                    txBox_width = pic_width
                    txBox_height = Cm(3)
                    txBox = slice_opju.shapes.add_textbox(txBox_left, txBox_top, txBox_width, txBox_height)
                    # tf.vertical_anchor = MSO_ANCHOR.TOP_LEFT
                    tf = txBox.text_frame
                    tf.auto_size = True

                    tf.text = "Competent cell:" + info[object][1] + "\n" + 'IPTG:' + info[object][2] + "\n" + "Temperature:" + info[object][3] + "\n" + "Induction time:" + info[object][4]
                    #####左对齐
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(10)
                        # paragraph.line_spacing = 1
                        paragraph.font.name = "Times New Roman"
                        paragraph.alignment = PP_ALIGN.LEFT
                        paragraph.font.color.rgb = RGBColor(0, 0, 0)
                        paragraph.font.name = "Times New Roman"

            if 3 < len(image_path_list) <= 6:
                pic_area_left = Cm(12)
                pic_area_top = Cm(4)
                spacing = Cm(3)
                # 循环添加图片
                pic_width = Cm(1.97)
                pic_aspect_ratio = 2.88 / 5.1
                pic_height = pic_width / pic_aspect_ratio

                for i, object, in enumerate(image_path_list):
                    image_path = os.path.join(path, object)

                    # 计算当前图片的左侧和顶部位置
                    if i >= 3:
                        pic_left = pic_area_left + (i - 3) * spacing + (i - 3) * pic_width
                        pic_top = Cm(10.5)
                    else:
                        pic_left = pic_area_left + i * spacing + i * pic_width
                        pic_top = pic_area_top
                    # 添加图片到幻灯片
                    slice_opju.shapes.add_picture(image_path, pic_left, pic_top, pic_width, pic_height)
                    txBox_left = pic_left-Cm(0.5)
                    txBox_top = pic_top + pic_height + Cm(0.2)
                    txBox_width = pic_width
                    txBox_height = Cm(3)
                    txBox = slice_opju.shapes.add_textbox(txBox_left, txBox_top, txBox_width, txBox_height)
                    tf = txBox.text_frame
                    tf.auto_size = True

                    tf.text = "Competent cell:" + info[object][1] + "\n" + 'IPTG:' + info[object][2] + "\n" + "Temperature:" + info[object][3] + "\n" + "Induction time:" + info[object][4]
                    #####左对齐
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(10)
                        # paragraph.line_spacing = 1
                        paragraph.font.name = "Times New Roman"
                        paragraph.alignment = PP_ALIGN.LEFT
                        paragraph.font.color.rgb = RGBColor(0, 0, 0)
                        paragraph.font.name = "Times New Roman"

            if 6 < len(image_path_list) <= 12:
                pic_area_left = Cm(12.23)
                pic_area_top = Cm(2.99)
                spacing = Cm(3.5)
                # 循环添加图片
                pic_width = Cm(1.55)
                pic_aspect_ratio = 2.88 / 5.1
                pic_height = pic_width / pic_aspect_ratio

                for i, object, in enumerate(image_path_list):
                    image_path = os.path.join(path, object)

                    # 计算当前图片的左侧和顶部位置
                    if i <=3:
                        pic_left = pic_area_left + i * spacing
                        pic_top = pic_area_top

                    if 3< i <= 7:
                        pic_left = pic_area_left + (i - 4) * spacing
                        pic_top = Cm(7.56)

                    if i >7:
                        pic_left = pic_area_left + (i - 8) * spacing
                        pic_top = Cm(11.98)

                    # 添加图片到幻灯片
                    slice_opju.shapes.add_picture(image_path, pic_left, pic_top, pic_width, pic_height)
                    txBox_left = pic_left-Cm(0.5)
                    txBox_top = pic_top + pic_height-Cm(0.2)
                    txBox_width = pic_width
                    txBox_height = Cm(3)
                    txBox = slice_opju.shapes.add_textbox(txBox_left, txBox_top, txBox_width, txBox_height)
                    tf = txBox.text_frame
                    tf.auto_size = True

                    tf.text = "Competent cell:" + info[object][1] + "\n" + 'IPTG:' + info[object][
                        2] + "\n" + "Temperature:" + info[object][3] + "\n" + "Induction time:" + info[object][4]
                    #####左对齐
                    for paragraph in tf.paragraphs:
                        paragraph.font.size = Pt(10)
                        # paragraph.line_spacing = 1
                        paragraph.font.name = "Times New Roman"
                        paragraph.alignment = PP_ALIGN.LEFT
                        paragraph.font.color.rgb = RGBColor(0, 0, 0)
                        paragraph.font.name = "Times New Roman"

            # 图片置于文本框下方
            # slice.shapes._spTree.insert(1,pic._element)
            # pptx_trans_membrane_helices_prediction_text = slice_opju.placeholders[21]

            # conclusion
            info_text_frame = opju_title3.text_frame
            info_paragraph = info_text_frame.add_paragraph()

            conclusion_name = info_text_frame.paragraphs[0].add_run()
            conclusion_name.text = "Conclusion: "
            conclusion_name.font.bold = True

            conclusion_info = info_text_frame.paragraphs[0].add_run()
            conclusion_info.text = "Based on the FSEC profile, ***"

        ppt_file_path = os.path.join(self.save_dir, "ppt", r'FSEC_Analysis_{}.pptx'.format(PublicScript.get_current_date()))
        self.pptx_template_file_path = ppt_file_path
        file.save(ppt_file_path)

    def opju2ppt(self):
        pythoncom.CoInitialize()

        ppt_app = Dispatch('PowerPoint.Application', pythoncom.CoInitialize())
        # ppt_app.Visible = False  # 显式打开PPT 调试设置True
        ppt = ppt_app.Presentations.Open(self.pptx_template_file_path, WithWindow=False)  # 打开ppt

        for page_index, column_name in enumerate(self.ppt_opju_name_list):
            print("    生成PPT: '{}'".format(column_name), flush=True)
            opju_file = os.path.join(self.save_dir, "opju", "{}.opju".format(column_name))

            ppt.Slides(page_index + 3).Shapes.AddOLEObject(Left=0.51 * 28.35, Top=5.14 * 28.35,
                FileName=opju_file,
                )  # 在ppt第一页插EXCEL表，设置为图标模式，并设置位置

            pids_list = psutil.pids()

            for i in pids_list:
                try:
                    p = psutil.Process(i)

                    if "origin".upper() in p.name().upper():
                        p.terminate()
                except:
                    pass
        opjuppt_file = os.path.join(self.save_dir, "ppt", r'FSEC_Analysis_{}.pptx'.format(PublicScript.get_current_date()))
        ppt.SaveAs(opjuppt_file)  # 保存
        os.remove(self.pptx_template_file_path)
        print("结果文件为{}".format(os.path.basename(opjuppt_file)), flush=True)

    def SubElement(self, parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def _set_cell_border(self, cell, border_color="000000", border_width='22500',
                         position_list=['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']):
        """ Hack function to enable the setting of border width and border color
            - bottom border only at present
            (c) Steve Canny
        """
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        # position_list = ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']
        for position in position_list:
            lnR = self.SubElement(
                tcPr, position, w=border_width, cap='flat', cmpd='sng', algn='ctr')
            solidFill = self.SubElement(lnR, 'a:solidFill')
            srgbClr = self.SubElement(solidFill, 'a:srgbClr', val=border_color)
            lnR_prstDash = self.SubElement(lnR, 'a:prstDash', val='solid')
            lnR_round_ = self.SubElement(lnR, 'a:round')
            lnR_headEnd = self.SubElement(lnR, 'a:headEnd', type='none', w='med', len='med')
            lnR_tailEnd = self.SubElement(lnR, 'a:tailEnd', type='none', w='med', len='med')

        return cell

@Gooey(
    #advanced=True,
    # tabbed_groups=True,
    # navigation='Tabbed',
    richtext_controls=True,  # 打开终端对颜色支持
    program_name="FSEC数据分析",  # 程序名称
    # encoding="gbk",  # 设置编码格式，打包的时候遇到问题
    # progress_regex=r"^progress: (\d+)%$",  # 正则，用于模式化运行时进度信息
    clear_before_run=True,
    required_cols=2
)
def main():

    settings_msg = '本程序提供FSEC数据分析'
    parser = GooeyParser(description=settings_msg)  # 添加上方的应用信息

    subs = parser.add_subparsers(help='commands', dest='command')

    # 导入配置文件
    data_file_path, pptx_template_file_path, complex_protein_info_xlsx, writer, a, b, peak_x_error, peak_prominence, Ex_Wavelength, save_dir,Free_GFP,ppt_file = Fsec.load_configuration()

    """
    FSEC数据分析
    """
    SingleSearchparser = subs.add_parser('FSEC数据分析')

    SingleSearchparser.add_argument("Fsec_txt_FilePath", metavar='Fsec数据文本',
                                    help="选择Fsec数据文本所在的目录(!!提醒：文本命名一定要以质粒编号开头)", widget='DirChooser',
                                    default=data_file_path)

    SingleSearchparser.add_argument("Ppt_Template_FilePath", metavar='PPT模板',
                                    widget='FileChooser', default= pptx_template_file_path if pptx_template_file_path else os.path.join(os.getcwd(), "Template.pptx"))

    SingleSearchparser.add_argument("ComplexProtein_FilePath", metavar='复合蛋白信息表',
                                    widget='FileChooser', default=complex_protein_info_xlsx if complex_protein_info_xlsx else r"\\192.168.1.52\佰翱得共享文件\生物部\生物公共资料\1. 项目管理平台\4.1 【项目】质粒构建\3 质粒入库\复合蛋白信息表.xlsx")

    SingleSearchparser.add_argument("writer", metavar='姓名', default=writer if writer else "填写姓名")

    SingleSearchparser.add_argument("standard_a", metavar='标准曲线参数: a', help="标准曲线形式为: Log(y)=a*(x - 0.7) + b", default=a)
    SingleSearchparser.add_argument("standard_b", metavar='标准曲线参数: b', help="标准曲线形式为: Log(y)=a*(x - 0.7) + b", default=b)

    SingleSearchparser.add_argument("peak_x_error", metavar='蛋白峰x轴的相对误差值(取值范围：[0, +∞)), 默认值为0.05', help="x轴误差小于指定值,判定为同一蛋白峰", default= peak_x_error if peak_x_error else 0.05)
    SingleSearchparser.add_argument("peak_prominence", metavar='蛋白峰的突起程度(取值范围：[0, +∞)), 默认值为0.5', help="取值越大, 峰的突起越显著; 峰的突起大于指定程度才判定为蛋白峰", default= peak_prominence if peak_prominence else 0.5)

    SingleSearchparser.add_argument("Ex_Wavelength", metavar='激发波长, 默认值为480', default=Ex_Wavelength if Ex_Wavelength else 480)

    SingleSearchparser.add_argument("outdir", metavar='结果保存文件夹',
                                    help="选择Fsec数据文本所在的目录", widget='DirChooser',
                                    default=save_dir if save_dir else os.getcwd())
    SingleSearchparser.add_argument("Free_GFP", metavar='标记名称',help='改成需要的标记名称，默认为His lite probe', default='His lite probe')

    ######加入胶图的pptx文件
    SingleSearchparser.add_argument("ppt_file", metavar='小试PPT文件夹',
                                    widget='DirChooser')
    ######将ppt胶图保存在文件夹里
    # SingleSearchparser.add_argument("image_file_path", metavar='小试ppt的图片保存路径',
    #                                  widget='DirChooser',
    #                                 default=os.getcwd())


    # 参数
    # data_file_path = r"C:\Users\Lenovo\Desktop\Fsec Data\FSEC PPT测试\20220715 18715 18926 19976 19714 18874 19715 1-ID0020"
    # pptx_template_file_path = r"D:\BaiduSyncdisk\Data\PyCharm\PyCharm_2018.1.4_Window\data\SZP\项目部门\徐霞\2022年\FSEC Analysis\Template.pptx"
    # complex_protein_info_xlsx = r"D:\BaiduSyncdisk\Data\PyCharm\Data\项目部门\嵇歆彧\Protein Catalog\产品目录编码更新\复合蛋白信息表保存目录\复合蛋白信息表_2022-07-21 10-52-11.xlsx"
    # standard_curve_dict = {"a": 1, "b": 1}
    # save_dir = r"D:\BaiduSyncdisk\Data\PyCharm\Data\项目部门\徐霞\2022年\FSEC Analysis\data"

    """
       程序功能实现
    """
    args = parser.parse_args()

    # 共享路径中蛋白编码查询系统保存文件的路径

    """界面获得的数据进行处理"""
    if args.command == "FSEC数据分析":
        data_file_path = args.Fsec_txt_FilePath
        pptx_template_file_path = args.Ppt_Template_FilePath
        complex_protein_info_xlsx = args.ComplexProtein_FilePath
        writer = args.writer
        standard_curve_dict = {"a": float(args.standard_a), "b": float(args.standard_b)}
        peak_x_error = float(args.peak_x_error)
        peak_prominence = float(args.peak_prominence)
        Ex_Wavelength = args.Ex_Wavelength
        save_dir = args.outdir
        Free_GFP = args.Free_GFP
        ppt_file=args.ppt_file
        # image_file_path=args.image_file_path


        # 保存设置
        configuration_dict = {"data_file_path": data_file_path,
                              "pptx_template_file_path": pptx_template_file_path,
                              "complex_protein_info_xlsx": complex_protein_info_xlsx,
                              "writer": writer,
                              "standard_curve_dict": standard_curve_dict,
                              "peak_x_error": peak_x_error,
                              "peak_prominence": peak_prominence,
                              "Ex_Wavelength": Ex_Wavelength,
                              "save_dir": save_dir,
                              "Free_GFP": Free_GFP,
                              "ppt_file": ppt_file}

        with open("configuration.json", "w", encoding="utf-8") as f:
            json.dump(configuration_dict, f, indent=4, ensure_ascii=False)

        # 生成结果保存目录
        #####加上一个image，储存小试ppt的图片
        result_file_name_list = ["ogs", "opju", "ppt","image"]
        for result_file_name in result_file_name_list:
            result_file_name_dir = os.path.join(save_dir, result_file_name)
            if not os.path.exists(result_file_name_dir):
                os.makedirs(result_file_name_dir)

        # 实例化
        fsec_obj = Fsec(data_file_path, pptx_template_file_path, complex_protein_info_xlsx, writer, standard_curve_dict, peak_x_error, peak_prominence, Ex_Wavelength, save_dir,Free_GFP,ppt_file)

        # 获取复合蛋白信息表数据
        print("获取复合蛋白信息表数据", flush=True)
        fsec_obj.get_complex_protein_info()

        # 获取txt数据
        print("获取FSEC数据", flush=True)
        fsec_obj.get_fsec_data()

        # 处理txt 数据
        print("处理FSEC数据", flush=True)
        fsec_obj.analysis_fsec_data()

        #####将小试ppt图片保存
        fsec_obj.generate_image()

        #######根据plasmidid 选择图片插入
        fsec_obj.choose_image()

        # 生成PPT
        # print("获取复合蛋白信息表数据", flush=True)
        fsec_obj.generate_ppt()


        # opju2ppt
        print("生成PPT", flush=True)
        fsec_obj.opju2ppt()


if __name__ == '__main__':
    if sys.stdout.encoding != 'UTF-8':
        sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer, 'strict')
    # if sys.stderr.encoding != 'UTF-8':
    #     sys.stderr = codecs.getwriter('utf-8')(sys.stderr.buffer, 'strict')
    # 编码修复
    
    main()