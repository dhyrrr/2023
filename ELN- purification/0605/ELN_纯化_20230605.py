from selenium.webdriver.chrome.service import Service
import time
import win32api, win32con
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.action_chains import ActionChains
import base64
from pptx import Presentation
import re
import os
from selenium.webdriver.common.keys import Keys
from PIL import Image
from pptx_tools import utils
from szpscript.Broswer_script import Broswer
import traceback
from gooey import Gooey, GooeyParser
import json
import  sys
import codecs
import pptx
from datetime import datetime

'''
4.3 在程序文件夹下创建image文件 content.json 
合并  click  sendkeys   table

复制模板的时候 一个一个复制

加异常处理
4.4  重复step 的获取及处理   
4.5  测试   copy的时候模板打开两次就会出现在中央了  print修改
4.6  1.ppt第一页信息获取if   len(text.splitlines())==3 
4.6_2   创建表单名称projectname  如果没有QC3的话把标识删了  图合并的话就把图标也合并 Affi-HP，ion  size

4.10-4.13 这周
4.10   改成'Ni Bestarose FF'  in style    复制模板改为===而非包含   项目编号不对的  pass，，，直接先那样吧，批量的时候再说
批量的脚本拿到各个ppt对应的内容和图片,  修改index的位置，确保拿到index   
4.11      Size Ion 等没有框的话在后面的那个加上去用if 'Conc' in text and 'A260/A280' in text 来判断是否存在东西   如果不含有那些东西的话，把模板相关内容删除  QC的图片进行处理，cell的问题   
4.12      多行表格和计算的问题 ok   包含体的模板ok  
4.13      qc的多个图片插入问题已解决   外部项目   还有那个dili模板的  写在下面  
********所有的都加异常处理 已解决*********
4.14 字母大小写的  ok  最终版本
4.18 Size  的冒号，cell&cy拿不到字典信息
4.20的修改了main里的内部、外部参数，之前的4.14版本就有问题，没有改过来，只不过最终的all没有调用main的内容所以没有影响，现在改过来 cell的拿不到信息？？？？拿不到是因为cell&cy的字典位置写的有问题，不用瞎改，，，
在表格里找到名字和项目对应的行然后点击
4.24 修改有些信息匹配不上的bug,没发现哪里有问题啊 加入Objective
4.27 Size那个1和2合在一起
4.28 Size信息对不上的问题
5.5 项目编号点不进去的问题 刷新 等一下  图片没有结论的  方框里面有多行的问题
5.9     1.  object部分只要蛋白名称，不要Bp号和日期  concentrate 加d
		2.亲和his的模板 把excel、talon等具体写清楚   3.Ion-exchange 的柱子类型写上    4.Size-exclusion chromatography 浓度的信息不对，，，，上次修改的时候注释掉了
		5.有份报告进不去,因为先读取了conclusion的内容，没有先获得experient，从而没有拿到contnet，导致失败，获取文本的时候ppt从上往下开始读取
		6.复制模板，模板改名字的时候会卡住，增加rang和timeout
		7.如果QC没有第二章图片，删除模板里的内容
		5.10   8.stop不管用的问题因为那个选项框是在外面的，所以失败后点击不管用
5.30  把内部和外部的合在一起，形成大脚本
6.1 夏凡的那个项目太少了，所以定位不到，导致一个都不行
    批量化，修改相对应的内容，如果项目编号不对或者创建无权限的直接pass
6.2   1.网络不好的直接关掉
      2.截图
6.5   1.gallery和非gallery的，不是外部和内部
      2. ppt第一次失败后尝试第二次、第三次
      3.clickoption   的关闭页面要写在外面，在函数里面只能结束当前函数的，不能直接结束
      4.driver.switch_to.default_content()   如果中间有问题导致没有回到默认的地方，导致后面的定位不到
      5.打包后没有提示了
	  6.打包后跑完也没有没有提示，点击stop不管用
	  7.clickoption 的尝试次数忘记改了，已经打包发走了，先这样吧，下次更新改过来
	  
'''
if sys.stdout.encoding != 'UTF-8':
    sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer, 'strict')
if sys.stderr.encoding != 'UTF-8':
    sys.stderr = codecs.getwriter('utf-8')(sys.stderr.buffer, 'strict')

def make_print_to_file():
    savename = "out.log"
    def write_to_log(message):
        with open(savename, "a", encoding='utf-8') as log_file:
            log_file.write(message)
        sys.__stdout__.write(message)

    sys.stdout.write = write_to_log

class Params:
    def __init__(self):
        pass
    website = r"https://scilligence.net/Biortus/ELN/Explorer.aspx"
    selenium_waittime = 3000
    broswer_name = "Chrome"
    chrome_driver_website = r"https://registry.npmmirror.com/binary.html?path=chromedriver/"
    chrome_driver_path = r"./chromedriver.exe"

class eln:

    __chrome_driver_path = Params.chrome_driver_path

    def __init__(self,account, password, ppt_file, image_file_path, name):
        self.account=account
        self.password=password
        self.ppt_file=ppt_file
        self.image_file_path=image_file_path
        self.name = name
        self.chrome_driver_path = eln.__chrome_driver_path
        self.driver = self.build_selenium_obj(self.chrome_driver_path)
        self.website = Params.website
        self.waittime = Params.selenium_waittime
        self.witness_finished_status = False  # 判定witness任务是否处理完, T
    #######静态方法通常用于定义与类相关但不依赖于类变量和实例变量的功能
    @staticmethod
    def check_broswer():
        # 尝试3次
        for i in range(5):
            broswer_version = Broswer.get_broswer_version(Params.broswer_name)

            if broswer_version:
                # 检查驱动是否存在
                chrome_driver_new_path = eln.__chrome_driver_path.replace(".exe",
                                                                              "_V{0}.exe".format(broswer_version))
                if not os.path.exists(chrome_driver_new_path):
                    print("下载驱动中!",flush=True)
                    # 下载驱动
                    Broswer.download_chrome_driver(broswer_version)
                    # 根据版本更改驱动名字
                    if os.path.exists(eln.__chrome_driver_path):
                        os.rename(eln.__chrome_driver_path, chrome_driver_new_path)
                        eln.__chrome_driver_path = chrome_driver_new_path
                        return True
                else:
                    # 驱动已经存在, 根据版本更改驱动名字
                    eln.__chrome_driver_path = chrome_driver_new_path
                    return True

            else:
                # 安装
                win32api.MessageBox(0, "Chrome浏览器未安装, 现在开始安装", "提醒", win32con.MB_TOPMOST)
                Broswer.install_chrome()
        else:
            # 安装浏览器或驱动失败
            return False
    @staticmethod
    def build_selenium_obj(chrome_driver_path):
        #######解决chrome正在收到测试软件监控的问题
        s = Service(executable_path=chrome_driver_path)
        options = webdriver.ChromeOptions()
        options.add_experimental_option("excludeSwitches",['enable-automation'])
        driver = webdriver.Chrome(service=s,options=options)
        driver.command_executor._commands["send_command"] = ("POST", '/session/$sessionId/chromium/send_command')
        params = {'cmd': 'Page.setDownloadBehavior',
                  'params': {'behavior': 'allow', 'downloadPath': os.getcwd()}}
        driver.execute("send_command", params)
        return driver
    ######读取pptx内容，保存图片及相关信息
    def read_pptfile_out(self):
        print("提取ppt信息",flush=True)
        ppt_file = self.ppt_file
        # image_file_path = os.path.abspath('image_file_path')
        image_file_path= self.image_file_path
        utils.save_pptx_as_png(image_file_path, ppt_file, overwrite_folder=True)
        prs = Presentation(ppt_file)
        slide1_dic = {}
        ProjectInformation = []
        purification = []
        storageBuffer_list = []
        content = {}

        for index, slide in enumerate(prs.slides):
            sortedShapes = sorted(slide.shapes, key=lambda x: (x.left))
            slide_height = prs.slide_height / 914400 * 2.54
            for shape in sortedShapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text = text_frame.text
                    if '项目编号' in text:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame
                                # 读取文本框内容
                                text = text_frame.text
                                if '项目编号' in text:
                                    try:
                                        try:
                                            ProjectName = text.split(":")[1].strip()
                                            slide1_dic['ProjectName'] = ProjectName

                                        except:
                                            ProjectName = text.split("：")[1].strip()
                                            slide1_dic['ProjectName'] = ProjectName
                                    except:
                                        print('检查项目编号处内容',flush=True)
                                        slide1_dic['ProjectName'] = ''
                                if len(text.splitlines())==3:
                                    try:
                                        code = text.splitlines()[0].split(",")[0]
                                        slide1_dic['code'] = code
                                    except:
                                        slide1_dic['code'] = ''
                                    try:
                                        条目名称 = text.splitlines()[0].split(",")[1].strip()
                                        slide1_dic['条目名称'] = 条目名称
                                    except:
                                        slide1_dic['条目名称'] = ''
                                    try:
                                        NotebookName = 'Protein purification'
                                        slide1_dic['NotebookName'] = NotebookName
                                    except:
                                        slide1_dic['NotebookName'] = ''
                                    try:
                                        ProteinNname = text.splitlines()[1].strip()
                                        slide1_dic['ProteinNname'] = ProteinNname
                                    except:
                                        slide1_dic['ProteinNname'] = ''

                    if text == 'Project Information':
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame
                                text = text_frame.text
                            if shape.has_table:
                                for row in shape.table.rows:  # 读每行
                                    for cell in row.cells:  # 读一行中的所有单元格
                                        c = cell.text
                                        ProjectInformation.append(c)
                    if "Source" and "Column used" and "Storage information" in text:
                        sortedShapes = sorted(slide.shapes, key=lambda x: (x.top))
                        for shape in sortedShapes:
                            if shape.has_table:
                                # 获取表格对象
                                table = shape.table
                                # 创建一个列表用于保存表格内容
                                table_data = []
                                # 遍历表格中的所有行和列
                                for i, row in enumerate(table.rows):
                                    row_data = []
                                    for j, cell in enumerate(row.cells):
                                        # 获取单元格中的文本
                                        cell_text = cell.text_frame.text.strip()
                                        row_data.append(cell_text)
                                    # 将该行添加到表格内容列表中
                                    table_data.append(row_data)
                                purification.append(table_data)

                    parts = re.split("Step \d*", text)
                    # print(parts)
                    if len(parts) > 1:
                        step_name = parts[1].split(":")[1].strip()
                        step_name = step_name.lower()

                        if step_name == 'Digestion & Affinity chromatography'.lower():
                            content[index] = {}
                            content[index]['Digestion & Affinity chromatography'] = {}
                            content[index]['Digestion & Affinity chromatography']['index'] = index
                            DiAff_frame_list = []
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'Conc' in text and 'Total amount' in text:
                                        DiAff_frame_list.append(text)
                                    if 'Conclusions' in text:
                                        try:
                                            conclusions = text.split("Conclusions:")[1].strip()
                                            content[index]['Digestion & Affinity chromatography'][
                                                'conclusions'] = conclusions
                                        except:
                                            print('检查Digestion_Affinity_chromatography 结论',flush=True)
                                            content[index]['Digestion & Affinity chromatography']['conclusions'] = ''

                                        try:
                                            protein = conclusions.split('tagged ')[1].split('protease')[0].strip() + ' protease'
                                            protein1 = protein.split('protease')[0].strip()
                                            content[index]['Digestion & Affinity chromatography']['protein'] = protein
                                            content[index]['Digestion & Affinity chromatography']['protein1'] = protein1
                                        except:
                                            print('检查Digestion & Affinity chromatography结论',flush=True)
                                            content[index]['Digestion & Affinity chromatography']['protein'] = ''
                                            content[index]['Digestion & Affinity chromatography']['protein1'] = ''
                                    if 'Experiment' in text:
                                        try:
                                            sample = 'from ' + text.splitlines()[1].split("Sample: From")[1].strip()
                                            content[index]['Digestion & Affinity chromatography'][
                                                'sample_value'] = sample
                                        except:
                                            print('检查Digestion & Affinity chromatography  Sample',flush=True)
                                            content[index]['Digestion & Affinity chromatography']['sample_value'] = ''
                                        try:
                                            column = text.split('Column:')[1].split('Sample')[0].strip()
                                            content[index]['Digestion & Affinity chromatography']['column'] = column
                                        except:
                                            print('检查Digestion & Affinity chromatography  column',flush=True)
                                            content[index]['Digestion & Affinity chromatography']['column'] = ''
                                    ######加上对应的内容

                                    result_top = ''
                                    conclusion_top = ''
                                    if "Results" in text:
                                        for shape in slide.shapes:
                                            if shape.has_text_frame:
                                                text_frame = shape.text_frame
                                                text = text_frame.text
                                                if "Experiment" in text:
                                                    result_top = (shape.height + shape.top) / 914400 * 2.54
                                                    # print(result_top)
                                                if "Conclusions" in text:
                                                    conclusion_top = shape.top / 914400 * 2.54
                                        try:
                                            if conclusion_top == '':
                                                conclusion_top = 16.5
                                                bottom_position = float(conclusion_top) / slide_height
                                            else:
                                                bottom_position = float(conclusion_top) / slide_height
                                            image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                            left_position = 0
                                            top_position = float(result_top) / slide_height
                                            # bottom_position = float(conclusion_top) / slide_height
                                            right_position = 1
                                            img = Image.open(image_name)
                                            img_size_width, img_size_height = img.size
                                            crop_left_position = left_position * img_size_width
                                            crop_top_position = top_position * img_size_height
                                            crop_right_position = right_position * img_size_width
                                            crop_bottom_position = bottom_position * img_size_height
                                            name = 'Digestion Affinity chromatography_' + str(index) + ".jpg"
                                            crop_image_path = os.path.join(image_file_path, name)
                                            crop_img = img.crop(
                                                (crop_left_position, crop_top_position, crop_right_position,
                                                 crop_bottom_position))
                                            if crop_img.mode == 'RGBA':
                                                crop_img = crop_img.convert("RGB")
                                            crop_img.save(crop_image_path, quality=95, subsampling=2)
                                        except:
                                            print('检查Digestion Affinity chromatography图片是否有问题',flush=True)
                                            pass
                                try:
                                    for k in list(content[index]['Digestion & Affinity chromatography'].keys()):
                                        if 'index' == k:
                                            Conc_list = []
                                            A260_A280_list = []
                                            Total_amount_list = []
                                            for i in DiAff_frame_list:

                                                try:
                                                    Conc = i.split('Conc.:')[1].split("mg/ml")[0].strip() + " "
                                                    Conc_list.append(Conc)
                                                    A260_A280 = i.split('A260/A280:')[1].split("Total")[0].strip()
                                                    A260_A280_list.append(A260_A280)
                                                    Total_amount = i.split("amount:")[1].strip()
                                                    Total_amount_list.append(Total_amount)
                                                    peak = i.split('Conc.:')[0].strip().replace('\n', ' ')
                                                    peak_list.append(peak)
                                                except:
                                                    print('检查Digestion & Affinity chromatography信息', flush=True)
                                            content[index]['Digestion & Affinity chromatography']['Conc'] = Conc_list
                                            content[index]['Digestion & Affinity chromatography']['A260_A280'] = A260_A280_list
                                            content[index]['Digestion & Affinity chromatography']['Total_amount_value'] = Total_amount_list
                                except:
                                    pass
                        if step_name == 'QCs'.lower():
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'QC1: SDS-PAGE' in text:
                                        content[index] = {}
                                        content[index]['QC'] = {}
                                        content[index]['QC']['index'] = index
                                        order = index
                                        for shape in slide.shapes:
                                            if shape.has_text_frame:
                                                text_frame = shape.text_frame
                                                text = text_frame.text
                                                if 'Batch No' in text:
                                                    try:
                                                        content[index]['QC']['Batch'] = text.split(':')[1].strip()
                                                        top = (shape.top + shape.height) / 914400 * 2.54
                                                        left_position = 0
                                                        top_position = top / slide_height
                                                        bottom = 16.5
                                                        bottom_position = bottom / slide_height
                                                        right_position = 1
                                                        image_name = os.path.join(image_file_path,"幻灯片{}.PNG".format(index + 1))
                                                        img = Image.open(image_name)
                                                        img_size_width, img_size_height = img.size
                                                        crop_left_position = left_position * img_size_width
                                                        crop_top_position = top_position * img_size_height
                                                        crop_right_position = right_position * img_size_width
                                                        crop_bottom_position = bottom_position * img_size_height
                                                        name = 'QC' + '_' + str(order) + ".jpg"
                                                        crop_image_path = os.path.join(image_file_path, name)
                                                        crop_img = img.crop(
                                                            (crop_left_position, crop_top_position, crop_right_position,
                                                             crop_bottom_position))
                                                        if crop_img.mode == 'RGBA':
                                                            crop_img = crop_img.convert("RGB")
                                                        crop_img.save(crop_image_path, quality=95, subsampling=2)
                                                    except:
                                                        print('检查QC图片是否有问题')
                                                        pass

                                    ####截QC2的图
                                    if 'QC2: LC-MS (zoom in)' in text:
                                        content[order]['QC'][index] = index
                                        if str(index - 1) != str(order):
                                            print('QC2: LC-MS (zoom in)包含多张图片，请注意检查', flush=True)
                                        for shape in prs.slides[index].shapes:
                                            if shape.has_text_frame:
                                                text_frame = shape.text_frame
                                                # 读取文本框内容
                                                text = text_frame.text
                                                if "Conclusions" in text:
                                                    try:
                                                        storageBuffer = \
                                                        text.split("buffer:")[1].strip().rpartition('.')[0]
                                                        storageBuffer_list.append(storageBuffer)
                                                    except:
                                                        print('检查QC Conclusions',flush=True)
                                                        storageBuffer = ''
                                                        storageBuffer_list.append(storageBuffer)
                                                try:
                                                    result_top = ''
                                                    conclusion_top = ''
                                                    if "zoom in" in text:
                                                        for shape in slide.shapes:
                                                            if shape.has_text_frame:
                                                                text_frame = shape.text_frame
                                                                text = text_frame.text
                                                                if "zoom in" in text:
                                                                    result_top = (shape.height + shape.top) / 914400 * 2.54
                                                                    # print(result_top)
                                                                if "Conclusions" in text:
                                                                    conclusion_top = shape.top / 914400 * 2.54

                                                        if conclusion_top == '':
                                                            conclusion_top = 16.5
                                                            bottom_position = float(conclusion_top) / slide_height
                                                        else:
                                                            bottom_position = float(conclusion_top) / slide_height

                                                        image_name = os.path.join(image_file_path,
                                                                                  "幻灯片{}.PNG".format(index + 1))
                                                        left_position = 0
                                                        top_position = float(result_top) / slide_height
                                                        bottom_position = float(conclusion_top) / slide_height
                                                        right_position = 1
                                                        img = Image.open(image_name)
                                                        img_size_width, img_size_height = img.size
                                                        crop_left_position = left_position * img_size_width
                                                        crop_top_position = top_position * img_size_height
                                                        crop_right_position = right_position * img_size_width
                                                        crop_bottom_position = bottom_position * img_size_height
                                                        name = 'QCzoomin' + '_' + str(order) + '_' + str(index) + ".jpg"
                                                        crop_image_path = os.path.join(image_file_path, name)
                                                        # print(crop_image_path)
                                                        crop_img = img.crop(
                                                            (crop_left_position, crop_top_position, crop_right_position,
                                                             crop_bottom_position))
                                                        if crop_img.mode == 'RGBA':
                                                            crop_img = crop_img.convert("RGB")
                                                        crop_img.save(crop_image_path, quality=95, subsampling=2)
                                                except:
                                                    pass
                                                    print('检查QCzoom格式是否正确',flush=True)
                        if step_name == "Cell lysis & Centrifugation".lower():
                            try:
                                for shape in prs.slides[index + 1].shapes:
                                    if shape.has_text_frame:
                                        text_frame = shape.text_frame
                                        text = text_frame.text
                                        if 'Step' in text:
                                            if 'Inclusion body preparation' in text:
                                                name = 'Cell lysis & Inclusion body preparation'
                                            else:
                                                name = 'Cell lysis & Centrifugation'
                            except:
                                name = 'Cell lysis & Centrifugation'

                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'Sample' in text:
                                        content[index] = {}
                                        content[index][name] = {}
                                        content[index][name]['index'] = index
                                        try:
                                            m = text.split(":")[1].split("g")[0].strip()
                                            content[index][name]['m_value'] = m
                                        except:
                                            print('检查Cell lysis & Centrifugation 质量',flush=True)
                                            m = ''
                                            content[index][name]['m_value'] = m
                                        try:
                                            v = text.splitlines()[0].split("from")[1].strip()
                                            content[index][name]['v_value'] = v
                                        except:
                                            print('检查Cell lysis & Centrifugation 体积',flush=True)
                                            v = ''
                                            content[index][name]['v_value'] = v
                                        try:
                                            procedures = text.split("was resuspended")[1].split(".")[0].strip()
                                            content[index][name]['procedures'] = procedures
                                        except:
                                            procedures = ''
                                            content[index][name]['procedures'] = procedures
                                            print('检查Cell lysis & Centrifugation Procedures', flush=True)
                                    if 'high pressure' in text:
                                        try:
                                            pressure = text.split('pressure homogenizer (')[1].split('Bar')[0]
                                            content[index][name]['pressure'] = pressure
                                        except:
                                            print('检查Cell lysis & Centrifugation 压力',flush=True)
                                            pressure = ''
                                            content[index][name]['pressure'] = pressure
                                    if 'Lysis buffer' in text:
                                        try:
                                            buff = text.split("mM")[0].split("Lysis buffer: ")[1].strip() + " ml"
                                            content[index][name]['buff_value'] = buff
                                        except:
                                            print('检查Cell lysis & Centrifugation Lysis buffer',flush=True)
                                            buff = ''
                                            content[index][name]['buff_value'] = buff
                                    if 'min' in text:
                                        try:
                                            ti = text.split("min at")[0].split("Centrifuged the lysate  with ")[1].strip() + " min"
                                            content[index][name]['ti_value'] = ti
                                        except:
                                            print('检查Cell lysis & Centrifugation  时间',flush=True)
                                            ti = ''
                                            content[index][name]['ti_value'] = ti

                        ######亲和住和离子住的内容一样
                        if step_name == 'Affinity chromatography'.lower():
                            sortedShapes = sorted(slide.shapes, key=lambda x: (x.top))
                            for shape in sortedShapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'Experiment' in text:
                                        if 'HiTrap Heparin HP' in text:
                                            content[index] = {}
                                            content[index]['Affinity chromatography (HiTrap Heparin HP)'] = {}
                                            content[index]['Affinity chromatography (HiTrap Heparin HP)']['index'] = index
                                            HP_Aff_frame_list = []
                                            for shape in slide.shapes:
                                                if shape.has_text_frame:
                                                    text_frame = shape.text_frame
                                                    text = text_frame.text
                                                    ######加上对应的内容
                                                    if 'Conc' in text and 'Total amount' in text:
                                                        HP_Aff_frame_list.append(text)

                                                    if 'Conclusions' in text:
                                                        try:
                                                            conclusions = text.split("Conclusions:")[1].strip()
                                                            content[index]['Affinity chromatography (HiTrap Heparin HP)']['conclusions'] = conclusions
                                                        except:
                                                            print('检查Affinity chromatography(HiTrap Heparin HP)结论',flush=True)
                                                            conclusions = ''
                                                            content[index][
                                                                'Affinity chromatography (HiTrap Heparin HP)'][
                                                                'conclusions'] = conclusions
                                                    if 'Experiment' in text:
                                                        try:
                                                            sample = 'from ' + text.split("Sample: From")[1].split(".")[0].strip()
                                                            content[index][
                                                                'Affinity chromatography (HiTrap Heparin HP)'][
                                                                'sample_value'] = sample
                                                        except:
                                                            print('检查Affinity chromatography(HiTrap Heparin HP)sample',flush=True)
                                                            sample = ''
                                                            content[index]['Affinity chromatography (HiTrap Heparin HP)']['sample_value'] = sample
                                                        try:
                                                            column = text.splitlines()[0].split('Column:')[1].strip()
                                                            content[index]['Affinity chromatography (HiTrap Heparin HP)']['column'] = column
                                                        except:
                                                            print(
                                                                '检查Affinity chromatography(HiTrap Heparin HP)   column ',flush=True)
                                                            column = ''
                                                            content[index]['Affinity chromatography (HiTrap Heparin HP)']['column'] = column
                                                    try:
                                                        result_top = ''
                                                        conclusion_top = ''
                                                        if "Results" in text:
                                                            for shape in slide.shapes:
                                                                if shape.has_text_frame:
                                                                    text_frame = shape.text_frame
                                                                    text = text_frame.text
                                                                    if "Experiment" in text:
                                                                        result_top = (shape.height + shape.top) / 914400 * 2.54
                                                                    if "Conclusions" in text:
                                                                        conclusion_top = shape.top / 914400 * 2.54
                                                            image_name = os.path.join(image_file_path,
                                                                                      "幻灯片{}.PNG".format(index + 1))
                                                            left_position = 0
                                                            top_position = float(result_top) / slide_height
                                                            bottom_position = float(conclusion_top) / slide_height
                                                            right_position = 1
                                                            img = Image.open(image_name)
                                                            img_size_width, img_size_height = img.size
                                                            crop_left_position = left_position * img_size_width
                                                            crop_top_position = top_position * img_size_height
                                                            crop_right_position = right_position * img_size_width
                                                            crop_bottom_position = bottom_position * img_size_height
                                                            name = 'Affinity chromatography亲和_' + str(index) + ".jpg"
                                                            crop_image_path = os.path.join(image_file_path, name)
                                                            crop_img = img.crop(
                                                                (crop_left_position, crop_top_position, crop_right_position,
                                                                 crop_bottom_position))
                                                            if crop_img.mode=='RGBA':
                                                                crop_img = crop_img.convert("RGB")
                                                            crop_img.save(crop_image_path, quality=95, subsampling=2)
                                                    except:
                                                        pass
                                            for k in list(content[index]['Affinity chromatography (HiTrap Heparin HP)'].keys()):
                                                if 'index' == k:
                                                    Conc_list = []
                                                    A260_A280_list = []
                                                    Total_amount_list = []
                                                    peak_list = []
                                                    for i in HP_Aff_frame_list:
                                                        if len(i.splitlines()) == 4:
                                                            Conc = i.splitlines()[1].split(":")[1].split("mg/ml")[0].strip() + " "
                                                            Conc_list.append(Conc)
                                                            A260_A280 = i.splitlines()[2].split(":")[1].strip()
                                                            A260_A280_list.append(A260_A280)
                                                            Total_amount = i.splitlines()[3].split(":")[1].strip()
                                                            Total_amount_list.append(Total_amount)
                                                            peak = i.splitlines()[0].strip()
                                                            peak_list.append(peak)
                                                        else:
                                                            print('检查Affinity chromatography 信息', flush=True)
                                                    content[index]['Affinity chromatography (HiTrap Heparin HP)'][
                                                        'Conc'] = Conc_list
                                                    content[index]['Affinity chromatography (HiTrap Heparin HP)'][
                                                        'A260_A280'] = A260_A280_list
                                                    content[index]['Affinity chromatography (HiTrap Heparin HP)'][
                                                        'Total_amount_value'] = Total_amount_list
                                                    content[index]['Affinity chromatography (HiTrap Heparin HP)'][
                                                        'Peak_value'] = peak_list

                                        else:
                                            Aff_frame_list = []
                                            content[index] = {}
                                            for shape in sortedShapes:
                                                if shape.has_text_frame:
                                                    text_frame = shape.text_frame
                                                    text = text_frame.text
                                                    if 'Experiment' in text:
                                                        try:
                                                            try:
                                                                style =text.splitlines()[0].split("Column:")[1].split(',')[0].strip()
                                                            except:
                                                                style =text.splitlines()[0].split("Resin:")[1].split(',')[0].strip()

                                                            if  'Ni Bestarose FF'.lower() in style  or  'Talon'.lower()  in style or  'excel'.lower()  in style or  'Protino'.lower()  in style or  'His FF'.lower()  in style or 'Excel'.lower()  in style:
                                                                column = ' (His)'
                                                            elif 'GST'.lower()  in style:
                                                                column = ' (GST)'
                                                            elif 'MBP'.lower()  in style:
                                                                column = ' (MBP)'
                                                            elif 'Strep'.lower()  in style:
                                                                column = ' (Strep)'
                                                            elif 'Flag'.lower()  in style:
                                                                column = ' (Flag)'
                                                            else:
                                                                column = ' (His)'
                                                        except:
                                                            column = ' (His)'

                                                        column_style = 'Affinity chromatography' + column

                                                        content[index][column_style] = {}
                                                        co = text.splitlines()[0].split(",")[1].strip() + ' '
                                                        content[index][column_style]['co_value'] = co
                                                        content[index][column_style]['index'] = index
                                                        content[index][column_style]['style']=style


                                                    if 'Conc' in text and 'Total amount' in text:
                                                        Aff_frame_list.append(text)

                                                    if 'Conclusions' in text:

                                                        try:
                                                            conclu = text.split(":")[1].replace('\n','')

                                                        except:
                                                            conclu = ''

                                                        content[index][column_style]['conclu'] = conclu


                                                    try:
                                                        result_top = ''
                                                        conclusion_top = ''
                                                        if "Results" in text:
                                                            for shape in slide.shapes:
                                                                slide_height = prs.slide_height / 914400 * 2.54
                                                                if shape.has_text_frame:
                                                                    text_frame = shape.text_frame
                                                                    text = text_frame.text
                                                                    if "Experiment" in text:
                                                                        result_top = (
                                                                                                 shape.height + shape.top) / 914400 * 2.54
                                                                    if "Conclusions" in text:
                                                                        conclusion_top = shape.top / 914400 * 2.54

                                                            if conclusion_top == '':
                                                                conclusion_top = 16.5
                                                                bottom_position = float(conclusion_top) / slide_height
                                                            else:
                                                                bottom_position = float(conclusion_top) / slide_height

                                                            image_name = os.path.join(image_file_path,
                                                                                      "幻灯片{}.PNG".format(index + 1))
                                                            left_position = 0
                                                            top_position = float(result_top) / slide_height
                                                            # bottom_position = float(conclusion_top) / slide_height
                                                            right_position = 1
                                                            img = Image.open(image_name)
                                                            img_size_width, img_size_height = img.size
                                                            crop_left_position = left_position * img_size_width
                                                            crop_top_position = top_position * img_size_height
                                                            crop_right_position = right_position * img_size_width
                                                            crop_bottom_position = bottom_position * img_size_height
                                                            name = 'Affinity chromatography_' + str(index) + ".jpg"
                                                            crop_image_path = os.path.join(image_file_path, name)
                                                            crop_img = img.crop(
                                                                (crop_left_position, crop_top_position, crop_right_position,
                                                                 crop_bottom_position))
                                                            if crop_img.mode=='RGBA':
                                                                crop_img = crop_img.convert("RGB")
                                                            crop_img.save(crop_image_path, quality=95, subsampling=2)
                                                    except:
                                                        pass
                                            for k in list(content[index][column_style].keys()):
                                                if 'index' == k:
                                                    Conc_list = []
                                                    A260_A280_list = []
                                                    Total_amount_list = []
                                                    peak_list = []
                                                    for i in Aff_frame_list:
                                                        try:
                                                            Conc = i.split('Conc.:')[1].split("mg/ml")[0].strip() + " "
                                                            Conc_list.append(Conc)
                                                            A260_A280 = i.split('A260/A280:')[1].split("Total")[
                                                                0].strip()
                                                            A260_A280_list.append(A260_A280)
                                                            Total_amount = i.split("amount:")[1].strip()
                                                            Total_amount_list.append(Total_amount)
                                                            peak = i.split('Conc.:')[0].strip().replace('\n', ' ')
                                                            peak_list.append(peak)
                                                        except:
                                                            print('检查Affinity chromatography  信息', flush=True)
                                                    content[index][column_style]['Conc'] = Conc_list
                                                    content[index][column_style]['A260_A280'] = A260_A280_list
                                                    content[index][column_style][
                                                        'Total_amount_value'] = Total_amount_list
                                                    content[index][column_style]['Peak_value'] = peak_list

                        if step_name == 'Size-exclusion chromatography'.lower():
                            content[index] = {}
                            content[index]['Size-exclusion chromatography'] = {}
                            size_frame_list = []
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'Experiment' in text:
                                        content[index]['Size-exclusion chromatography']['index'] = index
                                        try:
                                            column = text.splitlines()[0].split('Column:')[1].strip()
                                            content[index]['Size-exclusion chromatography']['column'] = column
                                        except:
                                            print('检查Size-exclusion chromatography   column',flush=True)
                                            column = ''
                                            content[index]['Size-exclusion chromatography']['column'] = column
                                        try:
                                            sample = text.split("Sample:")[1].split("Buffer")[0].strip()
                                            content[index]['Size-exclusion chromatography']['sample_value'] = sample
                                        except:
                                            print('检查Size-exclusion chromatography   sample',flush=True)
                                            sample = ''
                                            content[index]['Size-exclusion chromatography']['sample_value'] = sample
                                    if 'Conc' in text and 'Total amount' in text:
                                        size_frame_list.append(text)
                                    if 'Conclusions:' in text:
                                        try:
                                            conclu = text.split('Conclusions')[1].strip().replace('\n','')
                                            content[index]['Size-exclusion chromatography']['conclu'] = conclu
                                        except:
                                            print('检查Size-exclusion chromatography   结论',flush=True)
                                            conclu = ''
                                            content[index]['Size-exclusion chromatography']['conclu'] = conclu

                                    result_top = ''
                                    conclusion_top = ''
                                    if "Results" in text:
                                        for shape in slide.shapes:

                                            if shape.has_text_frame:
                                                text_frame = shape.text_frame
                                                text = text_frame.text
                                                if "Experiment" in text:
                                                    result_top = (shape.height + shape.top) / 914400 * 2.54
                                                    # print(result_top)
                                                if "Conclusions" in text:
                                                    conclusion_top = shape.top / 914400 * 2.54
                                        if conclusion_top == '':
                                            conclusion_top = 16.5
                                            bottom_position = float(conclusion_top) / slide_height
                                        else:
                                            bottom_position = float(conclusion_top) / slide_height
                                        image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                        left_position = 0
                                        top_position = float(result_top) / slide_height
                                        # bottom_position = float(conclusion_top) / slide_height
                                        right_position = 1
                                        img = Image.open(image_name)
                                        img_size_width, img_size_height = img.size
                                        crop_left_position = left_position * img_size_width
                                        crop_top_position = top_position * img_size_height
                                        crop_right_position = right_position * img_size_width
                                        crop_bottom_position = bottom_position * img_size_height
                                        name = 'Size exclusion chromatography_' + str(index) + ".jpg"
                                        crop_image_path = os.path.join(image_file_path, name)

                                        crop_img = img.crop(
                                            (crop_left_position, crop_top_position, crop_right_position,
                                             crop_bottom_position))
                                        if crop_img.mode == 'RGBA':
                                            crop_img = crop_img.convert("RGB")
                                        crop_img.save(crop_image_path, quality=95, subsampling=2)
                            for k in list(content[index]['Size-exclusion chromatography'].keys()):
                                if 'index' == k:
                                    Conc_list = []
                                    A260_A280_list = []
                                    Total_amount_list = []
                                    peak_list = []
                                    for i in size_frame_list:
                                        try:
                                            Conc = i.split('Conc.:')[1].split("mg/ml")[0].strip() + " "
                                            Conc_list.append(Conc)
                                            A260_A280 = i.split('A260/A280:')[1].split("Total")[0].strip()
                                            A260_A280_list.append(A260_A280)
                                            Total_amount = i.split("amount:")[1].strip()
                                            Total_amount_list.append(Total_amount)
                                            peak = i.split('Conc.:')[0].strip().replace('\n', ' ')
                                            peak_list.append(peak)
                                        except:
                                            print('检查Size-exclusion chromatography 信息',flush=True)

                                    content[index]['Size-exclusion chromatography']['Conc'] = Conc_list
                                    content[index]['Size-exclusion chromatography']['A260_A280'] = A260_A280_list
                                    content[index]['Size-exclusion chromatography']['Total_amount_value'] = Total_amount_list
                                    content[index]['Size-exclusion chromatography']['Peak_value'] = peak_list
                        if step_name == 'Ion-exchange chromatography'.lower():
                            content[index] = {}
                            content[index]['Ion-exchange chromatography'] = {}
                            frame_list = []
                            for shape in prs.slides[index].shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    content[index]['Ion-exchange chromatography']['index'] = index
                                    ######加上对应的内容
                                    if 'Conc' in text and 'Total amount' in text:
                                        frame_list.append(text)
                                    if 'Conclusions' in text:
                                        try:
                                            conclusions = text.split("Conclusions:")[1].strip()
                                            content[index]['Ion-exchange chromatography']['conclusions'] = conclusions
                                        except:
                                            print('检查Ion-exchange chromatography  结论',flush=True)
                                            conclusions = ''
                                            content[index]['Ion-exchange chromatography']['conclusions'] = conclusions
                                    if 'Experiment' in text:
                                        try:
                                            sample = 'from ' + text.splitlines()[1].split("Sample: From")[1].strip()
                                            content[index]['Ion-exchange chromatography']['sample_value'] = sample
                                        except:
                                            print('检查Ion-exchange chromatography  Sample',flush=True)
                                            sample = ''
                                            content[index]['Ion-exchange chromatography']['sample_value'] = sample

                                        try:
                                            column = text.split('Column:')[1].split('Sample')[0].strip()
                                            content[index]['Ion-exchange chromatography']['column'] = column
                                        except:
                                            print('检查Ion-exchange chromatography  column',flush=True)
                                            column = ''
                                            content[index]['Ion-exchange chromatography']['column'] = column
                                    try:
                                        result_top = ''
                                        conclusion_top = ''
                                        if "Results" in text:
                                            for shape in slide.shapes:

                                                if shape.has_text_frame:
                                                    text_frame = shape.text_frame
                                                    text = text_frame.text
                                                    if "Experiment" in text:
                                                        result_top = (shape.height + shape.top) / 914400 * 2.54

                                                    if "Conclusions" in text:
                                                        conclusion_top = shape.top / 914400 * 2.54
                                            if conclusion_top == '':
                                                conclusion_top = 16.5
                                                bottom_position = float(conclusion_top) / slide_height
                                            else:
                                                bottom_position = float(conclusion_top) / slide_height
                                            image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                            left_position = 0
                                            top_position = float(result_top) / slide_height
                                            # bottom_position = float(conclusion_top) / slide_height
                                            right_position = 1
                                            img = Image.open(image_name)
                                            img_size_width, img_size_height = img.size
                                            crop_left_position = left_position * img_size_width
                                            crop_top_position = top_position * img_size_height
                                            crop_right_position = right_position * img_size_width
                                            crop_bottom_position = bottom_position * img_size_height
                                            name = 'Ion exchange chromatography_' + str(index) + ".jpg"
                                            crop_image_path = os.path.join(image_file_path, name)
                                            crop_img = img.crop(
                                                (crop_left_position, crop_top_position, crop_right_position,
                                                 crop_bottom_position))
                                            if crop_img.mode == 'RGBA':
                                                crop_img = crop_img.convert("RGB")
                                            crop_img.save(crop_image_path, quality=95, subsampling=2)
                                    except:
                                        print("检查Ion exchange chromatography页内容")

                            for k in list(content[index]['Ion-exchange chromatography'].keys()):
                                if 'index' == k:
                                    Conc_list = []
                                    A260_A280_list = []
                                    Total_amount_list = []
                                    peak_list = []
                                    for i in frame_list:
                                        try:
                                            Conc = i.split('Conc.:')[1].split("mg/ml")[0].strip() + " "
                                            Conc_list.append(Conc)
                                            A260_A280 = i.split('A260/A280:')[1].split("Total")[0].strip()
                                            A260_A280_list.append(A260_A280)
                                            Total_amount = i.split("amount:")[1].strip()
                                            Total_amount_list.append(Total_amount)
                                            peak = i.split('Conc.:')[0].strip().replace('\n', ' ')
                                            peak_list.append(peak)
                                        except:
                                            print('检查Ion-exchange chromatography信息', flush=True)
                                    content[index]['Ion-exchange chromatography']['Conc'] = Conc_list
                                    content[index]['Ion-exchange chromatography']['A260_A280'] = A260_A280_list
                                    content[index]['Ion-exchange chromatography'][
                                        'Total_amount_value'] = Total_amount_list
                                    content[index]['Ion-exchange chromatography']['Peak_value'] = peak_list
                        if step_name == 'Diafiltration'.lower():
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'Sample' in text:
                                        content[index] = {}
                                        content[index]['Diafiltration'] = {}
                                        content[index]['Diafiltration']['index'] = index
                                        try:
                                            v = text.split("Sample:")[1].split('L')[0].strip() + ' L'
                                            content[index]['Diafiltration']['v_value'] = v
                                        except:
                                            print('检查Diafiltration 体积')
                                            v = ''
                                            content[index]['Diafiltration']['v_value'] = v
                        if step_name == 'Deadenylation'.lower():
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'Sample' in text:
                                        content[index] = {}
                                        content[index]['Deadenylation'] = {}
                                        content[index]['Deadenylation']['index'] = index
                                        order = index
                                        try:
                                            sample1 = text.split("Sample:")[1].split("Condition")[0].split('from')[
                                                0].strip()
                                            content[index]['Deadenylation']['sample_value1'] = sample1

                                        except:
                                            print('检查Deadenylation  Sample',flush=True)
                                            sample1 = ''
                                            content[index]['Deadenylation']['sample_value1'] = sample1
                                        try:
                                            sample2 = 'from ' +text.split("Sample:")[1].split("Condition")[0].split('from')[1].strip()
                                            content[index]['Deadenylation']['sample_value2'] = sample2
                                        except:
                                            print('检查Deadenylation  Sample',flush=True)
                                            sample2 = ''
                                            content[index]['Deadenylation']['sample_value2'] = sample2
                                        try:
                                            susrate1 = text.split('Subsrate mix:')[1].split('mg/ml')[0].strip()
                                            content[index]['Deadenylation']['susrate1'] = susrate1
                                        except:
                                            print('检查DeadenylationSubsrate mix',flush=True)
                                            susrate1 = ''
                                            content[index]['Deadenylation']['susrate1'] = susrate1
                                        try:

                                            susrate2 = text.split('Subsrate mix:')[1].split(',')[1].split('mg/ml')[
                                                0].strip()
                                            content[index]['Deadenylation']['susrate2'] = susrate2
                                        except:
                                            print('检查DeadenylationSubsrate mix',flush=True)
                                            susrate2 = ''
                                            content[index]['Deadenylation']['susrate2'] = susrate2
                                    try:
                                        result_top = ''
                                        conclusion_top = ''
                                        if "LC-MS results:" in text:
                                            for shape in slide.shapes:
                                                if shape.has_text_frame:
                                                    text_frame = shape.text_frame
                                                    text = text_frame.text
                                                    if "LC-MS results:" in text:
                                                        result_top = (shape.height + shape.top) / 914400 * 2.54
                                                    if "Conclusions" in text:
                                                        conclusion_top = shape.top / 914400 * 2.54
                                            # print(result_top,conclusion_top)
                                            if conclusion_top == '':
                                                conclusion_top = 16.5
                                                bottom_position = float(conclusion_top) / slide_height
                                            else:
                                                bottom_position = float(conclusion_top) / slide_height
                                            image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                            left_position = 0
                                            top_position = float(result_top) / slide_height
                                            # bottom_position = float(conclusion_top) / slide_height
                                            right_position = 1
                                            img = Image.open(image_name)
                                            img_size_width, img_size_height = img.size
                                            crop_left_position = left_position * img_size_width
                                            crop_top_position = top_position * img_size_height
                                            crop_right_position = right_position * img_size_width
                                            crop_bottom_position = bottom_position * img_size_height
                                            name = 'Deadenylation_' + str(order) + ".jpg"
                                            crop_image_path = os.path.join(image_file_path, name)
                                            crop_img = img.crop(
                                                (crop_left_position, crop_top_position, crop_right_position,
                                                 crop_bottom_position))
                                            if crop_img.mode == 'RGBA':
                                                crop_img = crop_img.convert("RGB")
                                            crop_img.save(crop_image_path, quality=95, subsampling=2)
                                    except:
                                        print('检查Deadenylation是否含有图片页',flush=True)
                        if step_name == 'Biotinylation'.lower():
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'Sample' in text:
                                        content[index] = {}
                                        content[index]['Biotinylation'] = {}
                                        content[index]['Biotinylation']['index'] = index
                                        order = index
                                        try:
                                            sample1 = text.split("Sample:")[1].split("Condition")[0].split('from')[
                                                0].strip()
                                            content[index]['Biotinylation']['sample_value1'] = sample1
                                        except:
                                            print('检查Biotinylation sample',flush=True)
                                            sample1 = ''
                                            content[index]['Biotinylation']['sample_value1'] = sample1
                                        try:
                                            sample2 = 'from ' + \
                                                      text.split("Sample:")[1].split("Condition")[0].split('from')[
                                                          1].split(
                                                          '.')[0].strip()
                                            content[index]['Biotinylation']['sample_value2'] = sample2
                                        except:
                                            print('检查Biotinylation sample',flush=True)
                                            sample2 = ''
                                            content[index]['Biotinylation']['sample_value2'] = sample2
                                        try:
                                            Protein_concentration = \
                                            text.split('Protein concentration:')[1].split('mg/ml')[
                                                0].strip()
                                            content[index]['Biotinylation'][
                                                'Protein_concentration'] = Protein_concentration
                                        except:
                                            print('检查Biotinylation Protein_concentration',flush=True)
                                            Protein_concentration = ''
                                            content[index]['Biotinylation'][
                                                'Protein_concentration'] = Protein_concentration
                                    try:
                                        result_top = ''
                                        conclusion_top = ''
                                        if "LC-MS results:" in text:
                                            for shape in slide.shapes:
                                                if shape.has_text_frame:
                                                    text_frame = shape.text_frame
                                                    text = text_frame.text
                                                    if "LC-MS results:" in text:
                                                        result_top = (shape.height + shape.top) / 914400 * 2.54
                                                    if "Conclusions" in text:
                                                        conclusion_top = shape.top / 914400 * 2.54

                                            if conclusion_top == '':
                                                conclusion_top = 16.5
                                                bottom_position = float(conclusion_top) / slide_height
                                            else:
                                                bottom_position = float(conclusion_top) / slide_height
                                            image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                            left_position = 0
                                            top_position = float(result_top) / slide_height
                                            # bottom_position = float(conclusion_top) / slide_height
                                            right_position = 1
                                            img = Image.open(image_name)
                                            img_size_width, img_size_height = img.size
                                            crop_left_position = left_position * img_size_width
                                            crop_top_position = top_position * img_size_height
                                            crop_right_position = right_position * img_size_width
                                            crop_bottom_position = bottom_position * img_size_height
                                            name = 'Biotinylation_' + str(order) + ".jpg"
                                            crop_image_path = os.path.join(image_file_path, name)
                                            crop_img = img.crop(
                                                (crop_left_position, crop_top_position, crop_right_position,
                                                 crop_bottom_position))
                                            if crop_img.mode == 'RGBA':
                                                crop_img = crop_img.convert("RGB")
                                            crop_img.save(crop_image_path, quality=95, subsampling=2)
                                    except:
                                        print('检查Biotinylation内容',flush=True)
                        if step_name == 'Digestion and biotinylation'.lower():
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'Sample' in text:
                                        content[index] = {}
                                        content[index]['Digestion and biotinylation'] = {}
                                        order = index
                                        content[index]['Digestion and biotinylation']['index'] = index
                                        try:
                                            sample1 = text.split("Sample:")[1].split("Condition")[0].split('from')[
                                                0].strip()
                                            content[index]['Digestion and biotinylation']['sample_value1'] = sample1
                                        except:
                                            print('检查Digestion_biotinylation sample',flush=True)
                                            sample1 = ''
                                            content[index]['Digestion and biotinylation']['sample_value1'] = sample1
                                        try:
                                            sample2 = 'from ' + \
                                                      text.split("Sample:")[1].split("Condition")[0].split('from')[
                                                          1].split(
                                                          '.')[0].strip()
                                            content[index]['Digestion and biotinylation']['sample_value2'] = sample2
                                        except:
                                            print('检查Digestion_biotinylation sample',flush=True)
                                            sample2 = ''
                                            content[index]['Digestion and biotinylation']['sample_value2'] = sample2
                                        try:
                                            Protein_concentration = \
                                            text.split('Protein concentration:')[1].split('mg/ml')[
                                                0].strip()
                                            content[index]['Digestion and biotinylation'][
                                                'Protein_concentration'] = Protein_concentration
                                        except:
                                            print('检查Digestion_biotinylation   Protein_concentration',flush=True)
                                            Protein_concentration = ''
                                            content[index]['Digestion and biotinylation'][
                                                'Protein_concentration'] = Protein_concentration
                                        try:
                                            BirA_Protein = \
                                                text.split('Add BirA enzyme to the mixture with a mass ratio of')[
                                                    1].split(';')[
                                                    0].strip()
                                            content[index]['Digestion and biotinylation']['BirA_Protein'] = BirA_Protein
                                        except:
                                            print('检查Digestion_biotinylation  BirA_Protein ',flush=True)
                                            BirA_Protein = ''
                                            content[index]['Digestion and biotinylation']['BirA_Protein'] = BirA_Protein
                                        try:
                                            add_enzyme = text.rpartition('Add')[2].split('enzyme')[0].strip()
                                            content[index]['Digestion and biotinylation']['add_enzyme'] = add_enzyme
                                        except:
                                            add_enzyme = ''
                                            content[index]['Digestion and biotinylation']['add_enzyme'] = add_enzyme
                                        try:
                                            TEVProtein = \
                                                text.rpartition('Add')[2].split('with a mass ratio of')[1].split(
                                                    ': Protein')[
                                                    0].strip()
                                            content[index]['Digestion and biotinylation']['TEVProtein'] = TEVProtein
                                        except:
                                            print('检查Digestion_biotinylation  TEVProtein',flush=True)
                                            TEVProtein = ''
                                            content[index]['Digestion and biotinylation']['TEVProtein'] = TEVProtein
                                    try:
                                        result_top = ''
                                        conclusion_top = ''
                                        if "LC-MS results:" in text:
                                            for shape in slide.shapes:

                                                if shape.has_text_frame:
                                                    text_frame = shape.text_frame
                                                    text = text_frame.text
                                                    if "LC-MS results:" in text:
                                                        result_top = (shape.height + shape.top) / 914400 * 2.54
                                                    if "Conclusions" in text:
                                                        conclusion_top = shape.top / 914400 * 2.54

                                            if conclusion_top == '':
                                                conclusion_top = 16.5
                                                bottom_position = float(conclusion_top) / slide_height
                                            else:
                                                bottom_position = float(conclusion_top) / slide_height
                                            image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                            left_position = 0
                                            top_position = float(result_top) / slide_height
                                            # bottom_position = float(conclusion_top) / slide_height
                                            right_position = 1
                                            img = Image.open(image_name)
                                            img_size_width, img_size_height = img.size
                                            crop_left_position = left_position * img_size_width
                                            crop_top_position = top_position * img_size_height
                                            crop_right_position = right_position * img_size_width
                                            crop_bottom_position = bottom_position * img_size_height
                                            name = 'Digestion biotinylation_' + str(order) + ".jpg"
                                            crop_image_path = os.path.join(image_file_path, name)
                                            crop_img = img.crop(
                                                (crop_left_position, crop_top_position, crop_right_position,
                                                 crop_bottom_position))
                                            if crop_img.mode == 'RGBA':
                                                crop_img = crop_img.convert("RGB")
                                            crop_img.save(crop_image_path, quality=95, subsampling=2)
                                    except:
                                        print('检查Digestion biotinylation内容')
                        if step_name == 'Dephosphorylation'.lower():
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'Sample' in text:
                                        content[index] = {}
                                        content[index]['Dephosphorylation'] = {}
                                        content[index]['Dephosphorylation']['index'] = index
                                        order = index
                                        try:
                                            sample1 = text.split("Sample:")[1].split("Condition")[0].split('from')[
                                                0].strip()
                                            content[index]['Dephosphorylation']['sample_value1'] = sample1
                                        except:
                                            print('检查Dephosphorylation Sample',flush=True)
                                            sample1 = ''
                                            content[index]['Dephosphorylation']['sample_value1'] = sample1
                                        try:
                                            sample2 = 'from ' +text.split("Sample:")[1].split("Condition")[0].split('from')[1].split('.')[0].strip()
                                            content[index]['Dephosphorylation']['sample_value2'] = sample2
                                        except:
                                            print('检查Dephosphorylation Sample',flush=True)
                                            sample2 = ''
                                            content[index]['Dephosphorylation']['sample_value2'] = sample2
                                        try:
                                            Protein_concentration = text.split('Protein concentration:')[1].split('mg/ml')[0].strip()
                                            content[index]['Dephosphorylation']['Protein_concentration'] = Protein_concentration
                                        except:
                                            print('检查Dephosphorylation Protein_concentration',flush=True)
                                            Protein_concentration = ''
                                            content[index]['Dephosphorylation']['Protein_concentration'] = Protein_concentration
                                    try:
                                        result_top = ''
                                        conclusion_top = ''
                                        if "LC-MS results:" in text:
                                            for shape in slide.shapes:

                                                if shape.has_text_frame:
                                                    text_frame = shape.text_frame
                                                    text = text_frame.text
                                                    if "LC-MS results:" in text:
                                                        result_top = (shape.height + shape.top) / 914400 * 2.54
                                                    if "Conclusions" in text:
                                                        conclusion_top = shape.top / 914400 * 2.54
                                            if conclusion_top == '':
                                                conclusion_top = 16.5
                                                bottom_position = float(conclusion_top) / slide_height
                                            else:
                                                bottom_position = float(conclusion_top) / slide_height
                                            image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                            left_position = 0
                                            top_position = float(result_top) / slide_height
                                            # bottom_position = float(conclusion_top) / slide_height
                                            right_position = 1
                                            img = Image.open(image_name)
                                            img_size_width, img_size_height = img.size
                                            crop_left_position = left_position * img_size_width
                                            crop_top_position = top_position * img_size_height
                                            crop_right_position = right_position * img_size_width
                                            crop_bottom_position = bottom_position * img_size_height
                                            name = 'Dephosphorylation_' + str(order) + ".jpg"
                                            crop_image_path = os.path.join(image_file_path, name)
                                            crop_img = img.crop(
                                                (crop_left_position, crop_top_position, crop_right_position,
                                                 crop_bottom_position))
                                            if crop_img.mode == 'RGBA':
                                                crop_img = crop_img.convert("RGB")
                                            crop_img.save(crop_image_path, quality=95, subsampling=2)
                                    except:
                                        print('检查Dephosphorylation 内容')
                        if step_name == 'Denaturation'.lower():
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    # print(index,text)
                                    if 'Procedures' in text:
                                        content[index] = {}
                                        content[index]['Denaturation and refolding'] = {}
                                        content[index]['Denaturation and refolding']['index'] = index
                                        try:
                                            pellet = text.split('pellet of')[1].split(',')[0].strip()
                                            content[index]['Denaturation and refolding']['pellet'] = pellet
                                        except:
                                            print('检查Denaturation Procedures内容',flush=True)
                                        try:
                                            try:
                                                v = text.split("resuspended in")[1].split('Denaturation')[0].strip()
                                                content[index]['Denaturation and refolding']['v'] = v
                                            except:
                                                v = text.split("resuspended in")[1].split('denaturation')[0].strip()
                                                content[index]['Denaturation and refolding']['v'] = v
                                        except:
                                            print('检查Denaturation Procedures内容',flush=True)
                        if step_name == 'Inclusion body preparation'.lower():
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    result_top = ''
                                    conclusion_top = ''
                                    if "Results" in text:
                                        for shape in slide.shapes:
                                            if shape.has_text_frame:
                                                text_frame = shape.text_frame
                                                text = text_frame.text
                                                if "Experiment" in text:
                                                    result_top = (shape.height + shape.top) / 914400 * 2.54
                                                if "Conclusions" in text:
                                                    conclusion_top = shape.top / 914400 * 2.54
                                        try:
                                            if conclusion_top == '':
                                                conclusion_top = 16.5
                                                bottom_position = float(conclusion_top) / slide_height
                                            else:
                                                bottom_position = float(conclusion_top) / slide_height
                                            image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                            left_position = 0
                                            top_position = float(result_top) / slide_height
                                            # bottom_position = float(conclusion_top) / slide_height
                                            right_position = 1
                                            img = Image.open(image_name)
                                            img_size_width, img_size_height = img.size
                                            crop_left_position = left_position * img_size_width
                                            crop_top_position = top_position * img_size_height
                                            crop_right_position = right_position * img_size_width
                                            crop_bottom_position = bottom_position * img_size_height
                                            name = 'Inclusion body preparation' + ".jpg"
                                            crop_image_path = os.path.join(image_file_path, name)
                                            crop_img = img.crop((crop_left_position, crop_top_position,
                                                                 crop_right_position, crop_bottom_position))
                                            if crop_img.mode == 'RGBA':
                                                crop_img = crop_img.convert("RGB")
                                            crop_img.save(crop_image_path, quality=95, subsampling=2)
                                        except:
                                            print('检查Inclusion body preparation内容',flush=True)

        #####获取步骤
        step_dic = {}
        for key, value in content.items():
            for k, v in value.items():
                try:
                    i = int(v['index'])
                    step_dic[i] = k
                except:
                    print('********注意检查'+k+',该步骤没有抓取到相关信息')
                    pass
        step_order = sorted(step_dic.items(), key=lambda x: x[0], reverse=False)
        step_list = []
        for i in step_order:
            step = i[1]
            step_list.append(step)
        # print(content)
        ####获取调整过的序号--content字典
        content_list = []
        content_dic = {}
        try:
            for key, value in content.items():
                content_list.append(value)
            for index, i in enumerate(content_list):
                content_dic[index + 1] = i
        except:
            pass
        try:
            dic = {}
            for i in range(1, len(purification[2])):
                conc = purification[2][i][2].split("mg/ml")[0]
                # print(conc)
                aliquot1 = []
                aliquot2 = []
                for i in purification[2][i][3].splitlines():
                    i = i.strip()
                    #####有的直接是空行有的又有分号
                    if ";" in i:
                        for j in i.split(";"):
                            if j != "":
                                aliquot1.append(j)
                    else:
                        if i != '':
                            aliquot2.append(i)
                aliquot = aliquot1 + aliquot2
                dic[conc] = aliquot
            table_list = []
            for key, value in dic.items():
                for i in value:
                    table_dic = {}
                    EachVolume = i.split('μl/tube')[0].strip()
                    coment = i.split(",")[1].strip()
                    try:
                        tubenum = i.split(",")[1].split("tube")[0].strip()
                    except:
                        tubenum = i.split(",")[1].split("tubes")[0].strip()
                    try:
                        EachAmount_num = float(key) * float(EachVolume)
                        EachAmount = '{:.0f}'.format(EachAmount_num)
                        quanity_num = float(EachAmount_num) * float(tubenum) / 1000
                        quanity = '{:.2f}'.format(quanity_num)
                        table_dic['EachVolume'] = str(EachVolume) + '  μL'
                        table_dic['EachAmount'] = str(EachAmount) + '  μg'
                        table_dic['quanity'] = str(quanity) + ' mg'
                        table_dic['key'] = str(key) + 'mg/ml'
                        table_dic['coment'] = coment
                        table_list.append(table_dic)
                    except:
                        print('检查Storage information表格单位格式是否正确,并手动补充相关内容',flush=True)
                    # print(EachVolume,EachAmount,key,quanity,coment)
            content_dic['tabel_list'] = table_list
        except:
            print('未抓取到Storage information信息，请手动填写',flush=True)
        content_file = os.path.join(image_file_path, 'content.json')

        with open(content_file, "w", encoding="utf-8") as f:
            json.dump(content_dic, f, indent=4, ensure_ascii=False)

        return slide1_dic, ProjectInformation, purification, storageBuffer_list, step_list, content_dic
    def saveoption(self):
        driver = self.driver
        for i in range(20):
            try:
                double_click_element = driver.find_element(By.XPATH,
                                                           r'//*[@id="_eformNaN2019042926.taskdetail"]/table/tbody/tr/td/table[2]/tbody/tr/td[1]/div')
                time.sleep(1)
                actions = ActionChains(driver)
                # 双击需要输入文本的框
                actions.double_click(double_click_element).perform()
                # 在框中输入文本
                active_element = driver.switch_to.active_element
                active_element.send_keys('2 hr')
                driver.switch_to.default_content()

                break

            except:
                time.sleep(5)
        else:

            print('********网络无法响应，无法保存修改内容，该份报告终止上传********', flush=True)
            sys.exit(0)  # Exit the program immediately


        ###点击save
        time.sleep(3)
        for i in range(20):
            try:
                element = WebDriverWait(driver, 10).until(
                    EC.presence_of_element_located(
                        (By.XPATH,
                         r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[1]/td/img[3]'))
                )
                actions = ActionChains(driver)
                for i in range(5):
                    actions.move_to_element(element).perform()
                    actions.click(element).perform()
                    time.sleep(0.1)
                for i in range(5):
                    element.click()
                    time.sleep(0.1)
                time.sleep(2)

                break

            except:
                time.sleep(5)
        else:
            print('********网络无法响应，无法保存修改内容，该份报告终止上传********', flush=True)
            sys.exit(0)  # Exit the program immediately
            # win32api.MessageBox(0, "网络存在问题，请检查", "提醒", win32con.MB_TOPMOST)
    ######根据ppt内容自动上传ELN
    def writetable(self,xpath_name,table_content):
        driver = self.driver
        for i in range(20):
            try:
                double_click_element = driver.find_element(By.XPATH,xpath_name)
                time.sleep(1)
                # 实例化ActionChains对象
                actions = ActionChains(driver)
                # 双击需要输入文本的框
                actions.double_click(double_click_element).perform()
                # 在框中输入文本
                active_element = driver.switch_to.active_element
                active_element.send_keys(table_content)
                break

            except:
                time.sleep(5)
        else:
            pass
    def clickoption(self,xpath_name):
        driver = self.driver

        for i in range(20):
            try:
                WebDriverWait(driver, 10).until(
                    EC.presence_of_element_located(
                        (By.XPATH,xpath_name))).click()
                time.sleep(1)
                break

            except:
                time.sleep(1)
        else:
            pass
    def sendkeys(self, xpath_name,keys):
        driver = self.driver
        for i in range(20):
            try:
                WebDriverWait(driver, 10).until(
                    EC.presence_of_element_located(
                        (By.XPATH, xpath_name))
                ).send_keys(keys)

                time.sleep(1)
                break

            except:
                time.sleep(5)
        else:
            pass
    def changecolor(self,xpath1,xpath2):
        driver = self.driver
        for i in range(20):
            try:
                elements = driver.find_elements(By.XPATH,
                                                xpath1)
                for element in elements:
                    # 设置字体颜色
                    driver.execute_script("arguments[0].style.color = 'black';", element)
                elements1 = driver.find_elements(By.XPATH,
                                                 xpath2)
                for element1 in elements1:
                    driver.execute_script("arguments[0].setAttribute('data-mce-style', arguments[1])", element1,
                                          'color: black;')

                break

            except:
                time.sleep(5)
        else:
            pass
    def insertimage(self,file,xpath1,xpath2):
        driver = self.driver
        # image_file_path = os.path.abspath('image_file_path')
        image_file_path = self.image_file_path
        with open(os.path.join(image_file_path, file), 'rb') as f:
            image_data = f.read()
            # 将图片数据编码为 base64 字符串
        image_base64 = base64.b64encode(image_data).decode('utf-8')
        # 将 base64 字符串插入到 HTML 代码中
        html_str = f'<img src="data:image/jpeg;base64,{image_base64}">'

        image_element = driver.find_element(By.XPATH, xpath1)
        ###删除
        driver.execute_script("arguments[0].remove()", image_element)

        ###添加
        start_offset = 0
        end_offset = 0
        #####图片删除后path没了，要重新看新的xpayh是什么
        element = driver.find_element(By.XPATH,xpath2)
        driver.execute_script("""
                     var range_obj = document.createRange();
                     range_obj.setStart(arguments[0], arguments[1]);
                     range_obj.setEnd(arguments[0], arguments[2]);
                     var fragment = range_obj.createContextualFragment(arguments[3]);
                     range_obj.insertNode(fragment);
                 """, element, start_offset, end_offset, html_str)
    def upload_ELN_out(self):
        slide1_dic, ProjectInformation, purification, storageBuffer_list, step_list, content_dic=self.read_pptfile_out()
        driver = self.driver
        driver.get(self.website)
        driver.maximize_window()
        account=self.account
        password=self.password
        image_file_path = self.image_file_path
        # image_file_path = os.path.abspath('image_file_path')
        ppt_file = self.ppt_file
        name = self.name
    #########登录
        # 等待输入框加载，输入账号

        self.sendkeys(r'//*[@id="__apppoint"]/table/tbody/tr[2]/td/table/tbody/tr[2]/td[2]/input',account+ "\n")

        # 等待输入框加载，输入密码
        self.sendkeys(r'//*[@id="__apppoint"]/table/tbody/tr[2]/td/table/tbody/tr[3]/td[2]/input', password+ "\n")

        print("登录成功",flush=True)

        ####测试用直接点进去
        #####点击My Notebooks

        # self.clickoption('//*[@id="__apppoint"]/table/tbody/tr/td[1]/table/tbody/tr/td/div/table/tbody/tr[3]/td/div/div[1]/div[1]/span')
        # # 点击NBK1022  123
        # self.clickoption('// *[ @ id = "__apppoint"] / table / tbody / tr / td[1] / table / tbody / tr / td / div / table / tbody / tr[3] / td / div / div[1] / div[2] / div[11] / div / span')
        #
        # self.clickoption('// *[ @ id = "__apppoint"] / table / tbody / tr / td[1] / table / tbody / tr / td / div / table / tbody / tr[3] / td / div / div[1] / div[2] / div[17] / div')
        #
        # #####创建功能可以实现,直接点进去,不重复创建,进到NBK1022-21
        # self.clickoption('//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div/table/tbody/tr[2]/td/div/div[1]/div/table/tbody/tr[4]/td[6]/div/a')
        #
        # self.clickoption('//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[1]/img[1]')


        ######创建project
        ######创建项目
        ##点击Home
        try:
            self.clickoption(r'// *[ @ id = "__scil_toolbar_left"] / table / tbody / tr / td[2] / table / tbody / tr[2] / td')
        except  Exception:
            print('********网络无法响应，该份报告终止上传********', flush=True)
            sys.exit(0)
        #####点击New进行创建
        try:
            self.clickoption(r'// *[ @ id = "div"] / table / tbody / tr[3] / td[3] / table / tbody / tr/ td / div / a[1]')
        except  Exception:
            print('********网络无法响应，该份报告终止上传********', flush=True)
            sys.exit(0)

        ######框里添加内容

        self.sendkeys(r'/ html / body / table[4] / tbody / tr[2] / td / div / div / table / tbody / tr[2] / td[2] / input', slide1_dic['ProjectName'])

        self.sendkeys(r'/html/body/table[4]/tbody/tr[2]/td/div/div/table/tbody/tr[3]/td[2]/input', 'Protein purification'+ "\n")

        #####点击创建、关闭
        try:
            self.clickoption(r'/ html / body / table[4] / tbody / tr[2] / td / div / div / table / tbody / tr[5] / td[2] / button')
        except  Exception:
            print('********网络无法响应，该份报告终止上传********', flush=True)
            sys.exit(0)
        time.sleep(3)
        for i in range(20):
            try:
                ok = driver.find_elements(By.XPATH, r"/ html / body / table[6] / tbody / tr[2] / td / div / div / div")
                time.sleep(1)
                break
            except:

                time.sleep(5)
        else:
            print('********网络无法响应，无法点击创建项目，该份报告终止上传********', flush=True)
            sys.exit(0)  # Exit the program immediately
        status = ''
        if len(ok) == 0:
            status = True
            print('项目创建成功',flush=True)
            time.sleep(2)
        if len(ok) != 0:
            oktext=WebDriverWait(driver, 5).until(EC.presence_of_element_located((By.XPATH,r'/ html / body / table[6] / tbody / tr[2] / td / div / div / div[1]'))).text
            if oktext=='[ERROR]: One project is allowed only one notebook for one user':
                status = True
                print('项目已存在',flush=True)
                try:
                    self.clickoption(r'/ html / body / table[6] / tbody / tr[2] / td / div / div / div[2] / button')

                    self.clickoption(r'/ html / body / table[4] / tbody / tr[1] / td[2] / img')

                    self.sendkeys(r'// *[ @ id = "div"] / table / tbody / tr[1] / td / table / tbody / tr[1] / td[1] / input', slide1_dic['ProjectName'])

                    self.clickoption(r'// *[ @ id = "div"] / table / tbody / tr[1] / td / table / tbody / tr[1] / td[2] / span')

                    time.sleep(2)

                    self.clickoption(r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / select[6]')
                    self.clickoption(r' // *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / select[6] / option[4]')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)
                # print('显示100')
                table = driver.find_element(By.XPATH,r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[1] / div / table')
                #####查找项目，先找第一页的并循环刷新多次，如果没有就翻页

                try:
                ######这个查找同时包含名字和项目编号的内容
                    ######先查找第一页是否含有含姓名和项目编号的内容，并刷新循环多次，如果没有的话就来翻页查找
                    for i in range(20):
                        # print('1')
                        row = table.find_element(By.XPATH,r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(slide1_dic['ProjectName'], name))
                        # print(row)
                        # print('2')
                        break
                    else:
                        try:
                            for i in range(20):
                                try:
                                    self.clickoption(r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / img[1]')
                                    # print('3')
                                    row = table.find_element(By.XPATH,r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(slide1_dic['ProjectName'], name))
                                    break
                                except:
                                    time.sleep(2)
                            else:
                                print('********1网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                                sys.exit(0)  # Exit the program immediately
                        except:
                            print('********2网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                            sys.exit(0)  # Exit the program immediately
                except:
                    #####翻页查找内容
                    for i in range(20):
                        try:
                            self.clickoption(r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[2] / div[contains(text(), "Next")]')
                            table = driver.find_element(By.XPATH,r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[1] / div / table')
                            row = table.find_element(By.XPATH,r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(slide1_dic['ProjectName'], name))
                            # print(i,flush=True)
                            break
                        except:
                            time.sleep(2)
                    else:
                        print('********3网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                time.sleep(3)
                for i in range(20):
                    try:
                        WebDriverWait(row, 5).until(
                            EC.presence_of_element_located(
                                (By.XPATH, r'./td[6]'))
                        ).click()
                        print('点击项目成功',flush=True)
                        time.sleep(1)
                        break
                    except:
                        time.sleep(5)
                else:
                    print('********4网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately

            if oktext=='[ERROR]: No privilege to create experiment in the project':
                status = False
                print('********无法创建该项目，请检查项目编号是否正确，该份报告终止上传',flush=True)
                pass
        ####再回到原来的位置
        # driver.execute_script("window.scrollTo(0, 0);")
        if status == True:
            #####创建新表单
            ####新表单名字
            # # 点击加号
            try:
                self.clickoption(r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div/table/tbody/tr[1]/td/span[4]/img')

                #####点击Biortus
                self.clickoption(r'/html/body/table[4]/tbody/tr/td/div/div[3]')

                #######点击纯化  xpath标签总是会变,用last函数
                self.clickoption(r'/html/body/table[last()]/tbody/tr/td/div/div[7]')
            except  Exception:
                print('********网络无法响应，该份报告终止上传********', flush=True)
                sys.exit(0)

            #######part1创建表单填写表格内容
            try:
                NewEnityname=slide1_dic['条目名称']+ ' ' + purification[2][1][1]
            except:
                NewEnityname=' '
                print('注意检查表单名称',flush=True)

            ###表单名称为NewEnityname  修改不了了暂时用这个name吧
            ######输入选项输入创建表单名称
            try:
                self.sendkeys(r'/ html / body / table[6] / tbody / tr[2] / td / div / div / table / tbody / tr[4] / td[2] / input', NewEnityname)
            except:
                self.sendkeys(r'/ html / body / table[6] / tbody / tr[2] / td / div / div / table / tbody / tr[4] / td[2] / input',' ')
                print('检查表单名称',flush=True)

            #####点击创建
            try:
                self.clickoption(r'/html/body/table[6]/tbody/tr[2]/td/div/div/table/tbody/tr[6]/td[2]/button')
            except  Exception:
                print('********网络无法响应，该份报告终止上传********', flush=True)
                sys.exit(0)

            ######填写表格内容
            ####添加title

            try:
                self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[1]/div',slide1_dic['条目名称'])
            except:
                print('检查表格内容')
                pass
            try:
                #####添加Sample ID
                self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[2]/div',slide1_dic['code'])
            except:
                pass
            try:
                Reg=slide1_dic['条目名称'].split("-")[0]
                Batch=slide1_dic['条目名称'].split("-")[1]
            except:
                Reg=''
                Batch=''
            #####添Reg.No(Parent ID)
            self.writetable( r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[3]/div',Reg)
            #####添Batch No(date)
            self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[4]/div', Batch)
            ######print(ProjectInformation[5],ProjectInformation[7],ProjectInformation[9])

            try:
                #####PlasmidName
                self.writetable( r'//*[@id="_eformNaN2019042926.protein1"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[1]/div',slide1_dic['ProteinNname'])
            except:

                pass

            #####MW
            try:
                self.writetable( r'//*[@id="_eformNaN2019042926.protein1"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[2]/div',ProjectInformation[5])
            except:
                pass
            try:
            ######1A280
                self.writetable(r'//*[@id="_eformNaN2019042926.protein1"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[3]/div',ProjectInformation[7])
            except:
                pass
            try:
                ######P1
                self.writetable(r'//*[@id="_eformNaN2019042926.protein1"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[4]/div',ProjectInformation[9])
            except:
                pass
            try:
                #######Cell pelletBatch No
                self.writetable(
                    r'//*[@id="_eformNaN2019042926.protein"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[1]/div',
                    purification[0][1][1])
            except:
                pass

            ########ExpressionVolume
            try:
                Volume = re.search("[1-9]\d*\.?\d*", purification[0][1][2]).group()
            except:
                Volume=' '

            self.writetable(r'//*[@id="_eformNaN2019042926.protein"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[4]/div',Volume + "L")

            ########Cell PelletAmount
            try:
                cellpelletamoutn= re.search("[1-9]\d*\.?\d*", purification[0][1][3]).group()
            except:
                cellpelletamoutn=' '
            self.writetable(r'//*[@id="_eformNaN2019042926.protein"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[5]/div',cellpelletamoutn+ "g")

            try:

                for index, i in enumerate(content_dic['tabel_list']):

                    if index == 0:
                        self.writetable(r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[6]/div',i['key'])
                        if len(storageBuffer_list) != 0:
                            self.writetable(
                                r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[8]/div',storageBuffer_list[0])
                        else:
                            print('缺少Storage Buffer信息',flush=True)
                            pass

                        #####Each Volume(ul)

                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[4]/div',i['EachVolume'])

                        #####comment

                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[9]/div',i['coment'])

                        #####EachAmount

                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[5]/div',i['EachAmount'])

                        #####quantity

                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[7]/div',i['quanity'])

                    else:
                        self.clickoption('//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[1]/td/table/tbody/tr/td[4]/img')
                        self.writetable(r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[6]/div'.format(index+1),i['key'])
                        if len(storageBuffer_list) != 0:
                            self.writetable(
                                r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[8]/div'.format(index+1),storageBuffer_list[0])
                        else:
                            pass

                        #####Each Volume(ul)
                        self.writetable(r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[4]/div'.format(index+1),i['EachVolume'])

                        #####comment
                        self.writetable(r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[9]/div'.format(index+1),i['coment'])

                        #####EachAmount
                        self.writetable(r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[5]/div'.format(index+1),i['EachAmount'])

                        self.writetable(r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[7]/div'.format(index+1),i['quanity'])
            except:
                pass
            ####添加object信息 只要标题的蛋白名称，不要Bp编号和日期
            try:
                iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                driver.switch_to.frame(iframe_list[0])
                #####修改内容
                element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/p')
                driver.execute_script("arguments[0].innerHTML = '{}';".format('purification of '+purification[2][1][1]), element)
                driver.switch_to.default_content()

            except:
                pass
            ######上面如果有错的话，没有回到默认的地方，可能就会定位不到后面的
            driver.switch_to.default_content()
            time.sleep(3)
            print("表格填写完毕",flush=True)
            try:
            #######part2复制 cart并改名字

                self.clickoption(r'//*[@id="__scil_toolbar_right"]/table/tbody/tr/td[2]/table/tbody/tr[1]/td/img')
                ######点击show 加上关键词匹配
                # self.clickoption(r'//*[@id="__scil_toolbar_right"]/table/tbody/tr/td[2]/table/tbody')
                self.clickoption(r'//span[contains(text(), "Show")]')
                #####关闭
                self.clickoption(r' / html / body / table[last()] / tbody / tr[1] / td[2] / img')
                ######再次点击
                self.clickoption(r'//*[@id="__scil_toolbar_right"]/table/tbody/tr/td[2]/table/tbody/tr[1]/td/img')
                ######点击show 加上关键词匹配
                self.clickoption(r'//span[contains(text(), "Show")]')
            except  Exception:
                print('********网络无法响应，该份报告终止上传********', flush=True)
                sys.exit(0)
            #####从页面中选择table
            table= driver.find_element(By.XPATH,r"/html/body/table[last()]/tbody/tr[2]/td/div/div/table/tbody/tr[1]/td[2]/div/div/table")

            #####根据列表元素复制模板
            print(step_list,flush=True)
            for step in step_list:
                for i in range(20):
                    try:
                        row=WebDriverWait(table, 5).until(
                            EC.presence_of_element_located(
                                (By.XPATH, r'.//tr[td[7]="{}"]'.format(step)))
                        )
                        time.sleep(1)
                        break

                    except:
                        time.sleep(5)
                else:
                    print('********1网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately
                time.sleep(2)


                for i in range(20):
                    try:
                        WebDriverWait(row, 5).until(
                            EC.presence_of_element_located(
                                (By.XPATH, r'./td[2]'))
                        ).click()
                        time.sleep(1)
                        break

                    except:
                        time.sleep(5)
                else:
                    print('********2网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately
                time.sleep(2)
                ####复制  滚动到这里
                self.clickoption(r'/html/body/table[last()]/tbody/tr[2]/td/div/div/table/tbody/tr[3]/td[2]/button[1]')
                time.sleep(1)
                for i in range(20):
                    try:
                        WebDriverWait(row, 5).until(
                            EC.presence_of_element_located(
                                (By.XPATH, r'./td[2]'))
                        ).click()
                        time.sleep(1)
                        break

                    except:
                        time.sleep(5)
                else:
                    print('********3网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately
                time.sleep(2)
            ######关闭
            try:
                self.clickoption(r'/html/body/table[last()]/tbody/tr[1]/td[2]/img')
            except  Exception:
                print('********网络无法响应，该份报告终止上传********', flush=True)
                sys.exit(0)
            print("复制模板完毕",flush=True)
            ###复制过来后改名字
            print("根据步骤修改标题",flush=True)
            for i in range(20):
                try:

                    NBKS = driver.find_elements(By.XPATH,
                                                r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[last()]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div/div')

                    time.sleep(1)
                    break

                except:
                    time.sleep(5)
            else:
                print('********网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                sys.exit(0)  # Exit the program immediately

            for j in range(1,len(NBKS)+1):
                try:
                    for i in range(20):
                        try:


                            WebDriverWait(driver, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH, r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[{}]/div/span'.format(j)))
                            ).click()

                            time.sleep(1)
                            break
                        except:
                            time.sleep(5)
                    else:
                        print('********1网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately


                    for i in range(20):
                        try:

                            WebDriverWait(driver, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[1]/td/img[5]'))
                            ).click()
                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********2网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                    for i in range(20):
                        try:
                            WebDriverWait(driver, 10).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     '/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[2]/input'))
                            ).clear()

                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********3网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                    ######文字
                    inner1_list=[]
                    for i in range(20):
                        try:
                            element1=WebDriverWait(driver, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH, r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[{}]/div/span'.format(j)))
                            )
                            if element1.is_displayed():
                                # 获取元素的文本信息
                                text = element1.text
                                # print(text)
                                inner1_list.append(text)
                            else:
                                # 元素不可见，获取元素内的HTML代码
                                inner1 = element1.get_attribute('innerHTML')
                                inner1_list.append(inner1)
                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********4网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                    #####改名字选项
                    ####需要改成什么
                    el=inner1_list[0].split("]")[1].split("- Copy")[0].strip()
                    time.sleep(1)
                    for index, i in enumerate(step_list):
                        if el == i:
                            el1=str(step_list.index(el)+1)+ ' '+inner1_list[0].split("]")[1].split("- Copy")[0].strip()
                            # print(el1)
                            step_list[index] = "matched"
                            # print(step_list)
                            ####用last函数
                            for i in range(20):
                                try:
                                    WebDriverWait(driver, 30).until(
                                        EC.presence_of_element_located(
                                            (By.XPATH,
                                             '/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[2]/input'))
                                    ).send_keys(el1)

                                    time.sleep(2)
                                    break

                                except:
                                    time.sleep(5)
                            else:
                                print('********5网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                                sys.exit(0)  # Exit the program immediately

                                ######save
                            for i in range(20):
                                try:
                                    WebDriverWait(driver, 30).until(
                                        EC.presence_of_element_located(
                                            (By.XPATH,
                                             r'/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[3]/button'))
                                    ).click()
                                    time.sleep(2)
                                    break

                                except:
                                    time.sleep(5)
                            else:
                                print('********6网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                                sys.exit(0)  # Exit the program immediately
                            # 修改元素值
                            # 在下一轮匹配中跳过已匹配的元素

                            break

                        else:
                            pass

                except:
                    pass

            time.sleep(2)
            #####修改完刷新一下
            try:
                self.clickoption(r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[1]/span')
                self.clickoption(r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[1]/td/img[1]')
            except  Exception:
                print('********网络无法响应，该份报告终止上传********', flush=True)
                sys.exit(0)

            time.sleep(2)


            ####part3  修改不同模板内容

            for j in range(1, len(NBKS) + 1):
                self.clickoption(r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[{}]/div/span'.format(j))
                time.sleep(1)
                ######文字
                inner2_list = []
                for i in range(20):
                    try:
                        element1 = WebDriverWait(driver, 5).until(
                            EC.presence_of_element_located(
                                (By.XPATH,
                                 r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[{}]/div/span'.format(
                                     j)))
                        )
                        if element1.is_displayed():
                            # 获取元素的文本信息
                            text = element1.text
                            # print(text)
                            inner2_list.append(text)
                        else:

                            # 元素不可见，获取元素内的HTML代码
                            inner2 = element1.get_attribute('innerHTML')
                            inner2_list.append(inner2)

                        time.sleep(1)
                        break

                    except:
                        time.sleep(5)
                else:
                    print('********7网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately
                #####根据steps进行匹配
                denum = inner2_list[0].split("]")[1].strip()
                print(denum,flush=True)
                for key, value in content_dic.items():
                    try:
                        for k, v in value.items():
                            if denum == str(key) + ' ' + str(k):
                            #######修改Cell lysis & Centrifugation
                                if k== 'Cell lysis & Centrifugation':
                                    time.sleep(2)
                                    iframe_list = driver.find_elements(By.XPATH,r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    print("修改模板Cell lysis & Centrifugation",flush=True)
                                    #####xpath一直变
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[3][contains(text(), "20")]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['m_value']), element)
                                    except:
                                        print('检查Cell lysis & Centrifugation 质量是否正确',flush=True)
                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[5][contains(text(), "5.0 L B")]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['v_value']), element)
                                    except:
                                        print('检查Cell lysis & Centrifugation 体积是否正确', flush=True)
                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[6]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(''),element)
                                    except:

                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[7]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(')'),element)
                                    except:
                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[12][contains(text(), "100 ml")]')
                                        driver.execute_script("arguments[0].innerHTML ='{}';".format(v['procedures']+str('.')), element)
                                    except:
                                        pass

                                    try:
                                        c4=v['ti_value']+" "
                                    except:
                                        print('检查Cell lysis & Centrifugation',flush=True)
                                        c4=' '
                                    element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[4]/span[5][contains(text(), "16000 rpm for 60 min")]')
                                    driver.execute_script("arguments[0].innerHTML = '{}';".format(c4), element)
                                    try:
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[3] / span[8]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['pressure']), element)
                                    except:
                                        pass
                                    try:
                                    #####循环删去内容
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[15][contains(text(), "Buffer ")]')
                                        driver.execute_script("arguments[0].innerHTML = 'Buffer.'", element)
                                        for i in range(13,22):
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[{}]'.format(i))
                                            driver.execute_script("arguments[0].innerHTML = ''", element)
                                    except:
                                        pass
                                    ######修改字体颜色
                                    try:
                                        self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]',r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]')
                                    except:
                                        pass
                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()
                                    print("模板Cell lysis & Centrifugation内容修改成功",flush=True)
                            ####修改Affinity chromatography
                                if k == 'Affinity chromatography (His)' :
                                    time.sleep(2)
                                    self.writetable(
                                        r'//*[@id="_eformNaN2019042926.column"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[2]/div',v['style'])
                                    iframe_list = driver.find_elements(By.XPATH,r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    print("修改模板Affinity chromatography",flush=True)

                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[7][contains(text(), "5 ml")]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['co_value']), element)
                                    except:
                                        print('检查Affinity chromatography',flush=True)
                                        pass
                                    try:
                                        if len(v['Conc'])==0:
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[7]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[6]')
                                            driver.execute_script("arguments[0].remove()", element)
                                        else:
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[7]/span[3][contains(text(), "4.93")]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Conc'][0]), element)

                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[7]/span[6][contains(text(), "0.84")]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['A260_A280'][0]), element)

                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[7]/span[9][contains(text(), "26 mg")]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Total_amount_value'][0]), element)
                                            for i in range(1, len(v['Conc'])):
                                                body = driver.find_element(By.TAG_NAME, 'body')

                                                # 将光标移到文本末尾处
                                                body.send_keys(Keys.END)
                                                body.send_keys(Keys.ENTER)
                                                body.send_keys('——Conc.: ' + v['Conc'][i] + 'mg/ml, A260/A280: ' +
                                                               v['A260_A280'][i] + ', hence, yields ' +
                                                               v['Total_amount_value'][i] + '.')
                                                body.send_keys(Keys.END)
                                    except:
                                        print('检查Affinity chromatography')
                                        pass


                                    # ######修改字体颜色
                                    self.changecolor( r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]',
                                                      r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    time.sleep(2)

                                    driver.switch_to.default_content()
                                    #
                                    ######插入图片部分Results & Discussion:

                                    iframe_list = driver.find_elements(By.XPATH,r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    # print(v,flush=True)
                                    # print(v['conclu'],flush=True)
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div[2]/p[3]/span[2][contains(text(), "Target protein was enriched by His FF columns. A 50-500 ")]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['conclu']), element)
                                    except:
                                        print('检查Affinity chromatography')
                                        pass
                                    time.sleep(5)

                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div[2]/p[3]/span[3][contains(text(), "mM")]')
                                        driver.execute_script("arguments[0].innerHTML = ''", element)
                                    except:
                                        pass
                                    time.sleep(5)
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div[2]/p[3]/span[4][contains(text(), " imidazole linear gradient elution was performed. Samples corresponding to lane 8-12 were collected.")]')
                                        driver.execute_script("arguments[0].innerHTML = ''", element)
                                    except:
                                        pass
                                    time.sleep(5)
                                    #######插入图片
                                    try:
                                        image = 'Affinity chromatography_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/div[2]/p[1]/img', '//*[@id="tinymce"]/div[2]/p[2]')
                                    except:
                                        print('检查Affinity chromatography图片',flush=True)
                                        pass

                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()

                                    print("模板Affinity chromatography内容修改成功",flush=True)
                                if k == 'Affinity chromatography (GST)':
                                    time.sleep(2)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    print("修改模板Affinity chromatography",flush=True)
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[3]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('GST ' + v['co_value']),element)
                                    except:
                                        print('检查Affinity chromatography体积')
                                    try:
                                        if len(v['Conc']) == 0:
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[7]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[6]')
                                            driver.execute_script("arguments[0].remove()", element)

                                        else:
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[7]/span[3]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Conc'][0]),element)

                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[7]/span[6]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['A260_A280'][0]),element)

                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[7]/span[9]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Total_amount_value'][0]), element)

                                            for i in range(1, len(v['Conc'])):
                                                body = driver.find_element(By.TAG_NAME, 'body')

                                                # 将光标移到文本末尾处
                                                body.send_keys(Keys.END)

                                                body.send_keys(Keys.ENTER)
                                                body.send_keys('——Conc.: ' + v['Conc'][i] + 'mg/ml, A260/A280: ' +
                                                               v['A260_A280'][i] + ', hence, yields ' +
                                                               v['Total_amount_value'][i] + '.')
                                                body.send_keys(Keys.END)
                                    except:
                                        pass

                                    # ######修改字体颜色

                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]'
                                                     ,r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    time.sleep(2)

                                    driver.switch_to.default_content()
                                    #
                                    ######插入图片部分Results & Discussion:

                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div[2]/p[3]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['conclu']),element)
                                    except:
                                        print('检查Affinity chromatography结论',flush=True)
                                    element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div[2]/p[3]/span[3]')
                                    driver.execute_script("arguments[0].innerHTML = ''", element)
                                    # 插入图片文件
                                    try:
                                        image = 'Affinity chromatography_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/div[2]/p[1]/img','//*[@id="tinymce"]/div[2]/p[2]')
                                    except:
                                        print('检查Affinity chromatography图片',flush=True)
                                        pass
                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()

                                    print("模板Affinity chromatography内容修改成功",flush=True)
                                if k == 'Affinity chromatography (Flag)':
                                    time.sleep(2)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    print("修改模板Affinity chromatography",flush=True)
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[5]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['co_value']),element)
                                    except:
                                        print('检查Affinity chromatography',flush=True)
                                        pass

                                    try:
                                        if len(v['Conc']) == 0:
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[7]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[6]')
                                            driver.execute_script("arguments[0].remove()", element)
                                        else:
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[7]/span[3][contains(text(), "4.93")]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Conc'][0]),element)

                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[7]/span[6]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['A260_A280'][0]),element)

                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[7]/span[9][contains(text(), "26 mg")]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Total_amount_value'][0]), element)

                                            for i in range(1, len(v['Conc'])):
                                                body = driver.find_element(By.TAG_NAME, 'body')

                                                # 将光标移到文本末尾处
                                                body.send_keys(Keys.END)

                                                body.send_keys(Keys.ENTER)
                                                body.send_keys('——Conc.: ' + v['Conc'][i] + 'mg/ml, A260/A280: ' +
                                                               v['A260_A280'][i] + ', hence, yields ' +
                                                               v['Total_amount_value'][i] + '.')
                                                body.send_keys(Keys.END)
                                    except:
                                        pass

                                    #
                                    # ######修改字体颜色

                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]'
                                    ,r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    time.sleep(2)

                                    driver.switch_to.default_content()
                                    #
                                    ######插入图片部分Results & Discussion:

                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div[2]/p[3]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['conclu']),element)
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[3]/span[3]')
                                        driver.execute_script("arguments[0].innerHTML = ''", element)
                                    except:
                                        pass
                                    # 图片文件
                                    try:
                                        image = 'Affinity chromatography_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/div[2]/p[1]/img','//*[@id="tinymce"]/div[2]/p[2]')
                                    except:
                                        print('检查Affinity chromatography图片',flush=True)
                                        pass
                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()

                                    print("模板Affinity chromatography内容修改成功",flush=True)
                                if k == 'Affinity chromatography (MBP)'  or k == 'Affinity chromatography (Strep)' :
                                    time.sleep(2)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    print("修改模板Affinity chromatography",flush=True)
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[1]/span[4]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['co_value']), element)
                                    except:
                                        pass
                                    try:
                                        if len(v['Conc']) == 0:
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[7]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[6]')
                                            driver.execute_script("arguments[0].remove()", element)
                                        else:
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[7]/span[3]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Conc'][0]),element)

                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[7]/span[6]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['A260_A280'][0]),element)
                                            if k == 'Affinity chromatography (MBP)':
                                                element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[7]/span[10]')
                                                driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Total_amount_value'][0]), element)
                                            if k == 'Affinity chromatography (Strep)':
                                                element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[7]/span[9]')
                                                driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Total_amount_value'][0]), element)

                                            for i in range(1, len(v['Conc'])):
                                                body = driver.find_element(By.TAG_NAME, 'body')

                                                # 将光标移到文本末尾处
                                                body.send_keys(Keys.END)

                                                body.send_keys(Keys.ENTER)
                                                body.send_keys('——Conc.: ' + v['Conc'][i] + 'mg/ml, A260/A280: ' +v['A260_A280'][i] + ', hence, yields ' + v['Total_amount_value'][i] + '.')
                                                body.send_keys(Keys.END)
                                    except:
                                        pass
                                    # ######修改字体颜色

                                    self.changecolor( r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]'
                                                     , r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')
                                    time.sleep(2)

                                    driver.switch_to.default_content()
                                    #
                                    ######插入图片部分Results & Discussion:

                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[3]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['conclu']), element)
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[3]/span[3]')
                                        driver.execute_script("arguments[0].innerHTML = ''", element)
                                    except:
                                        pass

                                    # 插入图片文件
                                    try:
                                        image = 'Affinity chromatography_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/div[2]/p[1]/img','//*[@id="tinymce"]/div[2]/p[2]')
                                    except:
                                        print('检查Affinity chromatography图片',flush=True)
                                        pass

                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()

                                    print("模板Affinity chromatography内容修改成功",flush=True)
                            #####修改Size-exclusion
                                if k == 'Size-exclusion chromatography':
                                    time.sleep(2)
                                    print("修改模板Size-exclusion chromatography",flush=True)
                                    # #####修改表格框的内容

                                    self.writetable(r'//*[@id="_eformNaN2019042926.column"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[2]/div',v['column'])

                                    ##########这个xpath一直变啊
                                    iframe_list = driver.find_elements(By.XPATH,r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])

                                    #####更改内容
                                    ###把第一行内容删了
                                    list=[2,3,4,5]

                                    for i in list:

                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/div/ol/li[1]/span[{}]'.format(i))
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(''), element)
                                    # try:
                                    #     element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/div/ol/li[1]/span[4]')
                                    #     driver.execute_script("arguments[0].innerHTML = '{}';".format(v['sample_value']), element)
                                    # except:
                                    #     pass
                                    #####按照新的模板进行修改内容
                                    try:
                                        content=' was concentrated and loaded onto '+ v['column']+ ' column.'
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/div/ol/li[1]/span[6]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(content),element)
                                    except:
                                        pass
                                    ####修改浓度信息
                                    try:

                                        element = driver.find_element(By.XPATH,'// *[ @ id = "tinymce"] / div / div / ol / li[7] / span')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('It was concentrated to '+purification[2][1][2]+'.'), element)
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/div/ol/li[8]/span[4][contains(text(), "10.5 mg/ml")]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(purification[2][1][2]), element)
                                    except:
                                        print("检查浓度等信息是否正确",flush=True)


                                    # try:
                                    #     element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/div/ol/li[2]/span[6]')
                                    #     driver.execute_script("arguments[0].innerHTML = '{}';".format(v['column'] + ' '), element)
                                    # except:
                                    #     pass
                                    # list=[7,8,9]
                                    # for i in list:
                                    #     element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/div/ol/li[2]/span[{}]'.format(i))
                                    #     driver.execute_script("arguments[0].innerHTML = '{}';".format(''), element)
                                    # try:
                                    #     comment2='It was concentrated to ' + purification[2][1][2]  + '.'
                                    #     element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/div/ol/li[7]/span')
                                    #     driver.execute_script("arguments[0].innerHTML = '{}';".format(comment2), element)
                                    #     element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/div/ol/li[8]/span[4][contains(text(), "10.5 mg/ml")]')
                                    #     driver.execute_script("arguments[0].innerHTML = '{}';".format(purification[2][1][2]), element)
                                    #
                                    # except:
                                    #     pass

                                    try:
                                    ######表格从表格里的内容填写
                                        aliquot1 = []
                                        aliquot2 = []
                                        for i in purification[2][1][3].splitlines():
                                            i = i.strip()
                                            #####有的直接是空行有的又有分号
                                            if ";" in i:

                                                for j in i.split(";"):
                                                    if j != "":
                                                        aliquot1.append(j)
                                            else:
                                                if i != '':
                                                    aliquot2.append(i)
                                        aliquo = aliquot1 + aliquot2

                                        list = []
                                        for i in aliquo:
                                            eachvolume = i.split(",")[0].split("/")[0].strip()
                                            comment = i.split(",")[1].strip()
                                            数字部分 = re.search("[1-9]\d*\.?\d*", eachvolume).group()
                                            # 数字部分 = re.search("^\d+", eachvolume).group()
                                            数字部分整理 = str(float(数字部分) / 1000) + str(' ml ')
                                            s = comment + " x " + 数字部分整理
                                            list.append(s)
                                            # print(s)
                                        # print(list)
                                        l2 = ' and '.join(list)
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/div/ol/li[8]/span[2][contains(text(), "6 tubes x 0.05 ml")]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(l2), element)

                                    except:
                                        pass

                                    try:
                                        if len(v['Conc'])!=0:

                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/div/ol/li[4]/span/span')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Peak_value'][0]), element)
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/div/ol/li[6]/span[3]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format( v['Conc'][0] ), element)
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/div/ol/li[6]/span[6]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['A260_A280'][0]), element)
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/div/ol/li[6]/span[8]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Total_amount_value'][0]),element)
                                            for i in range(1, len(v['Conc'])):
                                                body = driver.find_element(By.TAG_NAME, 'body')
                                                # 将光标移到文本末尾处
                                                body.send_keys(Keys.END)
                                                # 按下回车键换行
                                                body.send_keys(Keys.ENTER)
                                                # 输入新的文本内容
                                                body.send_keys(v['Peak_value'][i] + ' was collected.')
                                                body.send_keys(Keys.END)
                                                # 按下回车键换行
                                                body.send_keys(Keys.ENTER)
                                                body.send_keys('The protein amount was determined using A[280].')
                                                body.send_keys(Keys.END)
                                                body.send_keys(Keys.ENTER)
                                                body.send_keys('——Conc.: ' + v['Conc'][i] + 'mg/ml, A260/A280: ' +v['A260_A280'][i] + ', hence, yields ' +v['Total_amount_value'][i] + '.')
                                                body.send_keys(Keys.END)
                                        else:
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / div / ol / li[6]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / div / ol / li[5]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / div / ol / li[4]')
                                            driver.execute_script("arguments[0].remove()", element)
                                    except:
                                        pass
                                    #####删除第二段
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/div/ol/li[2]')
                                    driver.execute_script("arguments[0].remove()", element)
                                    # ######修改字体颜色
                                    self.changecolor(r'//*[@id="tinymce"]/div/div/ol/li//*[contains(@style, "color: ")]'
                                                     , r'//*[@id ="tinymce"]/div/div/ol/li//*[contains(@data-mce-style, "color: ")]')


                                    driver.switch_to.default_content()

                                    iframe_list = driver.find_elements(By.XPATH,r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])

                                    ######改变截图方式后只需要放一个，把原来的图删掉即可
                                    try:
                                        image = 'Size exclusion chromatography_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/div[4]/p[1]/img','//*[@id="tinymce"]/div[4]/p[1]')
                                    except:
                                        print('检查Size_exclusion_chromatography图片',flush=True)
                                        pass

                                    ######修改结论
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div[4]/p[2]/span[2][contains(text(), "Based on the SDS-PAGE result, peak 1 was collected. ")]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['conclu']), element)
                                    except:
                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[3] / p / span')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('FPLC Chromatography map & SDS-PAGE picture:'),element)
                                    except:
                                        pass
                                    try:
                                        image_element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]')
                                        driver.execute_script("arguments[0].remove()", image_element)
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[1] ')
                                        driver.execute_script("arguments[0].remove()", element)
                                    except:
                                        pass
                                    driver.switch_to.default_content()

                                    time.sleep(1)
                                    self.saveoption()
                                    print("模板Size-exclusion chromatography内容修改成功",flush=True)
                            ######修改QCs
                                if k == 'QC':
                                    time.sleep(2)
                                    print("修改模板QC",flush=True)
                                    iframe_list = driver.find_elements(By.XPATH,r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    #####修改内容
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/p[1]/span[contains(text(), "Batch No.: 20201201-BP9971-1")]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('Batch No.: '+v['Batch']), element)
                                    except:
                                        pass
                                    ######修改字体颜色
                                    try:
                                        self.changecolor(r'//*[@id="tinymce"]/div/p/span[contains(@style, "color: ")]'
                                                         , r'//*[@id ="tinymce"]/div/p/span[contains(@data-mce-style, "color: ")]')
                                    except:
                                        print('检查QC',flush=True)
                                        pass

                                    #####删除QC的内容
                                    try:
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / p[11]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / p[10]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / p[9]')
                                        driver.execute_script("arguments[0].remove()", element)
                                    except:
                                        pass
                                    try:
                                        i_list = []
                                        for i, j in v.items():
                                            if type(i) == int:
                                                # print(i)
                                                i_list.append(i)
                                    except:
                                        pass
                                    # print(i_list)

                                    ###插入zoomin
                                    try:
                                        image = 'QCzoomin_' + str(v['index']) + '_' + str(i_list[0]) +  '.jpg'
                                        # print(image,flush=True)
                                        self.insertimage(image, '//*[@id="tinymce"]/div[1]/p[7]/img','//*[@id="tinymce"]/div[1]/p[7]')
                                    except:
                                        print('检查QC (zoom in) 图片', flush=True)
                                        pass
                                    element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / p[5]')
                                    driver.execute_script("arguments[0].remove()", element)
                                    element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / p[4]')
                                    driver.execute_script("arguments[0].remove()", element)

                                    #####插入第一张图
                                    try:
                                        image = 'QC_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/div[1]/p[3]/img',
                                                         '//*[@id="tinymce"]/div[1]/p[3]')
                                    except:
                                        print('检查QC1,QC2,QC3图片是否有问题',flush=True)
                                        pass
                                    element = driver.find_element(By.XPATH,'// *[ @ id = "tinymce"] / div / p[2] / span')
                                    driver.execute_script("arguments[0].innerHTML = 'QC 1: SDS-PAGE & QC 2: LC-MS  & QC 3: Analytical SEC (Superdex 200 Increase 5/150 GL) QC buffer: 50 mM Tris-HCl (pH 7.5), 300 mM NaCl';",element)
                                    try:
                                        if len(i_list)>1:
                                            # print('需要添加图片')
                                            for i in range(1, len(i_list)):
                                                body = driver.find_element(By.TAG_NAME, 'body')
                                                # 将光标移到文本末尾处
                                                body.send_keys(Keys.END)
                                                # 按下回车键换行
                                                body.send_keys(Keys.ENTER)
                                                image_file_path = self.image_file_path
                                                # image_file_path = os.path.abspath('image_file_path')
                                                file = 'QCzoomin_' + str(v['index']) + '_' + str(i_list[i]) + '.jpg'
                                                # print(file,flush=True)
                                                with open(os.path.join(image_file_path, file), 'rb') as f:
                                                    image_data = f.read()
                                                    # 将图片数据编码为 base64 字符串
                                                image_base64 = base64.b64encode(image_data).decode('utf-8')
                                                # 将 base64 字符串插入到 HTML 代码中
                                                html_str = f'<img src="data:image/jpeg;base64,{image_base64}">'
                                                ###添加
                                                start_offset = 0
                                                end_offset = 0
                                                #####图片删除后path没了，要重新看新的xpayh是什么
                                                element = driver.find_element(By.XPATH,r'//*[@id="tinymce"]/div/p[{}]'.format(int(5 + i)))
                                                driver.execute_script("""
                                                                                                             var range_obj = document.createRange();
                                                                                                             range_obj.setStart(arguments[0], arguments[1]);
                                                                                                             range_obj.setEnd(arguments[0], arguments[2]);
                                                                                                             var fragment = range_obj.createContextualFragment(arguments[3]);
                                                                                                             range_obj.insertNode(fragment);
                                                                                                         """, element,start_offset, end_offset, html_str)
                                    except:
                                        print('检查QC (zoom in) 图片', flush=True)

                                    if len(i_list)==0:
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / p[5]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / p[4]')
                                        driver.execute_script("arguments[0].remove()", element)

                                    driver.switch_to.default_content()
                                    self.saveoption()

                                    print("模板QC内容修改成功",flush=True)
                            ########修改Digestion & Affinity chromatography
                                if k== 'Digestion & Affinity chromatography':
                                    time.sleep(2)
                                    print("修改模板Digestion & Affinity chromatography",flush=True)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[1][contains(text(), "TEV pro")]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['protein']), element)
                                    except:
                                        print('检查Digestion & Affinity chromatography 蛋白',flush=True)
                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH,' // *[ @ id = "tinymce"] / div / ol / li[1] / span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(' was added to the sample ('+v['sample_value']), element)
                                    except:
                                        print('检查Digestion & Affinity chromatography 体积',flush=True)
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[5]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(' ('+v['protein1']+': '),element)
                                    except:
                                        print('检查Digestion & Affinity chromatography 蛋白', flush=True)

                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[4]/span[1]')
                                        u= v['column'].splitlines()[0].split(',')[0].strip()
                                        if 'Ni Bestarose FF'.lower()  in u or 'Talon'.lower()  in u or 'excel'.lower()  in u or 'Protino'.lower()  in u or 'His FF'.lower()  in u or 'Excel'.lower()  in u:
                                            u_style = 'His FF'
                                        elif  'GST'.lower()  in u:
                                            u_style = 'GST'
                                        elif 'MBP'.lower()  in u:
                                            u_style = 'MBP'
                                        elif 'Strep'.lower()  in u:
                                            u_style = 'Strep'
                                        elif 'Flag'.lower()  in u:
                                            u_style = 'Flag'
                                        else:
                                            u_style = 'His FF'
                                        column = u_style + ' (' +v['column'].splitlines()[0].split(',')[1] + ') '
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(column), element)
                                    except:
                                        print('检查Digestion_Affinity_chromatography 柱子类型',flush=True)
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[4]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(''), element)
                                    except:
                                        pass
                                    try:
                                        if len(v['Conc']) == 0:
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[10]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[9]')
                                            driver.execute_script("arguments[0].remove()", element)
                                        else:
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[10]/span[3]/span[contains(text(), "4.36")]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Conc'][0]), element)
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[10]/span[6][contains(text(), "0.84")]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['A260_A280'][0]),element)
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[10]/span[8][contains(text(), "26 mg")]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Total_amount_value'][0]), element)

                                            for i in range(1, len(v['Conc'])):
                                                body = driver.find_element(By.TAG_NAME, 'body')

                                                # 将光标移到文本末尾处
                                                body.send_keys(Keys.END)

                                                body.send_keys(Keys.ENTER)
                                                body.send_keys('——Conc.: ' + v['Conc'][i] + 'mg/ml, A260/A280: ' +
                                                               v['A260_A280'][i] + ', hence, yields ' +
                                                               v['Total_amount_value'][i] + '.')
                                                body.send_keys(Keys.END)
                                    except:
                                        pass

                                    # ######修改字体颜色
                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]'
                                                     , r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    time.sleep(2)

                                    driver.switch_to.default_content()
                                    #
                                    ######插入图片部分Results & Discussion:

                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/p/span[2]')
                                        driver.execute_script( "arguments[0].innerHTML = '{}';".format(v['conclusions']),element)
                                    except:
                                        print('检查Digestion_Affinity_chromatography conclusions', flush=True)
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/p/span[3]')
                                    driver.execute_script("arguments[0].innerHTML = '{}';".format(''),element)
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/p/span[4]')
                                    driver.execute_script("arguments[0].innerHTML = '{}';".format(''), element)


                                    # 插入本地图片文件
                                    try:
                                        image = 'Digestion Affinity chromatography_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/div[2]/p/img',
                                                     '//*[@id="tinymce"]/div[2]/p')
                                    except:
                                        print('检查Digestion & Affinity chromatography图片',flush=True)
                                        pass
                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()
                                    print('模板Digestion & Affinity chromatography内容修改成功',flush=True)
                            ########修改Diafiltration
                                if k == 'Diafiltration':
                                    time.sleep(2)
                                    print('修改模板Diafiltration',flush=True)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')

                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[1]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['v_value']+' '),element)
                                    except:
                                        print('检查Diafiltration 体积',flush=True)
                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]'
                                                     , r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()

                                    print('模板Diafiltration内容修改成功',flush=True)
                            ########修改Deadenylation
                                if k == 'Deadenylation':
                                    time.sleep(2)
                                    print('修改模板Deadenylation',flush=True)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')

                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[1]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['sample_value1'] + ' ('+v['susrate1']+' mg/ml) ' +v['sample_value2']), element)
                                    except:
                                        print('检查Deadenylation', flush=True)
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[3]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['susrate2']+' mg/ml'), element)
                                    except:
                                        print('检查Deadenylation', flush=True)

                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]'
                                                     , r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    driver.switch_to.default_content()
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    # 插入本地图片文件
                                    try:
                                        image = 'Deadenylation_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/p/img',
                                                     '//*[@id="tinymce"]/p')
                                    except:
                                        print('检查Deadenylation图片',flush=True)
                                        pass

                                    driver.switch_to.default_content()

                                    time.sleep(1)
                                    self.saveoption()
                                    print('模板Deadenylation内容修改成功',flush=True)
                            ########修改Biotinylation
                                if k==  'Biotinylation':
                                    time.sleep(2)
                                    print('修改模板Biotinylation',flush=True)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[1]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['sample_value1'] + ' (' + v['Protein_concentration'] + ' mg/ml) ' + v['sample_value2']+'.'), element)
                                    except:
                                        print('检查Biotinylation',flush=True)

                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]', r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    driver.switch_to.default_content()
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    # 插入本地图片文件
                                    try:
                                        image = 'Biotinylation_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/p[1]/img','//*[@id="tinymce"]/p[1]')
                                    except:
                                        print('检查Deadenylation图片',flush=True)
                                        pass
                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()
                                    print('模板Biotinylation内容修改成功',flush=True)
                            ########修改Dephosphorylation
                                if k== 'Dephosphorylation':
                                    time.sleep(2)
                                    print('修改模板Dephosphorylation',flush=True)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[1]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['sample_value1'] + ' (' + v['Protein_concentration'] + ' mg/ml) ' + v['sample_value2']+'.'), element)
                                    except:
                                        print('检查Dephosphorylation',flush=True)

                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]'
                                                     , r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    driver.switch_to.default_content()
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    # 插入本地图片文件
                                    try:
                                        image = 'Dephosphorylation_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/p[1]/img',
                                                     '//*[@id="tinymce"]/p[1]')
                                    except:
                                        print('检查Deadenylation图片',flush=True)
                                        pass

                                    driver.switch_to.default_content()

                                    time.sleep(1)
                                    self.saveoption()
                                    print('模板Dephosphorylation修改成功',flush=True)
                            ########修改Digestion and biotinylation
                                if k==  'Digestion and biotinylation':
                                    time.sleep(2)
                                    print('修改模板Digestion and biotinylation',flush=True)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[1]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['sample_value1'] + ' (' + v['Protein_concentration'] + ' mg/ml) ' + v['sample_value2'] + '.'), element)
                                    except:
                                        print('检查Digestion and biotinylation')
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[1]/span[3]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(''), element)
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[4]/span[5]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format( ' mg/ml) ' + v['BirA_Protein']),element)
                                    except:
                                        print('检查Digestion and biotinylation')

                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[4]/span[6]')
                                    driver.execute_script("arguments[0].innerHTML = '{}';".format(''),element)
                                    try:
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[4] / span[7]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('; '), element)

                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[4] / span[8]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['add_enzyme'] + ' '), element)
                                    except:
                                        print('检查Digestion and biotinylation')
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[4]/span[12]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(' ' + v['TEVProtein']),element)
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[4]/span[13]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(''),element)
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[4]/span[14]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(''),element)
                                    except:
                                        print('检查Digestion and biotinylation')

                                    ###修改字体颜色
                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]'
                                                     , r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    driver.switch_to.default_content()
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    # 插入本地图片文件
                                    try:
                                        image = 'Digestion biotinylation_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/p[1]/img',
                                                     '//*[@id="tinymce"]/p[1]')
                                    except:
                                        print('检查Deadenylation图片',flush=True)
                                        pass
                                    driver.switch_to.default_content()

                                    time.sleep(1)
                                    self.saveoption()
                                    print('模板Digestion and biotinylation内容修改成功',flush=True)
                            ########修改Ion-exchange chromatography
                                if k==  'Ion-exchange chromatography':
                                    time.sleep(2)
                                    self.writetable(r'//*[@id="_eformNaN2019042926.column"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[2]/div',v['column'])
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[1]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('('+v['sample_value']+', ~ 10.00 mg)'), element)
                                    except:
                                        print('检查Ion-exchange chromatography sample内容',flush=True)
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[4]/span[1]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['column']),element)
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[4]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(' '), element)
                                    except:
                                        print('检查Ion-exchange chromatography 体积', flush=True)

                                    try:
                                    # print(Ion_exchange_chromatography['Conc'])
                                        if len(v['Conc']) == 0:
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[9]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[8]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[7]')
                                            driver.execute_script("arguments[0].remove()", element)
                                        else:
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[7]/span/span')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Peak_value'][0]), element)
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[9]/span[2]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format('Conc.: '+ v['Conc'][0]+' mg/ml, '), element)
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[9]/span[4]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['A260_A280'][0]), element)
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[9]/span[6]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Total_amount_value'][0]), element)

                                            for  i in range(1,len(v['Conc'])):

                                                body = driver.find_element(By.TAG_NAME,'body')

                                                # 将光标移到文本末尾处
                                                body.send_keys(Keys.END)

                                                # 按下回车键换行
                                                body.send_keys(Keys.ENTER)

                                                # 输入新的文本内容
                                                body.send_keys(v['Peak_value'][i]+' was collected.')
                                                body.send_keys(Keys.END)

                                                # 按下回车键换行
                                                body.send_keys(Keys.ENTER)
                                                body.send_keys('The amount of protein was determined using A[280].')
                                                body.send_keys(Keys.END)
                                                body.send_keys(Keys.ENTER)
                                                body.send_keys('——Conc.: '+ v['Conc'][i]+'mg/ml, A260/A280: '+ v['A260_A280'][i]+', hence, yields ' + v['Total_amount_value'][i] + '.')
                                                body.send_keys(Keys.END)
                                    except:
                                        pass

                                    # ######修改字体颜色
                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]'
                                                     , r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    time.sleep(2)

                                    driver.switch_to.default_content()

                                    ####先修改图片，不然的话会有问题找不到路径
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])

                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[4]/p[2]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['conclusions']),element)
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[4]/p[2]/span[3]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(''), element)
                                    except:
                                        print('检查Ion-exchange chromatography conclusions', flush=True)

                                    # 插入本地图片文件
                                    try:
                                        image = 'Ion exchange chromatography_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '// *[ @ id = "tinymce"] / div[4] / p[1] / img',
                                                         '//*[@id="tinymce"]/div[4]/p[1]')
                                    except:
                                        print('检查Deadenylation图片', flush=True)
                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[3] / p / span')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('FPLC Chromatography map & SDS-PAGE picture:'),element)
                                        image_element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[2]')
                                        driver.execute_script("arguments[0].remove()", image_element)
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[1]')
                                        driver.execute_script("arguments[0].remove()", element)
                                    except:
                                        pass

                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()

                                    print("模板Ion-exchange chromatography内容修改成功",flush=True)
                                if k==  'Affinity chromatography (HiTrap Heparin HP)':
                                    print('修改模板Affinity chromatography (HiTrap Heparin HP)',flush=True)
                                    time.sleep(2)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    #####更改内容
                                    try:
                                        if v['sample_value']=='':
                                            print('检查sample内容')
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[1]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('('+v['sample_value']+', ~ 10.00 mg)'), element)
                                    except:
                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[4]/span[1]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['column']),element)
                                    except:
                                        pass
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[4]/span[2]')
                                    driver.execute_script("arguments[0].innerHTML = '{}';".format(' '),element)
                                    try:
                                        if len(v['Conc']) == 0:
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[9]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[8]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[7]')
                                            driver.execute_script("arguments[0].remove()", element)
                                        else:
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[7]/span/span')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Peak_value'][0]), element)
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[9]/span[2]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format('Conc.: '+ v['Conc'][0]+' mg/ml, '), element)
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[9]/span[4]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['A260_A280'][0]), element)
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[9]/span[6]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format(v['Total_amount_value'][0]), element)

                                            for  i in range(1,len(v['Conc'])):

                                                body = driver.find_element(By.TAG_NAME,'body')

                                                # 将光标移到文本末尾处
                                                body.send_keys(Keys.END)

                                                # 按下回车键换行
                                                body.send_keys(Keys.ENTER)

                                                # 输入新的文本内容
                                                body.send_keys(v['Peak_value'][i]+' was collected.')
                                                body.send_keys(Keys.END)

                                                # 按下回车键换行
                                                body.send_keys(Keys.ENTER)
                                                body.send_keys('The amount of protein was determined using A[280].')
                                                body.send_keys(Keys.END)
                                                body.send_keys(Keys.ENTER)
                                                body.send_keys('——Conc.: '+ v['Conc'][i]+'mg/ml, A260/A280: '+ v['A260_A280'][i]+', hence, yields ' + v['Total_amount_value'][i] + '.')
                                                body.send_keys(Keys.END)
                                    except:
                                        pass
                                    # ######修改字体颜色
                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]'
                                                     , r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')
                                    time.sleep(2)
                                    driver.switch_to.default_content()

                                    ####先修改图片
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[4]/p[2]/span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['conclusions']), element)
                                    except:
                                        print("检查Affinity chromatography (HiTrap Heparin HP) conclusions",flush=True)
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[4]/p[2]/span[3]')
                                    driver.execute_script("arguments[0].innerHTML = '{}';".format(''), element)

                                    # 插入本地图片文件
                                    try:
                                        image = 'Affinity chromatography亲和_' + str(v['index']) + '.jpg'
                                        self.insertimage(image, '// *[ @ id = "tinymce"] / div[4] / p[1] / img',
                                                     '//*[@id="tinymce"]/div[4]/p[1]')
                                    except:
                                        print('检查Deadenylation图片',flush=True)
                                        pass
                                        ###删除
                                    try:
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[3] / p / span')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('FPLC Chromatography map & SDS-PAGE picture:'),element)
                                        image_element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]')
                                        driver.execute_script("arguments[0].remove()", image_element)
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[1]')
                                        driver.execute_script("arguments[0].remove()", element)
                                    except:
                                        pass
                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()

                                    print("模板Affinity chromatography (HiTrap Heparin HP) 内容修改成功",flush=True)
                                if k== 'Denaturation and refolding':
                                    time.sleep(2)
                                    print('修改模板Denaturation and refolding', flush=True)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')

                                    driver.switch_to.frame(iframe_list[1])
                                    try:
                                    #####更改内容
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div / ol / li[1] / span[2]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['pellet']), element)
                                    except:
                                        print("检查Denaturation and refolding",flush=True)
                                    try:
                                        element = driver.find_element(By.XPATH, ' // *[ @ id = "tinymce"] / div / ol / li[1] / span[8]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['v']+' '), element)
                                    except:
                                        print("检查Denaturation and refolding", flush=True)
                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]', r'//*[@id ="tinymce"]/div/ol/li//*[contains(@data-mce-style, "color: ")]')
                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()
                                    print('模板Denaturation and refolding内容修改成功', flush=True)
                                if k== 'Cell lysis & Inclusion body preparation':
                                    time.sleep(2)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    print("修改模板Cell lysis & Inclusion body preparation", flush=True)
                                    #####xpath一直变
                                    try:
                                        element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/div/ol/li[1]/span[1]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['m_value']+' g '), element)
                                    except:
                                        print("检查Cell lysis & Inclusion body preparation", flush=True)
                                    try:
                                        volume=v['procedures'].split('with')[1].split('of lysis')[0].strip()+' '
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[1]/span[6]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(volume), element)
                                    except:
                                        print('检查Cell lysis & Inclusion body preparation Procedures体积')

                                    try:
                                        c4 = v['ti_value'] + " "
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[4]/span[5]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(c4), element)
                                    except:

                                        print('检查Cell lysis & Inclusion body preparation Procedures 时间', flush=True)

                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[3]/span[8]')
                                    driver.execute_script("arguments[0].innerHTML = '{}';".format(v['pressure']), element)

                                    ######修改字体颜色

                                    self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]',
                                                     r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]')
                                    time.sleep(2)
                                    driver.switch_to.default_content()
                                    ####先修改图片，不然的话会有问题找不到路径
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    # 插入本地图片文件
                                    try:
                                        image = 'Inclusion body preparation' + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/p/span/img',
                                                         '//*[@id="tinymce"]/p/span')
                                    except:
                                        print('检查Cell lysis & Inclusion body preparation图片', flush=True)
                                        pass
                                    time.sleep(1)

                                    driver.switch_to.default_content()
                                    self.saveoption()

                                    print("模板Cell lysis & Inclusion body preparation内容修改成功", flush=True)
                    except:
                        pass
            print("完成修改模板内容任务",flush=True)
            for i in range(20):
                try:
                    WebDriverWait(driver, 5).until(
                        EC.presence_of_element_located(
                            (By.XPATH,
                             '//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[last()]/table/tbody/tr[1]/td/span[3]/img'))
                    ).click()
                    time.sleep(1)
                    break

                except:
                    time.sleep(5)
            else:
                print('********1网络无法响应，上传ppt失败，该份报告终止上传********', flush=True)
                sys.exit(0)  # Exit the program immediately

            for i in range(20):
                try:
                    WebDriverWait(driver, 10).until(
                        EC.presence_of_element_located(
                            (By.XPATH,
                             '/html/body/table[last()]/tbody/tr/td/div/div[2]/table/tbody/tr/td[1]/span'))
                    ).click()
                    time.sleep(1)
                    break

                except:
                    time.sleep(5)
            else:
                print('********2网络无法响应，上传ppt失败，该份报告终止上传********', flush=True)
                sys.exit(0)  # Exit the program immediately

            for i in range(20):
                try:
                    WebDriverWait(driver, 10).until(
                        EC.presence_of_element_located(
                            (By.XPATH,
                             '/html/body/table[last()]/tbody/tr/td/div/div[6]/span'))
                    ).click()
                    time.sleep(1)
                    break

                except:
                    time.sleep(5)
            else:
                print('********3网络无法响应，上传ppt失败，该份报告终止上传********', flush=True)
                sys.exit(0)  # Exit the program immediately

            file_dialog = driver.find_element(By.XPATH, r'//input[@type="file"]')
            time.sleep(2)

            file_dialog.send_keys(ppt_file)

            for i in range(20):
                try:
                    WebDriverWait(driver, 5).until(
                        EC.presence_of_element_located(
                            (By.XPATH,
                             '/ html / body / table[last()] / tbody / tr[2] / td / div / div / div / form / table / tbody / tr[7] / td[2] / button'))
                    ).click()
                    time.sleep(1)
                    break

                except:
                    time.sleep(5)
            else:
                print('********4网络无法响应，上传ppt失败，该份报告终止上传********', flush=True)
                sys.exit(0)  # Exit the program immediately
            print("pptx文件上传成功",flush=True)
            time.sleep(5)
        else:
            pass
    def read_pptfile_gallery(self):
        print("提取ppt信息", flush=True)
        ppt_file = self.ppt_file
        # image_file_path = os.path.abspath('image_file_path')
        image_file_path = self.image_file_path
        utils.save_pptx_as_png(image_file_path, ppt_file, overwrite_folder=True)
        prs = Presentation(ppt_file)
        slide1_dic = {}
        purification = []
        storageBuffer_list = []
        content = {}
        for index, slide in enumerate(prs.slides):
            sortedShapes = sorted(slide.shapes, key=lambda x: (x.left))
            slide_height = prs.slide_height / 914400 * 2.54
            for shape in sortedShapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text = text_frame.text
                    if 'ELN' in text:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame
                                # 读取文本框内容
                                text = text_frame.text
                                if 'ELN' in text:
                                    try:
                                        try:
                                            ProjectName = text.split(":")[1].strip()
                                            slide1_dic['ProjectName'] = ProjectName
                                        except:
                                            ProjectName = text.split("：")[1].strip()
                                            slide1_dic['ProjectName'] = ProjectName
                                    except:
                                        print('检查项目编号处内容', flush=True)
                                        slide1_dic['ProjectName'] = ''

                                if len(text.splitlines()) == 3:
                                    try:
                                        code = text.splitlines()[0].split(",")[0]
                                        slide1_dic['code'] = code
                                    except:
                                        slide1_dic['code'] = ''
                                    try:
                                        条目名称 = text.splitlines()[0].split(",")[1].strip()
                                        slide1_dic['条目名称'] = 条目名称
                                    except:
                                        slide1_dic['条目名称'] = ''
                                    try:
                                        NotebookName = 'Protein purification'
                                        slide1_dic['NotebookName'] = NotebookName
                                    except:
                                        slide1_dic['NotebookName'] = ''
                                    try:
                                        ProteinNname = text.splitlines()[1].strip()
                                        slide1_dic['ProteinNname'] = ProteinNname
                                    except:
                                        slide1_dic['ProteinNname'] = ''
                    if "Project Information" in text:
                        sortedShapes = sorted(slide.shapes, key=lambda x: (x.top))
                        for shape in sortedShapes:
                            if shape.has_table:
                                # 获取表格对象
                                table = shape.table
                                # 创建一个列表用于保存表格内容
                                table_data = []
                                # 遍历表格中的所有行和列
                                for i, row in enumerate(table.rows):
                                    row_data = []
                                    for j, cell in enumerate(row.cells):
                                        # 获取单元格中的文本
                                        cell_text = cell.text_frame.text.strip()
                                        row_data.append(cell_text)
                                    # 将该行添加到表格内容列表中
                                    table_data.append(row_data)
                                purification.append(table_data)
                    if 'Purification Scheme' in text:
                        sortedShapes = sorted(slide.shapes, key=lambda x: (x.left))
                        text_list = []
                        for shape in sortedShapes:
                            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                                # 遍历组合图中的所有形状
                                for s in shape.shapes:
                                    if hasattr(s, 'text') and s.text.strip():
                                        # 输出每个形状的文本内容
                                        # print(s.text)
                                        s = s.text.splitlines()
                                        text_list.append(s)
                        # Purification_Scheme['Purification_Scheme']=text_list
                    parts = re.split("Step \d*", text)
                    if len(parts) > 1:
                        step_name = parts[1].split(":")[1].strip()
                        step_name = step_name.lower()
                        # print(step_name)
                        if step_name == "Cell lysis & Centrifugation".lower():
                            try:
                                for shape in prs.slides[index + 1].shapes:
                                    if shape.has_text_frame:
                                        text_frame = shape.text_frame
                                        text = text_frame.text
                                        if 'Step' in text:
                                            if 'Inclusion body preparation' in text:
                                                name = 'Cell lysis & Inclusion body preparation'
                                            else:
                                                name = 'Cell lysis & Centrifugation'
                            except:
                                name = 'Cell lysis & Centrifugation'

                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'Sample' in text:
                                        content[index] = {}
                                        content[index][name] = {}
                                        content[index][name]['index'] = index
                                        try:
                                            m = text.split(":")[1].split("g")[0].strip()
                                            content[index][name]['m_value'] = m
                                        except:
                                            print('检查Cell lysis & Centrifugation 质量', flush=True)
                                            m = ''
                                            content[index][name]['m_value'] = m
                                        try:
                                            v = text.splitlines()[0].split("from")[1].strip()
                                            content[index][name]['v_value'] = v
                                        except:
                                            print('检查Cell lysis & Centrifugation 体积', flush=True)
                                            v = ''
                                            content[index][name]['v_value'] = v
                                        try:
                                            procedures = text.split("was resuspended")[1].split(".")[0].strip()
                                            content[index][name]['procedures'] = procedures
                                        except:
                                            procedures = ''
                                            content[index][name]['procedures'] = procedures
                                            print('检查Cell lysis & Centrifugation Procedures', flush=True)
                                    if 'high pressure' in text:
                                        try:
                                            pressure = text.split('pressure homogenizer (')[1].split('Bar')[0]
                                            content[index][name]['pressure'] = pressure
                                        except:
                                            print('检查Cell lysis & Centrifugation 压力', flush=True)
                                            pressure = ''
                                            content[index][name]['pressure'] = pressure
                                    if 'Lysis buffer' in text:
                                        try:
                                            buff = text.split("mM")[0].split("Lysis buffer: ")[1].strip() + " ml"
                                            content[index][name]['buff_value'] = buff
                                        except:
                                            print('检查Cell lysis & Centrifugation Lysis buffer', flush=True)
                                            buff = ''
                                            content[index][name]['buff_value'] = buff
                                    if 'min' in text:
                                        try:
                                            ti = text.split("min at")[0].split("Centrifuged the lysate  with ")[
                                                     1].strip() + " min"
                                            content[index][name]['ti_value'] = ti
                                        except:
                                            print('检查Cell lysis & Centrifugation  时间', flush=True)
                                            ti = ''
                                            content[index][name]['ti_value'] = ti
                        if step_name == 'Affinity & SEC'.lower():
                            content[index] = {}
                            content[index]['Gallery protocol (AC-AC-SEC)'] = {}
                            content[index]['Gallery protocol (AC-AC-SEC)']['index'] = index

                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    # 读取文本框内容
                                    text = text_frame.text
                                    try:
                                        result_top = ''
                                        conclusion_top = ''
                                        if "Elution results" in text:
                                            for shape in slide.shapes:
                                                if shape.has_text_frame:
                                                    text_frame = shape.text_frame
                                                    text = text_frame.text
                                                    if "Elution results" in text:
                                                        result_top = (shape.height + shape.top) / 914400 * 2.54
                                                    if "Conclusions" in text:
                                                        conclusion_top = shape.top / 914400 * 2.54
                                                        try:
                                                            try:
                                                                peak = text.split('Conclusions:')[1].split('were')[0].strip()
                                                                content[index]['Gallery protocol (AC-AC-SEC)'][
                                                                    'peak'] = peak
                                                            except:
                                                                peak = text.split('Conclusions:')[1].split('was')[
                                                                    0].strip()
                                                                content[index]['Gallery protocol (AC-AC-SEC)'][
                                                                    'peak'] = peak
                                                        except:
                                                            content[index]['Gallery protocol (AC-AC-SEC)']['peak'] = ''
                                                            print('检查Affinity & SEC conclusion', flush=True)
                                            # print(result_top,conclusion_top)
                                            if conclusion_top == '':
                                                conclusion_top = 16.5
                                                bottom_position = float(conclusion_top) / slide_height
                                            else:
                                                bottom_position = float(conclusion_top) / slide_height
                                            image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                            left_position = 0
                                            top_position = float(result_top) / slide_height
                                            # bottom_position = float(conclusion_top) / slide_height
                                            right_position = 1
                                            img = Image.open(image_name)
                                            img_size_width, img_size_height = img.size
                                            crop_left_position = left_position * img_size_width
                                            crop_top_position = top_position * img_size_height
                                            crop_right_position = right_position * img_size_width
                                            crop_bottom_position = bottom_position * img_size_height
                                            name = 'Affinity & SEC_' + str(index) + ".jpg"
                                            crop_image_path = os.path.join(image_file_path, name)
                                            crop_img = img.crop(
                                                (crop_left_position, crop_top_position, crop_right_position,
                                                 crop_bottom_position))
                                            if crop_img.mode == 'RGBA':
                                                crop_img = crop_img.convert("RGB")
                                            crop_img.save(crop_image_path, quality=95, subsampling=2)
                                    except:
                                        print('检查Affinity & SEC', flush=True)
                        if step_name == 'Affinity & Digestion & SEC'.lower():
                            content[index] = {}
                            content[index]['Gallery protocol (AC-Digestion-AC-SEC)'] = {}
                            content[index]['Gallery protocol (AC-Digestion-AC-SEC)']['index'] = index
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    # 读取文本框内容
                                    text = text_frame.text
                                    try:
                                        result_top = ''
                                        conclusion_top = ''
                                        if "Elution results" in text:
                                            for shape in slide.shapes:
                                                if shape.has_text_frame:
                                                    text_frame = shape.text_frame
                                                    text = text_frame.text
                                                    if "Elution results" in text:
                                                        result_top = (shape.height + shape.top) / 914400 * 2.54
                                                    if "Conclusions" in text:
                                                        conclusion_top = shape.top / 914400 * 2.54
                                                        try:
                                                            try:
                                                                peak = text.split('Conclusions:')[1].split('were')[
                                                                    0].strip()
                                                                content[index][
                                                                    'Gallery protocol (AC-Digestion-AC-SEC)'][
                                                                    'peak'] = peak
                                                            except:
                                                                peak = text.split('Conclusions:')[1].split('was')[
                                                                    0].strip()
                                                                content[index][
                                                                    'Gallery protocol (AC-Digestion-AC-SEC)'][
                                                                    'peak'] = peak
                                                        except:
                                                            content[index]['Gallery protocol (AC-Digestion-AC-SEC)'][
                                                                'peak'] = ' '
                                                            print('检查Affinity & SEC conclusion', flush=True)
                                            # print(result_top,conclusion_top)
                                            if conclusion_top == '':
                                                conclusion_top = 16.5
                                                bottom_position = float(conclusion_top) / slide_height
                                            else:
                                                bottom_position = float(conclusion_top) / slide_height
                                            image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                            left_position = 0
                                            top_position = float(result_top) / slide_height
                                            # bottom_position = float(conclusion_top) / slide_height
                                            right_position = 1
                                            img = Image.open(image_name)
                                            img_size_width, img_size_height = img.size
                                            crop_left_position = left_position * img_size_width
                                            crop_top_position = top_position * img_size_height
                                            crop_right_position = right_position * img_size_width
                                            crop_bottom_position = bottom_position * img_size_height
                                            name = 'Affinity & Digestion & SEC_' + str(index) + ".jpg"
                                            crop_image_path = os.path.join(image_file_path, name)
                                            crop_img = img.crop(
                                                (crop_left_position, crop_top_position, crop_right_position,
                                                 crop_bottom_position))
                                            if crop_img.mode == 'RGBA':
                                                crop_img = crop_img.convert("RGB")
                                            crop_img.save(crop_image_path, quality=95, subsampling=2)
                                    except:
                                        print('检查Affinity & Digestion & SEC', flush=True)
                        if step_name == 'QCs'.lower():
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    text_frame = shape.text_frame
                                    text = text_frame.text
                                    if 'QC1: SDS-PAGE' in text:
                                        content['QC'] = {}
                                        content['QC']['index'] = index
                                        order = index
                                        for shape in slide.shapes:
                                            if shape.has_text_frame:
                                                text_frame = shape.text_frame
                                                text = text_frame.text
                                                if 'Batch No' in text:
                                                    try:
                                                        content['QC']['Batch'] = text.split(':')[1].strip()
                                                        top = (shape.top + shape.height) / 914400 * 2.54
                                                        left_position = 0
                                                        top_position = top / slide_height
                                                        bottom = 16.5
                                                        bottom_position = bottom / slide_height
                                                        right_position = 1
                                                        image_name = os.path.join(image_file_path,
                                                                                  "幻灯片{}.PNG".format(index + 1))
                                                        img = Image.open(image_name)
                                                        img_size_width, img_size_height = img.size
                                                        crop_left_position = left_position * img_size_width
                                                        crop_top_position = top_position * img_size_height
                                                        crop_right_position = right_position * img_size_width
                                                        crop_bottom_position = bottom_position * img_size_height
                                                        name = 'QC' + '_' + str(order) + ".jpg"
                                                        crop_image_path = os.path.join(image_file_path, name)
                                                        crop_img = img.crop(
                                                            (crop_left_position, crop_top_position, crop_right_position,
                                                             crop_bottom_position))
                                                        if crop_img.mode == 'RGBA':
                                                            crop_img = crop_img.convert("RGB")
                                                        crop_img.save(crop_image_path, quality=95, subsampling=2)
                                                    except:
                                                        print('检查QC图片是否有问题')
                                                        pass
                                    ####截QC2的图
                                    if 'QC2: LC-MS (zoom in)' in text:
                                        content['QC'][index] = index
                                        if str(index - 1) != str(order):
                                            print('QC2: LC-MS (zoom in)包含多张图片，请注意检查', flush=True)
                                        for shape in prs.slides[index].shapes:
                                            if shape.has_text_frame:
                                                text_frame = shape.text_frame
                                                # 读取文本框内容
                                                text = text_frame.text
                                                if "Conclusions" in text:
                                                    try:
                                                        storageBuffer = \
                                                            text.split("buffer:")[1].strip().rpartition('.')[0]
                                                        storageBuffer_list.append(storageBuffer)
                                                    except:
                                                        print('检查QC Conclusions', flush=True)
                                                        storageBuffer = ''
                                                        storageBuffer_list.append(storageBuffer)
                                                try:
                                                    result_top = ''
                                                    conclusion_top = ''
                                                    if "zoom in" in text:
                                                        for shape in slide.shapes:
                                                            if shape.has_text_frame:
                                                                text_frame = shape.text_frame
                                                                text = text_frame.text
                                                                if "zoom in" in text:
                                                                    result_top = (shape.height + shape.top) / 914400 * 2.54
                                                                    # print(result_top)
                                                                if "Conclusions" in text:
                                                                    conclusion_top = shape.top / 914400 * 2.54
                                                        image_name = os.path.join(image_file_path,
                                                                                  "幻灯片{}.PNG".format(index + 1))

                                                        if conclusion_top == '':
                                                            conclusion_top = 16.5
                                                            bottom_position = float(conclusion_top) / slide_height
                                                        else:
                                                            bottom_position = float(conclusion_top) / slide_height

                                                        left_position = 0
                                                        top_position = float(result_top) / slide_height
                                                        # bottom_position = float(conclusion_top) / slide_height
                                                        right_position = 1
                                                        img = Image.open(image_name)
                                                        img_size_width, img_size_height = img.size
                                                        crop_left_position = left_position * img_size_width
                                                        crop_top_position = top_position * img_size_height
                                                        crop_right_position = right_position * img_size_width
                                                        crop_bottom_position = bottom_position * img_size_height
                                                        name = 'QCzoomin' + '_' + str(order) + '_' + str(index) + ".jpg"
                                                        crop_image_path = os.path.join(image_file_path, name)
                                                        # print(crop_image_path)
                                                        crop_img = img.crop(
                                                            (crop_left_position, crop_top_position, crop_right_position,
                                                             crop_bottom_position))
                                                        if crop_img.mode == 'RGBA':
                                                            crop_img = crop_img.convert("RGB")
                                                        crop_img.save(crop_image_path, quality=95, subsampling=2)
                                                except:
                                                    pass
                                                    print('检查QCzoom格式是否正确', flush=True)
        #####获取步骤
        step_dic = {}
        ####获取step_list
        for key, value in content.items():
            for k, v in value.items():
                try:
                    i = int(v['index'])
                    step_dic[i] = k
                except:
                    pass
        step_order = sorted(step_dic.items(), key=lambda x: x[0], reverse=False)
        step_list = []
        for i in step_order:
            step = i[1]
            step_list.append(step)
        ####获取调整过的序号--content字典
        # print(content)
        int_keys_sorted = sorted([k for k in content.keys() if isinstance(k, int)])
        # 遍历列表，更改字典的key
        for i, k in enumerate(int_keys_sorted):
            content[i + 1] = content.pop(k)
        # print(content)
        dic = {}
        try:
            for i in range(1, len(purification[1])):
                conc = purification[1][i][2].split("mg/ml")[0]
                # print(conc)
                aliquot1 = []
                aliquot2 = []
                for i in purification[1][i][3].splitlines():
                    i = i.strip()
                    #####有的直接是空行有的又有分号
                    if ";" in i:
                        for j in i.split(";"):
                            if j != "":
                                aliquot1.append(j)
                    else:
                        if i != '':
                            aliquot2.append(i)
                aliquot = aliquot1 + aliquot2
                dic[conc] = aliquot
            table_list = []
            for key, value in dic.items():
                for i in value:
                    table_dic = {}
                    EachVolume = i.split('μl/tube')[0].strip()
                    coment = i.split(",")[1].strip()
                    try:
                        tubenum = i.split(",")[1].split("tube")[0].strip()
                    except:
                        tubenum = i.split(",")[1].split("tubes")[0].strip()
                    try:
                        EachAmount_num = float(key) * float(EachVolume)
                        EachAmount = '{:.0f}'.format(EachAmount_num)
                        quanity_num = float(EachAmount_num) * float(tubenum) / 1000
                        quanity = '{:.2f}'.format(quanity_num)
                        table_dic['EachVolume'] = str(EachVolume) + '  μL'
                        table_dic['EachAmount'] = str(EachAmount) + '  μg'
                        table_dic['quanity'] = str(quanity) + ' mg'
                        table_dic['key'] = str(key) + 'mg/ml'
                        table_dic['coment'] = coment
                        table_list.append(table_dic)
                    except:
                        print('检查Storage information表格单位格式是否正确', flush=True)

            content['tabel_list'] = table_list
        except:
            pass
        content_dic = content
        content_file = os.path.join(image_file_path, 'content.json')
        with open(content_file, "w", encoding="utf-8") as f:
            json.dump(content_dic, f, indent=4, ensure_ascii=False)
        #print(content_dic)
        return slide1_dic, purification, storageBuffer_list, step_list, content_dic, text_list
    def upload_ELN_gallery(self):
        slide1_dic, purification, storageBuffer_list, step_list, content_dic, text_list = self.read_pptfile_gallery()
        driver = self.driver
        driver.get(self.website)
        driver.maximize_window()
        account = self.account
        password = self.password
        image_file_path = self.image_file_path
        # image_file_path = os.path.abspath('image_file_path')
        ppt_file = self.ppt_file
        name = self.name
        #########登录
        # 等待输入框加载，输入账号

        self.sendkeys(r'//*[@id="__apppoint"]/table/tbody/tr[2]/td/table/tbody/tr[2]/td[2]/input', account + "\n")

        # 等待输入框加载，输入密码
        self.sendkeys(r'//*[@id="__apppoint"]/table/tbody/tr[2]/td/table/tbody/tr[3]/td[2]/input', password + "\n")

        print("登录成功", flush=True)

        ####测试用直接点进去
        #####点击My Notebooks

        # self.clickoption('//*[@id="__apppoint"]/table/tbody/tr/td[1]/table/tbody/tr/td/div/table/tbody/tr[3]/td/div/div[1]/div[1]/span')
        # # 点击NBK1022  123
        # self.clickoption('// *[ @ id = "__apppoint"] / table / tbody / tr / td[1] / table / tbody / tr / td / div / table / tbody / tr[3] / td / div / div[1] / div[2] / div[11] / div / span')
        #
        # self.clickoption('// *[ @ id = "__apppoint"] / table / tbody / tr / td[1] / table / tbody / tr / td / div / table / tbody / tr[3] / td / div / div[1] / div[2] / div[18] / div')
        #
        # #####创建功能可以实现,直接点进去,不重复创建,进到NBK1022-21
        # self.clickoption('//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div/table/tbody/tr[2]/td/div/div[1]/div/table/tbody/tr[4]/td[6]/div/a')
        #
        # self.clickoption('//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[1]/img[1]')

        ######创建project
        ######创建项目
        ##点击Home
        try:
            self.clickoption(
                r'// *[ @ id = "__scil_toolbar_left"] / table / tbody / tr / td[2] / table / tbody / tr[2] / td')
            #####点击New进行创建
            self.clickoption(
                r'// *[ @ id = "div"] / table / tbody / tr[3] / td[3] / table / tbody / tr/ td / div / a[1]')
        except  Exception:
            print('********网络无法响应，该份报告终止上传********', flush=True)
            sys.exit(0)
        ######框里添加内容

        self.sendkeys(
            r'/ html / body / table[4] / tbody / tr[2] / td / div / div / table / tbody / tr[2] / td[2] / input',
            slide1_dic['ProjectName'])
        self.sendkeys(r'/html/body/table[4]/tbody/tr[2]/td/div/div/table/tbody/tr[3]/td[2]/input',
                      'Protein purification' + "\n")

        #####点击创建、关闭
        try:
            self.clickoption(
            r'/ html / body / table[4] / tbody / tr[2] / td / div / div / table / tbody / tr[5] / td[2] / button')
        except  Exception:
            print('********网络无法响应，该份报告终止上传********', flush=True)
            sys.exit(0)

        time.sleep(3)
        for i in range(20):
            try:
                ok = driver.find_elements(By.XPATH, r"/ html / body / table[6] / tbody / tr[2] / td / div / div / div")
                time.sleep(1)
                break
            except:

                time.sleep(5)
        else:
            print('********网络无法响应，无法点击创建项目，该份报告终止上传********', flush=True)
            sys.exit(0)  # Exit the program immediately


        status = ''
        if len(ok) == 0:
            status = True
            print('项目创建成功', flush=True)
            time.sleep(2)

        if len(ok) != 0:
            oktext = WebDriverWait(driver, 5).until(EC.presence_of_element_located(
                (By.XPATH, r'/ html / body / table[6] / tbody / tr[2] / td / div / div / div[1]'))).text
            if oktext == '[ERROR]: One project is allowed only one notebook for one user':
                print('项目已存在', flush=True)
                status = True
                try:
                    self.clickoption(r'/ html / body / table[6] / tbody / tr[2] / td / div / div / div[2] / button')

                    self.clickoption(r'/ html / body / table[4] / tbody / tr[1] / td[2] / img')

                    self.sendkeys(
                        r'// *[ @ id = "div"] / table / tbody / tr[1] / td / table / tbody / tr[1] / td[1] / input',
                        slide1_dic['ProjectName'])

                    self.clickoption(
                        r'// *[ @ id = "div"] / table / tbody / tr[1] / td / table / tbody / tr[1] / td[2] / span')

                    time.sleep(2)

                    self.clickoption(
                        r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / select[6]')
                    self.clickoption(
                        r' // *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / select[6] / option[4]')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)
                table = driver.find_element(By.XPATH,
                                            r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[1] / div / table')
                ######这个查找同时包含名字和项目编号的内容
                try:
                    ######这个查找同时包含名字和项目编号的内容
                    ######先查找第一页是否含有含姓名和项目编号的内容，并刷新循环多次，如果没有的话就来翻页查找
                    for i in range(20):
                        # print('1')
                        row = table.find_element(By.XPATH,
                                                 r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(
                                                     slide1_dic['ProjectName'], name))
                        # print('2')
                        break
                    else:
                        try:
                            for i in range(20):
                                try:
                                    self.clickoption(
                                        r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / img[1]')

                                    row = table.find_element(By.XPATH,
                                                             r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(
                                                                 slide1_dic['ProjectName'], name))
                                    break
                                except:
                                    time.sleep(2)
                            else:
                                print('********1网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                                sys.exit(0)  # Exit the program immediately
                        except:
                            print('********2网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                            sys.exit(0)  # Exit the program immediately


                except:
                    #####翻页查找内容
                    for i in range(20):
                        try:
                            self.clickoption(
                                r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[2] / div[contains(text(), "Next")]')
                            table = driver.find_element(By.XPATH,
                                                        r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[1] / div / table')
                            row = table.find_element(By.XPATH,
                                                     r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(
                                                         slide1_dic['ProjectName'], name))
                            # print(i,flush=True)
                            break
                        except:
                            time.sleep(2)
                    else:
                        print('********3网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                time.sleep(3)
                for i in range(20):
                    try:
                        WebDriverWait(row, 10).until(
                            EC.presence_of_element_located(
                                (By.XPATH, r'./td[6]'))
                        ).click()
                        # print('点击项目成功',flush=True)
                        time.sleep(1)
                        break
                    except:
                        time.sleep(5)
                else:
                    print('********4网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately

            if oktext == '[ERROR]: No privilege to create experiment in the project':
                status = False
                print('********无法创建该项目，请检查项目编号是否正确，该份报告终止上传', flush=True)
                pass
        if status == True:

            ####再回到原来的位置
            driver.execute_script("window.scrollTo(0, 0);")
            #####创建新表单
            ####新表单名字
            #
            # # 点击加号
            try:
                self.clickoption(
                    r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div/table/tbody/tr[1]/td/span[4]/img')

                #####点击Biortus
                self.clickoption(r'/html/body/table[4]/tbody/tr/td/div/div[3]')

                #######点击纯化  xpath标签总是会变,用last函数
                self.clickoption(r'/html/body/table[last()]/tbody/tr/td/div/div[7]')
            except  Exception:
                print('********网络无法响应，该份报告终止上传********', flush=True)
                sys.exit(0)

            #######part1创建表单填写表格内容
            try:
                NewEnityname = slide1_dic['条目名称'] + ' ' + purification[1][1][1]
            except:
                NewEnityname = ' '
                print('注意检查表单名称', flush=True)

            ###表单名称为NewEnityname  修改不了了暂时用这个name吧
            ######输入选项输入创建表单名称
            try:
                self.sendkeys(
                    r'/ html / body / table[6] / tbody / tr[2] / td / div / div / table / tbody / tr[4] / td[2] / input',
                    NewEnityname)
            except:
                self.sendkeys(
                    r'/ html / body / table[6] / tbody / tr[2] / td / div / div / table / tbody / tr[4] / td[2] / input',
                    ' ')
                print('检查表单名称', flush=True)

            #####点击创建
            try:
                self.clickoption(r'/html/body/table[6]/tbody/tr[2]/td/div/div/table/tbody/tr[6]/td[2]/button')
            except  Exception:
                print('********网络无法响应，该份报告终止上传********', flush=True)
                sys.exit(0)
            ######填写表格内容
            ####添加title

            try:
                self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[1]/div',
                                slide1_dic['条目名称'])
            except:

                print('检查表格内容')
                pass
            try:
                #####添加Sample ID
                self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[2]/div',
                                slide1_dic['code'])
            except:
                pass
            try:
                Reg = slide1_dic['条目名称'].split("-")[0]
                Batch = slide1_dic['条目名称'].split("-")[1]
            except:
                Reg = ''
                Batch = ''
            #####添Reg.No(Parent ID)
            self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[3]/div', Reg)
            #####添Batch No(date)
            self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[4]/div',
                            Batch)
            ######print(ProjectInformation[5],ProjectInformation[7],ProjectInformation[9])

            try:
                #####PlasmidName
                self.writetable(r'//*[@id="_eformNaN2019042926.protein1"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[1]/div',
                                slide1_dic['ProteinNname'])
            except:

                pass

            try:
                for index, i in enumerate(content_dic['tabel_list']):
                    if index == 0:
                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[6]/div',
                            i['key'])
                        if len(storageBuffer_list) != 0:
                            self.writetable(
                                r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[8]/div',
                                storageBuffer_list[0])
                        else:
                            print('缺少Storage Buffer信息', flush=True)
                            pass

                        #####Each Volume(ul)

                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[4]/div',
                            i['EachVolume'])

                        #####comment

                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[9]/div',
                            i['coment'])

                        #####EachAmount

                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[5]/div',
                            i['EachAmount'])

                        #####quantity

                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[1]/td[7]/div',
                            i['quanity'])

                    else:
                        self.clickoption(
                            '//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[1]/td/table/tbody/tr/td[4]/img')
                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[6]/div'.format(
                                index + 1), i['key'])
                        if len(storageBuffer_list) != 0:
                            self.writetable(
                                r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[8]/div'.format(
                                    index + 1), storageBuffer_list[0])
                        else:
                            pass

                        #####Each Volume(ul)
                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[4]/div'.format(
                                index + 1), i['EachVolume'])

                        #####comment
                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[9]/div'.format(
                                index + 1), i['coment'])

                        #####EachAmount
                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[5]/div'.format(
                                index + 1), i['EachAmount'])

                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein2"]/table/tbody/tr[2]/td/table[2]/tbody/tr[{}]/td[7]/div'.format(
                                index + 1), i['quanity'])
            except:
                pass

            time.sleep(3)
            print("表格填写完毕", flush=True)

            #######part2复制 cart并改名字
            try:
                self.clickoption(r'//*[@id="__scil_toolbar_right"]/table/tbody/tr/td[2]/table/tbody/tr[1]/td/img')
                ######点击show 加上关键词匹配
                self.clickoption(r'//span[contains(text(), "Show")]')
                #####关闭

                self.clickoption(r' / html / body / table[last()] / tbody / tr[1] / td[2] / img')
                ######再次点击
                self.clickoption(r'//*[@id="__scil_toolbar_right"]/table/tbody/tr/td[2]/table/tbody/tr[1]/td/img')

                ######点击show 加上关键词匹配
                self.clickoption(r'//span[contains(text(), "Show")]')
            except  Exception:
                print('********网络无法响应，该份报告终止上传********', flush=True)
                sys.exit(0)

            #####从页面中选择table
            table = driver.find_element(By.XPATH,
                                        r"/html/body/table[last()]/tbody/tr[2]/td/div/div/table/tbody/tr[1]/td[2]/div/div/table")

            #####根据列表元素复制模板
            print(step_list, flush=True)

            for step in step_list:
                for i in range(20):
                    try:
                        row = WebDriverWait(table, 10).until(
                            EC.presence_of_element_located(
                                (By.XPATH, r'.//tr[td[7]="{}"]'.format(step)))
                        )
                        time.sleep(1)
                        break

                    except:
                        time.sleep(5)
                else:
                    print('********1网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately
                time.sleep(2)

                for i in range(20):
                    try:
                        WebDriverWait(row, 10).until(
                            EC.presence_of_element_located(
                                (By.XPATH, r'./td[2]'))
                        ).click()
                        time.sleep(1)
                        break

                    except:
                        time.sleep(5)
                else:
                    print('********2网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately
                time.sleep(2)
                ####复制  滚动到这里
                try:
                    self.clickoption(r'/html/body/table[last()]/tbody/tr[2]/td/div/div/table/tbody/tr[3]/td[2]/button[1]')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)
                time.sleep(1)
                for i in range(20):
                    try:
                        WebDriverWait(row, 10).until(
                            EC.presence_of_element_located(
                                (By.XPATH, r'./td[2]'))
                        ).click()
                        time.sleep(1)
                        break

                    except:
                        time.sleep(5)
                else:
                    print('********3网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately
                time.sleep(2)
            ######关闭
            self.clickoption(r'/html/body/table[last()]/tbody/tr[1]/td[2]/img')
            print("复制模板完毕", flush=True)
            ###复制过来后改名字
            print("根据步骤修改标题", flush=True)

            for i in range(20):
                try:

                    NBKS = driver.find_elements(By.XPATH,
                                                r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[last()]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div/div')

                    time.sleep(1)
                    break

                except:
                    time.sleep(5)
            else:
                print('********1网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                sys.exit(0)  # Exit the program immediately

            for j in range(1, len(NBKS) + 1):
                try:
                    for i in range(20):
                        try:

                            WebDriverWait(driver, 10).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[{}]/div/span'.format(
                                         j)))
                            ).click()

                            time.sleep(1)
                            break
                        except:
                            time.sleep(5)
                    else:
                        print('********2网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                    for i in range(20):
                        try:

                            WebDriverWait(driver, 10).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[1]/td/img[5]'))
                            ).click()
                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********3网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                    for i in range(20):
                        try:
                            WebDriverWait(driver, 10).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     '/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[2]/input'))
                            ).clear()

                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********4网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                    ######文字
                    inner1_list = []
                    for i in range(20):
                        try:
                            element1 = WebDriverWait(driver, 10).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[{}]/div/span'.format(
                                         j)))
                            )
                            if element1.is_displayed():
                                # 获取元素的文本信息
                                text = element1.text
                                # print(text)
                                inner1_list.append(text)
                            else:
                                # 元素不可见，获取元素内的HTML代码
                                inner1 = element1.get_attribute('innerHTML')
                                inner1_list.append(inner1)
                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********5网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                    #####改名字选项
                    ####需要改成什么
                    el = inner1_list[0].split("]")[1].split("- Copy")[0].strip()
                    time.sleep(1)
                    for index, i in enumerate(step_list):
                        if el == i:
                            el1 = str(step_list.index(el) + 1) + ' ' + inner1_list[0].split("]")[1].split("- Copy")[
                                0].strip()
                            # print(el1)
                            step_list[index] = "matched"
                            # print(step_list)
                            ####用last函数
                            for i in range(20):
                                try:
                                    WebDriverWait(driver, 10).until(
                                        EC.presence_of_element_located(
                                            (By.XPATH,
                                             '/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[2]/input'))
                                    ).send_keys(el1)

                                    time.sleep(2)
                                    break

                                except:
                                    time.sleep(5)
                            else:
                                print('********6网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                                sys.exit(0)  # Exit the program immediately

                                ######save
                            for i in range(20):
                                try:
                                    WebDriverWait(driver, 10).until(
                                        EC.presence_of_element_located(
                                            (By.XPATH,
                                             r'/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[3]/button'))
                                    ).click()
                                    time.sleep(2)
                                    break

                                except:
                                    time.sleep(5)
                            else:
                                print('********7网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                                sys.exit(0)  # Exit the program immediately
                            # 修改元素值
                            # 在下一轮匹配中跳过已匹配的元素

                            break

                        else:
                            pass

                except:
                    pass

            time.sleep(2)
            try:
            #####修改完刷新一下
                self.clickoption(
                    r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[1]/span')
                self.clickoption(
                    r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[1]/td/img[1]')
            except  Exception:
                print('********网络无法响应，该份报告终止上传********', flush=True)
                sys.exit(0)

            time.sleep(2)

            ####part3  修改不同模板内容

            for j in range(1, len(NBKS) + 1):
                self.clickoption(
                    r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[{}]/div/span'.format(
                        j))
                time.sleep(1)
                ######文字
                inner2_list = []
                for i in range(20):
                    try:
                        element1 = WebDriverWait(driver, 10).until(
                            EC.presence_of_element_located(
                                (By.XPATH,
                                 r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[{}]/div/span'.format(
                                     j)))
                        )
                        if element1.is_displayed():
                            # 获取元素的文本信息
                            text = element1.text
                            # print(text)
                            inner2_list.append(text)
                        else:
                            # 元素不可见，获取元素内的HTML代码
                            inner2 = element1.get_attribute('innerHTML')
                            inner2_list.append(inner2)

                        time.sleep(1)
                        break

                    except:
                        time.sleep(5)
                else:
                    print('********8网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately
                #####根据steps进行匹配
                denum = inner2_list[0].split("]")[1].strip()
                print(denum, flush=True)
                for key, value in content_dic.items():
                    try:
                        for k, v in value.items():
                            if denum == str(key) + ' ' + str(k):
                                #######修改Cell lysis & Centrifugation
                                if k == 'Cell lysis & Centrifugation':
                                    time.sleep(2)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    print("修改模板Cell lysis & Centrifugation", flush=True)
                                    #####xpath一直变
                                    try:
                                        element = driver.find_element(By.XPATH,
                                                                      '//*[@id="tinymce"]/div/ol/li[1]/span[3][contains(text(), "20")]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['m_value']),
                                                              element)
                                    except:
                                        print('检查Cell lysis & Centrifugation 质量是否正确', flush=True)
                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH,
                                                                      '//*[@id="tinymce"]/div/ol/li[1]/span[5][contains(text(), "5.0 L B")]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['v_value']),
                                                              element)
                                    except:
                                        print('检查Cell lysis & Centrifugation 体积是否正确', flush=True)
                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[1]/span[6]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(''), element)
                                    except:

                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/ol/li[1]/span[7]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(')'), element)
                                    except:
                                        pass
                                    try:
                                        element = driver.find_element(By.XPATH,
                                                                      '//*[@id="tinymce"]/div/ol/li[1]/span[12][contains(text(), "100 ml")]')
                                        driver.execute_script(
                                            "arguments[0].innerHTML ='{}';".format(v['procedures'] + str('.')), element)
                                    except:
                                        pass

                                    try:
                                        c4 = v['ti_value'] + " "
                                    except:
                                        print('检查Cell lysis & Centrifugation', flush=True)
                                        c4 = ' '
                                    element = driver.find_element(By.XPATH,
                                                                  '//*[@id="tinymce"]/div/ol/li[4]/span[5][contains(text(), "16000 rpm for 60 min")]')
                                    driver.execute_script("arguments[0].innerHTML = '{}';".format(c4), element)
                                    try:
                                        element = driver.find_element(By.XPATH,
                                                                      '// *[ @ id = "tinymce"] / div / ol / li[3] / span[8]')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['pressure']),
                                                              element)
                                    except:
                                        pass
                                    try:
                                        #####循环删去内容
                                        element = driver.find_element(By.XPATH,
                                                                      '//*[@id="tinymce"]/div/ol/li[1]/span[15][contains(text(), "Buffer ")]')
                                        driver.execute_script("arguments[0].innerHTML = 'Buffer.'", element)
                                        for i in range(13, 22):
                                            element = driver.find_element(By.XPATH,
                                                                          '//*[@id="tinymce"]/div/ol/li[1]/span[{}]'.format(
                                                                              i))
                                            driver.execute_script("arguments[0].innerHTML = ''", element)
                                    except:
                                        pass
                                    ######修改字体颜色
                                    try:
                                        self.changecolor(r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]',
                                                         r'//*[@id="tinymce"]/div/ol/li//*[contains(@style, "color: ")]')
                                    except:
                                        pass
                                    driver.switch_to.default_content()
                                    time.sleep(1)
                                    self.saveoption()
                                    print("模板Cell lysis & Centrifugation内容修改成功", flush=True)
                                if k == 'Gallery protocol (AC-AC-SEC)':
                                    time.sleep(2)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    print("修改模板Gallery protocol (AC-AC-SEC)", flush=True)
                                    ######修改第一部分
                                    try:
                                        ######修改峰
                                        element = driver.find_element(By.XPATH,
                                                                      r'// *[ @ id = "tinymce"] / ol[3] / li[4] / span / span')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['peak'] + ' '),
                                                              element)

                                    except:
                                        print('检查Gallery protocol (AC-AC-SEC)  conclusion', flush=True)

                                    try:
                                        element = driver.find_element(By.XPATH,
                                                                      r'// *[ @ id = "tinymce"] / ol[3] / li[2] / span[7]')
                                        driver.execute_script(
                                            "arguments[0].innerHTML = '{}';".format(text_list[-1][0] + ' '), element)
                                    except:
                                        print('检查Gallery protocol (AC-AC-SEC)  分子筛类型', flush=True)
                                    #####循环删除分子筛其他部分
                                    try:

                                        element = driver.find_element(By.XPATH,
                                                                      '// *[ @ id = "tinymce"] / ol[3] / li[2] / span[9]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        element = driver.find_element(By.XPATH,
                                                                      '// *[ @ id = "tinymce"] / ol[3] / li[2] / span[8]')
                                        driver.execute_script("arguments[0].remove()", element)
                                    except:
                                        pass
                                    try:
                                        if len(text_list) == 2:

                                            if 'His' in text_list[0][0]:
                                                print('His')
                                                element = driver.find_element(By.XPATH,
                                                                              r'//*[@id="tinymce"]/ol[1]/li[1]/span[3]')
                                                driver.execute_script(
                                                    "arguments[0].innerHTML = '{}';".format(text_list[0][0] + ' '), element)
                                                element = driver.find_element(By.XPATH,
                                                                              r'//*[@id="tinymce"]/ol[1]/li[1]/span[4]')
                                                driver.execute_script(
                                                    "arguments[0].innerHTML = '{}';".format(text_list[0][1] + ' '), element)
                                                ###循环删去strep的
                                                element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / ol[2]')
                                                driver.execute_script("arguments[0].remove()", element)
                                                element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / p[3]')
                                                driver.execute_script("arguments[0].remove()", element)
                                            if 'Strep' in text_list[0][0]:
                                                element = driver.find_element(By.XPATH,
                                                                              r'// *[ @ id = "tinymce"] / ol[2] / li[1] / div / p / span[1] / span')
                                                driver.execute_script(
                                                    "arguments[0].innerHTML = '{}';".format('strep ' + text_list[0][1]),
                                                    element)
                                                element = driver.find_element(By.XPATH,
                                                                              '// *[ @ id = "tinymce"] / p[1] / strong')
                                                driver.execute_script("arguments[0].remove()", element)
                                                element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / ol[1]')
                                                driver.execute_script("arguments[0].remove()", element)
                                    except:
                                        print('检查Purification Scheme流程图是否有问题', flush=True)

                                    try:
                                        if len(text_list) == 3:
                                            if 'His' in text_list[0][0] and 'Strep' in text_list[1]:
                                                element = driver.find_element(By.XPATH,
                                                                              r'//*[@id="tinymce"]/ol[1]/li[1]/span[3]')
                                                driver.execute_script(
                                                    "arguments[0].innerHTML = '{}';".format(text_list[0][0]), element)
                                                element = driver.find_element(By.XPATH,
                                                                              r'//*[@id="tinymce"]/ol[1]/li[1]/span[4]')
                                                driver.execute_script(
                                                    "arguments[0].innerHTML = '{}';".format(text_list[0][1]), element)
                                                element = driver.find_element(By.XPATH,
                                                                              r'// *[ @ id = "tinymce"] / ol[2] / li[1] / div / p / span[1] / span')
                                                driver.execute_script(
                                                    "arguments[0].innerHTML = '{}';".format('strep ' + text_list[0][1]),
                                                    element)
                                            else:
                                                element = driver.find_element(By.XPATH,
                                                                              r'//*[@id="tinymce"]/ol[1]/li[1]/span[3]')
                                                driver.execute_script(
                                                    "arguments[0].innerHTML = '{}';".format(text_list[0][0]), element)
                                                element = driver.find_element(By.XPATH,
                                                                              r'//*[@id="tinymce"]/ol[1]/li[1]/span[4]')
                                                driver.execute_script(
                                                    "arguments[0].innerHTML = '{}';".format(text_list[0][1]), element)
                                                element = driver.find_element(By.XPATH,
                                                                              r'// *[ @ id = "tinymce"] / ol[2] / li[1] / div / p / span[1] / span')
                                                driver.execute_script(
                                                    "arguments[0].innerHTML = '{}';".format('strep ' + text_list[0][1]),
                                                    element)
                                    except:
                                        print('检查Purification Scheme流程图是否有问题', flush=True)

                                    self.changecolor(r'//*[@id="tinymce"]/ol/li//*[contains(@style, "color: ")]',
                                                     r'//*[@id ="tinymce"]/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    time.sleep(2)
                                    driver.switch_to.default_content()

                                    ######插入QC部分Results & Discussion:
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    #####修改内容
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[3]/span')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(
                                            'Batch No.: ' + content_dic['QC']['Batch']), element)
                                    except:
                                        pass
                                    ######修改字体颜色
                                    try:
                                        self.changecolor(r'//*[@id="tinymce"]/div/p/span[contains(@style, "color: ")]'
                                                         ,
                                                         r'//*[@id ="tinymce"]/div/p/span[contains(@data-mce-style, "color: ")]')
                                    except:
                                        print('检查QC', flush=True)
                                        pass

                                    #####删除QC3那一部分的内容
                                    try:
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[2] / p[13]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[2] / p[12]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[2] / p[11]')
                                        driver.execute_script("arguments[0].remove()", element)
                                    except:
                                        pass
                                    try:
                                        i_list = []
                                        for i, j in content_dic['QC'].items():
                                            if type(i) == int:
                                                # print(i)
                                                i_list.append(i)
                                    except:
                                        pass

                                    ###插入zoomin
                                    try:
                                        image = 'QCzoomin_' + str(content_dic['QC']['index']) + '_' + str(
                                            i_list[0]) + '.jpg'
                                        # print(image,flush=True)
                                        self.insertimage(image, '//*[@id="tinymce"]/div[2]/p[9]/img',
                                                         '//*[@id="tinymce"]/div[2]/p[9]')
                                    except:
                                        print('检查QC (zoom in) 图片', flush=True)
                                        pass
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[7]')
                                    driver.execute_script("arguments[0].remove()", element)
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[6]')
                                    driver.execute_script("arguments[0].remove()", element)
                                    #####插入第一张图
                                    try:
                                        image = 'QC_' + str(content_dic['QC']['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/div[2]/p[5]/img',
                                                         '//*[@id="tinymce"]/div[2]/p[5]')
                                    except:
                                        print('检查QC1,QC2,QC3图片是否有问题', flush=True)
                                        pass
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[4]/span')
                                    driver.execute_script(
                                        "arguments[0].innerHTML = 'QC 1: SDS-PAGE & QC 2: LC-MS  & QC 3: Analytical SEC (Superdex 200 Increase 5/150 GL) QC buffer: 50 mM Tris-HCl (pH 7.5), 300 mM NaCl';",
                                        element)
                                    try:
                                        if len(i_list) > 1:
                                            # print('需要添加图片')
                                            for i in range(1, len(i_list)):
                                                body = driver.find_element(By.TAG_NAME, 'body')
                                                # 将光标移到文本末尾处
                                                body.send_keys(Keys.END)
                                                # 按下回车键换行
                                                body.send_keys(Keys.ENTER)
                                                # image_file_path = os.path.abspath('image_file_path')
                                                image_file_path = self.image_file_path
                                                file = 'QCzoomin_' + str(content_dic['QC']['index']) + '_' + str(
                                                    i_list[i]) + '.jpg'
                                                # print(file,flush=True)
                                                with open(os.path.join(image_file_path, file), 'rb') as f:
                                                    image_data = f.read()
                                                    # 将图片数据编码为 base64 字符串
                                                image_base64 = base64.b64encode(image_data).decode('utf-8')
                                                # 将 base64 字符串插入到 HTML 代码中
                                                html_str = f'<img src="data:image/jpeg;base64,{image_base64}">'
                                                ###添加
                                                start_offset = 0
                                                end_offset = 0

                                                #####图片删除后path没了，要重新看新的xpayh是什么
                                                element = driver.find_element(By.XPATH,
                                                                              r'//*[@id="tinymce"]/div[2]/p[{}]'.format(
                                                                                  int(7 + i)))
                                                driver.execute_script("""
                                                                                                                                            var range_obj = document.createRange();
                                                                                                                                            range_obj.setStart(arguments[0], arguments[1]);
                                                                                                                                            range_obj.setEnd(arguments[0], arguments[2]);
                                                                                                                                            var fragment = range_obj.createContextualFragment(arguments[3]);
                                                                                                                                            range_obj.insertNode(fragment);
                                                                                                                                        """,
                                                                      element, start_offset, end_offset, html_str)
                                    except:
                                        print('检查QC (zoom in) 图片', flush=True)

                                    try:
                                        file = 'Affinity & SEC_' + str(v['index']) + '.jpg'
                                        with open(os.path.join(image_file_path, file), 'rb') as f:
                                            image_data = f.read()
                                            # 将图片数据编码为 base64 字符串
                                        image_base64 = base64.b64encode(image_data).decode('utf-8')
                                        # 将 base64 字符串插入到 HTML 代码中
                                        html_str = f'<img src="data:image/jpeg;base64,{image_base64}">'
                                        ###添加
                                        start_offset = 0
                                        end_offset = 0

                                        #####图片删除后path没了，要重新看新的xpayh是什么
                                        element = driver.find_element(By.XPATH, r'//*[@id="tinymce"]/p[1]')
                                        driver.execute_script("""
                                                                                                                                    var range_obj = document.createRange();
                                                                                                                                    range_obj.setStart(arguments[0], arguments[1]);
                                                                                                                                    range_obj.setEnd(arguments[0], arguments[2]);
                                                                                                                                    var fragment = range_obj.createContextualFragment(arguments[3]);
                                                                                                                                    range_obj.insertNode(fragment);
                                                                                                                                """,
                                                              element, start_offset, end_offset, html_str)
                                    except:
                                        print('检查AC-AC-SEC图片', flush=True)
                                    if len(i_list) == 0:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[7]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[6]')
                                        driver.execute_script("arguments[0].remove()", element)
                                    driver.switch_to.default_content()
                                    self.saveoption()

                                    print("模板Gallery protocol (AC-AC-SEC)内容修改成功", flush=True)
                                if k == 'Gallery protocol (AC-Digestion-AC-SEC)':
                                    time.sleep(2)
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[1])
                                    print("修改模板Gallery protocol (AC-Digestion-AC-SEC)", flush=True)
                                    ######修改第一部分
                                    try:
                                        ######修改峰
                                        element = driver.find_element(By.XPATH, r'//*[@id="tinymce"]/ol[3]/li[4]/span/span')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(v['peak'] + ' '),
                                                              element)

                                    except:
                                        print('检查Gallery protocol (AC-Digestion-AC-SEC) conclusion', flush=True)

                                    try:
                                        element = driver.find_element(By.XPATH, r'//*[@id="tinymce"]/ol[3]/li[2]/span[7]')
                                        driver.execute_script(
                                            "arguments[0].innerHTML = '{}';".format(text_list[-1][0] + ' '), element)
                                    except:
                                        print('检查Gallery protocol (AC-AC-SEC)  分子筛类型', flush=True)
                                    #####循环删除分子筛其他部分
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/ol[3]/li[2]/span[9]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/ol[3]/li[2]/span[8]')
                                        driver.execute_script("arguments[0].remove()", element)
                                    except:
                                        pass
                                    try:
                                        if len(text_list) == 3:
                                            element = driver.find_element(By.XPATH,
                                                                          r'//*[@id="tinymce"]/ol[1]/li[1]/span[3]')
                                            driver.execute_script(
                                                "arguments[0].innerHTML = '{}';".format(text_list[0][0] + ' '), element)
                                            element = driver.find_element(By.XPATH,
                                                                          r'//*[@id="tinymce"]/ol[1]/li[1]/span[4]')
                                            driver.execute_script(
                                                "arguments[0].innerHTML = '{}';".format(text_list[0][1] + ' '), element)

                                        else:
                                            element = driver.find_element(By.XPATH,
                                                                          r'//*[@id="tinymce"]/ol[1]/li[1]/span[3]')
                                            driver.execute_script(
                                                "arguments[0].innerHTML = '{}';".format(text_list[0][0] + ' '), element)
                                            element = driver.find_element(By.XPATH,
                                                                          r'//*[@id="tinymce"]/ol[1]/li[1]/span[4]')
                                            driver.execute_script(
                                                "arguments[0].innerHTML = '{}';".format(text_list[0][1] + ' '), element)
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/ol[1]/li[5]/span')
                                            element.location_once_scrolled_into_view
                                            element.send_keys(Keys.ENTER)
                                            element.send_keys(
                                                'Clarified supernatant was loaded onto Binding Buffer A pre-equilibrated ' +
                                                text_list[1][0] + ' ' + text_list[1][1] + ' ' + 'column.')
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/ol[1]/li[6]/span')
                                            element.location_once_scrolled_into_view
                                            element.send_keys(Keys.ENTER)
                                            element.send_keys('Column was washed with 10 CV of Binding Buffer A')
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/ol[1]/li[7]/span')
                                            element.location_once_scrolled_into_view
                                            element.send_keys(Keys.ENTER)
                                            element.send_keys('Columnwas eluted with Buffer B.')
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/ol[1]/li[8]/span')
                                            element.location_once_scrolled_into_view
                                            element.send_keys(Keys.ENTER)
                                            element.send_keys('Samples were prepared for SDS-PAGE')
                                            element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/ol[1]/li[9]/span')
                                            element.location_once_scrolled_into_view
                                            element.send_keys(Keys.ENTER)
                                            element.send_keys('Fraction Buffer B Elution was collected')
                                    except:
                                        print('检查Purification Scheme流程图是否有问题', flush=True)

                                    self.changecolor(r'//*[@id="tinymce"]/ol/li//*[contains(@style, "color: ")]',
                                                     r'//*[@id ="tinymce"]/ol/li//*[contains(@data-mce-style, "color: ")]')

                                    time.sleep(2)
                                    driver.switch_to.default_content()

                                    ######插入QC部分Results & Discussion:
                                    iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                    driver.switch_to.frame(iframe_list[2])
                                    #####修改内容
                                    try:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[3]/span')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format(
                                            'Batch No.: ' + content_dic['QC']['Batch']), element)
                                    except:
                                        pass
                                    ######修改字体颜色
                                    try:
                                        self.changecolor(r'//*[@id="tinymce"]/div/p/span[contains(@style, "color: ")]'
                                                         ,
                                                         r'//*[@id ="tinymce"]/div/p/span[contains(@data-mce-style, "color: ")]')
                                    except:
                                        print('检查QC', flush=True)
                                        pass

                                    #####删除QC3那一部分的内容
                                    try:
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[2] / p[13]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[2] / p[12]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / div[2] / p[11]')
                                        driver.execute_script("arguments[0].remove()", element)
                                    except:
                                        pass
                                    try:
                                        i_list = []
                                        for i, j in content_dic['QC'].items():
                                            if type(i) == int:
                                                # print(i)
                                                i_list.append(i)
                                    except:
                                        pass

                                    ###插入zoomin
                                    try:
                                        image = 'QCzoomin_' + str(content_dic['QC']['index']) + '_' + str(
                                            i_list[0]) + '.jpg'
                                        # print(image,flush=True)
                                        self.insertimage(image, '//*[@id="tinymce"]/div[2]/p[9]/img',
                                                         '//*[@id="tinymce"]/div[2]/p[9]')
                                    except:
                                        print('检查QC (zoom in) 图片', flush=True)
                                        pass
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[7]')
                                    driver.execute_script("arguments[0].remove()", element)
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[6]')
                                    driver.execute_script("arguments[0].remove()", element)
                                    #####插入第一张图
                                    try:
                                        image = 'QC_' + str(content_dic['QC']['index']) + '.jpg'
                                        self.insertimage(image, '//*[@id="tinymce"]/div[2]/p[5]/img',
                                                         '//*[@id="tinymce"]/div[2]/p[5]')
                                    except:
                                        print('检查QC1,QC2,QC3图片是否有问题', flush=True)
                                        pass
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[4]/span')
                                    driver.execute_script(
                                        "arguments[0].innerHTML = 'QC 1: SDS-PAGE & QC 2: LC-MS  & QC 3: Analytical SEC (Superdex 200 Increase 5/150 GL) QC buffer: 50 mM Tris-HCl (pH 7.5), 300 mM NaCl';",
                                        element)
                                    try:
                                        if len(i_list) > 1:
                                            # print('需要添加图片')
                                            for i in range(1, len(i_list)):
                                                body = driver.find_element(By.TAG_NAME, 'body')
                                                # 将光标移到文本末尾处
                                                body.send_keys(Keys.END)
                                                # 按下回车键换行
                                                body.send_keys(Keys.ENTER)
                                                # image_file_path = os.path.abspath('image_file_path')
                                                image_file_path = self.image_file_path
                                                file = 'QCzoomin_' + str(content_dic['QC']['index']) + '_' + str(
                                                    i_list[i]) + '.jpg'
                                                # print(file,flush=True)
                                                with open(os.path.join(image_file_path, file), 'rb') as f:
                                                    image_data = f.read()
                                                    # 将图片数据编码为 base64 字符串
                                                image_base64 = base64.b64encode(image_data).decode('utf-8')
                                                # 将 base64 字符串插入到 HTML 代码中
                                                html_str = f'<img src="data:image/jpeg;base64,{image_base64}">'
                                                ###添加
                                                start_offset = 0
                                                end_offset = 0

                                                #####图片删除后path没了，要重新看新的xpayh是什么
                                                element = driver.find_element(By.XPATH,
                                                                              r'//*[@id="tinymce"]/div[2]/p[{}]'.format(
                                                                                  int(7 + i)))
                                                driver.execute_script("""
                                                                                                                                            var range_obj = document.createRange();
                                                                                                                                            range_obj.setStart(arguments[0], arguments[1]);
                                                                                                                                            range_obj.setEnd(arguments[0], arguments[2]);
                                                                                                                                            var fragment = range_obj.createContextualFragment(arguments[3]);
                                                                                                                                            range_obj.insertNode(fragment);
                                                                                                                                        """,
                                                                      element, start_offset, end_offset, html_str)
                                    except:
                                        print('检查QC (zoom in) 图片', flush=True)

                                    try:
                                        file = 'Affinity & Digestion & SEC_' + str(v['index']) + '.jpg'
                                        with open(os.path.join(image_file_path, file), 'rb') as f:
                                            image_data = f.read()
                                            # 将图片数据编码为 base64 字符串
                                        image_base64 = base64.b64encode(image_data).decode('utf-8')
                                        # 将 base64 字符串插入到 HTML 代码中
                                        html_str = f'<img src="data:image/jpeg;base64,{image_base64}">'
                                        ###添加
                                        start_offset = 0
                                        end_offset = 0

                                        #####图片删除后path没了，要重新看新的xpayh是什么
                                        element = driver.find_element(By.XPATH, r'//*[@id="tinymce"]/p[1]')
                                        driver.execute_script("""
                                                                                                                                    var range_obj = document.createRange();
                                                                                                                                    range_obj.setStart(arguments[0], arguments[1]);
                                                                                                                                    range_obj.setEnd(arguments[0], arguments[2]);
                                                                                                                                    var fragment = range_obj.createContextualFragment(arguments[3]);
                                                                                                                                    range_obj.insertNode(fragment);
                                                                                                                                """,
                                                              element, start_offset, end_offset, html_str)
                                    except:
                                        print('检查AC-AC-SEC图片', flush=True)

                                    if len(i_list) == 0:
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[7]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div[2]/p[6]')
                                        driver.execute_script("arguments[0].remove()", element)
                                    driver.switch_to.default_content()
                                    self.saveoption()

                                    print("模板Gallery protocol (AC-Digestion-AC-SEC)内容修改成功", flush=True)

                    except:
                        pass
            print("完成修改模板内容任务", flush=True)
            for i in range(20):
                try:
                    WebDriverWait(driver, 10).until(
                        EC.presence_of_element_located(
                            (By.XPATH,
                             '//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[last()]/table/tbody/tr[1]/td/span[3]/img'))
                    ).click()
                    time.sleep(1)
                    break

                except:
                    time.sleep(5)
            else:
                print('********1网络无法响应，无法上传ppt，该份报告终止上传********', flush=True)
                sys.exit(0)  # Exit the program immediately

            for i in range(20):
                try:
                    WebDriverWait(driver, 10).until(
                        EC.presence_of_element_located(
                            (By.XPATH,
                             '/html/body/table[last()]/tbody/tr/td/div/div[2]/table/tbody/tr/td[1]/span'))
                    ).click()
                    time.sleep(1)
                    break

                except:
                    time.sleep(5)
            else:
                print('********2网络无法响应，无法上传ppt，该份报告终止上传********', flush=True)
                sys.exit(0)  # Exit the program immediately

            for i in range(20):
                try:
                    WebDriverWait(driver, 10).until(
                        EC.presence_of_element_located(
                            (By.XPATH,
                             '/html/body/table[last()]/tbody/tr/td/div/div[6]/span'))
                    ).click()
                    time.sleep(1)
                    break

                except:
                    time.sleep(5)
            else:
                print('********3网络无法响应，无法上传ppt，该份报告终止上传********', flush=True)
                sys.exit(0)  # Exit the program immediately

            file_dialog = driver.find_element(By.XPATH, r'//input[@type="file"]')
            time.sleep(2)
            file_dialog.send_keys(ppt_file)

            for i in range(20):
                try:
                    WebDriverWait(driver, 10).until(
                        EC.presence_of_element_located(
                            (By.XPATH,
                             '/ html / body / table[last()] / tbody / tr[2] / td / div / div / div / form / table / tbody / tr[7] / td[2] / button'))
                    ).click()
                    time.sleep(1)
                    break

                except:
                    time.sleep(5)
            else:
                print('********4网络无法响应，无法上传ppt，该份报告终止上传********', flush=True)
                sys.exit(0)  # Exit the program immediately
            print("pptx文件上传成功", flush=True)
            time.sleep(5)
        else:
            pass
@Gooey(
    richtext_controls=True,  # 打开终端对颜色支持
    program_name="纯化ELN报告上传",  # 程序名称
    clear_before_run=True,
    required_cols=2
)
def main():
    settings_msg = '纯化ELN报告上传'
    parser = GooeyParser(description=settings_msg)  # 添加上方的应用信息
    subs = parser.add_subparsers(help='commands', dest='command')
    # 导入配置文件
    # account,password, ppt_file, name= eln.load_configuration()
    SingleSearchparser = subs.add_parser('纯化ELN报告上传')
    SingleSearchparser.add_argument("script", choices=["非gallery", "gallery"], help="选择上传报告类型")
    SingleSearchparser.add_argument("account", metavar='ELN账号')
    SingleSearchparser.add_argument("password", metavar='ELN账号密码')
    SingleSearchparser.add_argument("name", metavar='ELN Author')
    SingleSearchparser.add_argument("ppt_files", metavar='ppt文件夹',widget='DirChooser',help="选择ppt文件夹")

    """
       程序功能实现
    """
    args = parser.parse_args()

    """界面获得的数据进行处理"""
    if args.command == "纯化ELN报告上传":
        account = args.account.strip()
        password = args.password.strip()
        ppt_path = args.ppt_files.strip()
        name=args.name

        if not os.path.exists('content'):
            os.makedirs('content')
        if os.path.exists('content'):
            # print('删除content文件夹下的所有文件夹', flush=True)
            for root, dirs, files in os.walk('content', topdown=False):
                for file in files:
                    # print('Deleting file:', os.path.join(root, file))
                    os.remove(os.path.join(root, file))
                for dir in dirs:
                    # print('Deleting folder:', os.path.join(root, dir))
                    os.rmdir(os.path.join(root, dir))
    ppt_files = [file for file in os.listdir(ppt_path) if file.endswith('.pptx')]
    total_count = len(ppt_files)


    if args.script == "非gallery":
        ####检查Chrome
        if eln.check_broswer():
            make_print_to_file()
            current_time = datetime.now()
            print('-------------程序开始运行时间：' + str(current_time) + '-------------')
            for index, file in enumerate(ppt_files):
                ppt_file = os.path.join(ppt_path, file)
                ppt_file_name = file.split('.pptx')[0]
                # print(index, ppt_file_name,ppt_file)
                ppt_file_name = ppt_file_name.replace('.', '-')
                ######获取存放文件的路径
                image_file = os.path.abspath('content')
                ######获取每一个pptx存放图片的路径
                image_file_path = os.path.join(image_file, ppt_file_name)
                # print(image_file_path)
                if not os.path.exists(image_file_path):
                    os.makedirs(image_file_path)
                if os.path.exists(image_file_path):
                    for roots, dirs, files in os.walk(image_file_path):
                        for file in files:
                            os.remove(os.path.join(image_file_path, file))
                # 实例化对象

                print('*********正在处理第{}/{}个文件*********'.format(index + 1, total_count))
                print(ppt_file)

                ########尝试三次上传，如果失败的话，就跳过

                try_count = 0
                max_try_count = 3

                while try_count < max_try_count:
                    try:
                        #####实例化，重新进行上传，而不是在原来的基础上
                        eln_auto = eln(account, password, ppt_file, image_file_path, name)
                        eln_auto.upload_ELN_out()
                        break  # Break out of the loop if upload is successful
                    except Exception:
                        time.sleep(2)
                        print(f'*********{ppt_file_name}+第{try_count+1}次上传失败*********', flush=True)
                        try_count += 1
                        if try_count < max_try_count:
                            print(f'尝试进行第{try_count+1} 次上传...')
                            # eln_auto.upload_ELN_out()
                            time.sleep(1)  # Add a delay before retrying

                if try_count == max_try_count:
                    print(f'*********{ppt_file_name}+上传最终失败*********', flush=True)

            print('*********所有内容上传完成,注意检查核对*********', flush=True)
            current_time = datetime.now()
            print('-------------程序结束运行时间：' + str(current_time) + '-------------', flush=True)

        else:
            win32api.MessageBox(0, "Chrome浏览器或驱动安装失败, 请手动安装重新运行程序, 或联系技术人员", "提醒", win32con.MB_TOPMOST)

    if args.script == "gallery":
        ####检查Chrome
        make_print_to_file()
        if eln.check_broswer():
            current_time = datetime.now()
            print('-------------程序开始运行时间：' + str(current_time) + '-------------')
            for index, file in enumerate(ppt_files):
                ppt_file = os.path.join(ppt_path, file)
                ppt_file_name = file.split('.pptx')[0]
                # print(index, ppt_file_name,ppt_file)
                ppt_file_name = ppt_file_name.replace('.', '-')
                ######获取存放文件的路径
                image_file = os.path.abspath('content')
                ######获取每一个pptx存放图片的路径
                image_file_path = os.path.join(image_file, ppt_file_name)
                if not os.path.exists(image_file_path):
                    os.makedirs(image_file_path)
                if os.path.exists(image_file_path):
                    for roots, dirs, files in os.walk(image_file_path):
                        for file in files:
                            os.remove(os.path.join(image_file_path, file))

                print('*********正在处理第{}/{}个文件*********'.format(index + 1, total_count))
                print(ppt_file)
                ########尝试三次上传，如果失败的话，就跳过

                try_count = 0
                max_try_count = 3

                while try_count < max_try_count:
                    try:
                        #####实例化，重新进行上传，而不是在原来的基础上
                        eln_auto = eln(account, password, ppt_file, image_file_path, name)
                        eln_auto.upload_ELN_gallery()
                        time.sleep(3)
                        break  # Break out of the loop if upload is successful
                    except:
                        time.sleep(1)
                        print(f'*********{ppt_file_name}+第{try_count + 1}次上传失败*********', flush=True)
                        try_count += 1
                        if try_count < max_try_count:
                            print(f'尝试进行第{try_count+1} 次上传...')
                            time.sleep(1)  # Add a delay before retrying
                if try_count == max_try_count:
                    print(f'*********{ppt_file_name}+上传最终失败*********', flush=True)
            print('*********所有内容上传完成,注意检查核对*********', flush=True)
            current_time = datetime.now()
            print('-------------程序结束运行时间：' + str(current_time) + '-------------', flush=True)
        else:
            win32api.MessageBox(0, "Chrome浏览器或驱动安装失败, 请手动安装重新运行程序, 或联系技术人员", "提醒", win32con.MB_TOPMOST)
if __name__ == '__main__':
    try:
        main()
        ######打包后运行完没有提示，加上这个
        win32api.MessageBox(0, "任务完成", "提醒", win32con.MB_TOPMOST)
        input('\n')
    except Exception as e:
        traceback.print_exc()
        input("任务出错，联系计算结构平台")
        input("按任意键退出:")
