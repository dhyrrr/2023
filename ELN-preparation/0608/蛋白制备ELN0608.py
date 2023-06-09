from selenium.webdriver.chrome.service import Service
import time
import win32api, win32con
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.action_chains import ActionChains
import base64
from pptx.shapes.picture import Picture
import re
import os
from selenium.webdriver.common.keys import Keys
from PIL import Image
from pptx_tools import utils
from szpscript.Broswer_script import Broswer
import traceback
from gooey import Gooey, GooeyParser
import json
import  sys
import codecs
from datetime import datetime
import pptx
from pptx import Presentation
import shutil

'''
0524 把QC和WB的脚本合并  QC的应该没问题，WB的需要再修改测试
0525 上传WB 多个抗体的，多个质粒的情况上传表格,    QC 和WB批量
0526 批量并行上传 
0531  1.批量不并行上传 
      2.wb的保存填title,WB的都根据ppt的页面来确定，每个模板的名称都是ppt上的标题 
      3.QC的第一页由页码2来确定，不含有多个，只可能出现一页。且把之前写的根据QC和wab类型来确定什么的删掉，直接分开写。把含有多个步骤的QCs的改名字代码删掉
      4.WB的表单名都和ppt上的标题一样，不要搞那么复杂，直接复制过来，QC的和纯化的一样，不需要改动
0601  1.如果有网络异常的直接关掉
      2.将输出信息写入log文件
      3.测试
      4.和纯化的一样，创建项目的时候有的项目没那么多//*[@id="div"]/table/tbody/tr[3]/td[3]/table/tbody/tr/td/div/a[1]
      5.RGB

0605 1.clickoption    的关闭页面，调整
	 2.打包后跑完没有提示，点击stop不管用
0606  1.wb第一页的表格，如果没有的话，就不管
      2.western 和qc按照标准模板来，都是用纯化的来， wb的改成，  三行提取，wb sample time   wb的截取整页信息，  
0607  1.填项目编号的时候，有的会出现下拉的框，设置个异常处理
      2.有的是时候会卡住，把那个截图的移到前面
      3.已经上传完的和失败的项目编号有问题的分开
0608  1.移动文件（上传完成和项目编号有问题）的时候，如果文件已经存在，就删除
      2.填写那个项目编号的时候，等个几秒，防止明明是已经存在或有问题的缺提示创建成功
'''
if sys.stdout.encoding != 'UTF-8':
    sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer, 'strict')
if sys.stderr.encoding != 'UTF-8':
    sys.stderr = codecs.getwriter('utf-8')(sys.stderr.buffer, 'strict')

def make_print_to_file():
    savename = "out.log"
    def write_to_log(message):
        with open(savename, "a", encoding='utf-8') as log_file:
            log_file.write(message)
        sys.__stdout__.write(message)

    sys.stdout.write = write_to_log

def delete_text_boxes_with_keyword(ppt_file, image_file_path):
        # ppt_file = self.ppt_file
        # image_file_path = self.image_file_path
        # print('生成copy文件',image_file_path)
        # 打开 PowerPoint 文件
        prs = Presentation(ppt_file)
        # 遍历每个幻灯片
        for slide in prs.slides:
            # 遍历每个形状
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 判断文本是否包含关键字
                            if 'Total amount' in run.text:
                                # 删除包含关键字的文本框
                                slide.shapes._spTree.remove(shape._element)
                                break
        # 保存更改后的 PowerPoint 文件
        new_ppt_file = os.path.join(image_file_path, 'copy.pptx')
        prs.save(new_ppt_file)
        time.sleep(2)
        utils.save_pptx_as_png(image_file_path, new_ppt_file, overwrite_folder=True)

class Params:
    def __init__(self):
        pass
    website = r"https://scilligence.net/Biortus/ELN/Explorer.aspx"
    selenium_waittime = 3000
    broswer_name = "Chrome"
    chrome_driver_website = r"https://registry.npmmirror.com/binary.html?path=chromedriver/"
    chrome_driver_path = r"./chromedriver.exe"

class out:

    __chrome_driver_path = Params.chrome_driver_path

    def __init__(self,account,password,ppt_file,image_file_path,name):
        self.account=account
        self.password=password
        self.ppt_file=ppt_file
        self.name = name
        self.image_file_path=image_file_path
        self.chrome_driver_path = out.__chrome_driver_path
        self.driver = self.build_selenium_obj(self.chrome_driver_path)
        self.website = Params.website
        self.waittime = Params.selenium_waittime
        self.witness_finished_status = False  # 判定witness任务是否处理完, T
    #######静态方法通常用于定义与类相关但不依赖于类变量和实例变量的功能
    @staticmethod
    def check_broswer():
        # 尝试3次
        for i in range(5):
            broswer_version = Broswer.get_broswer_version(Params.broswer_name)

            if broswer_version:
                # 检查驱动是否存在
                chrome_driver_new_path = out.__chrome_driver_path.replace(".exe",
                                                                              "_V{0}.exe".format(broswer_version))
                if not os.path.exists(chrome_driver_new_path):
                    print("下载驱动中!",flush=True)
                    # 下载驱动
                    Broswer.download_chrome_driver(broswer_version)

                    # 根据版本更改驱动名字
                    if os.path.exists(out.__chrome_driver_path):
                        os.rename(out.__chrome_driver_path, chrome_driver_new_path)
                        out.__chrome_driver_path = chrome_driver_new_path
                        return True
                else:
                    # 驱动已经存在, 根据版本更改驱动名字
                    out.__chrome_driver_path = chrome_driver_new_path
                    return True

            else:
                # 安装
                win32api.MessageBox(0, "Chrome浏览器未安装, 现在开始安装", "提醒", win32con.MB_TOPMOST)
                Broswer.install_chrome()
        else:
            # 安装浏览器或驱动失败
            return False
    @staticmethod
    def load_configuration():
        if os.path.exists("configuration.json"):
            with open("configuration.json", "r", encoding="utf-8") as f:
                configuration_dict = json.load(f)

            try:
                account  = configuration_dict["account "]
            except:
                account  = ""

            try:
                password = configuration_dict["password"]
            except:
                password = ""

            try:
                ppt_file = configuration_dict["ppt_file"]
            except:
                ppt_file = ""
            try:
                name = configuration_dict["name"]
            except:
                name = ""

            # try:
            #     image_file_path = configuration_dict["image_file_path"]
            # except:
            #     image_file_path = ""

        else:
            account = ""
            password = ""
            ppt_file = ""
            name = ""

        return account,password, ppt_file,name
    @staticmethod
    def build_selenium_obj(chrome_driver_path):
        #######解决chrome正在收到测试软件监控的问题
        s = Service(executable_path=chrome_driver_path)
        options = webdriver.ChromeOptions()
        options.add_experimental_option("excludeSwitches",['enable-automation'])
        driver = webdriver.Chrome(service=s,options=options)
        driver.command_executor._commands["send_command"] = ("POST", '/session/$sessionId/chromium/send_command')
        params = {'cmd': 'Page.setDownloadBehavior',
                  'params': {'behavior': 'allow', 'downloadPath': os.getcwd()}}
        driver.execute("send_command", params)
        return driver
    ######读取pptx内容，保存图片及相关信息
    ######根据ppt内容自动上传ELN
    def writetable(self,xpath_name,table_content):
        driver = self.driver
        for i in range(20):
            try:
                double_click_element = driver.find_element(By.XPATH,xpath_name)
                time.sleep(1)
                # 实例化ActionChains对象
                actions = ActionChains(driver)
                # 双击需要输入文本的框
                actions.double_click(double_click_element).perform()
                # 在框中输入文本
                active_element = driver.switch_to.active_element
                active_element.send_keys(table_content)
                break

            except:
                time.sleep(5)
        else:
            pass
            # print('********无法创建该项目，请检查项目编号是否正确，该份报告终止上传********', flush=True)
            # sys.exit(0)  # Exit the program immediately
    def clickoption(self,xpath_name):
        driver = self.driver

        for i in range(20):
            try:
                WebDriverWait(driver, 10).until(
                    EC.presence_of_element_located(
                        (By.XPATH,xpath_name))).click()
                time.sleep(1)
                break

            except:
                time.sleep(5)
        else:
            pass
    def sendkeys(self, xpath_name,keys):
        driver = self.driver
        for i in range(20):
            try:
                WebDriverWait(driver, 10).until(
                    EC.presence_of_element_located(
                        (By.XPATH, xpath_name))
                ).send_keys(keys)

                time.sleep(1)
                break

            except:
                time.sleep(5)
        else:
            pass
            # print('********无法创建该项目，请检查项目编号是否正确，该份报告终止上传********', flush=True)
            # sys.exit(0)  # Exit the program immediately
    def changecolor(self,xpath1,xpath2):
        driver = self.driver
        for i in range(20):
            try:
                elements = driver.find_elements(By.XPATH,
                                                xpath1)
                for element in elements:
                    # 设置字体颜色
                    driver.execute_script("arguments[0].style.color = 'black';", element)
                elements1 = driver.find_elements(By.XPATH,
                                                 xpath2)
                for element1 in elements1:
                    driver.execute_script("arguments[0].setAttribute('data-mce-style', arguments[1])", element1,
                                          'color: black;')

                break

            except:
                time.sleep(5)
        else:
            pass
            # print('********无法创建该项目，请检查项目编号是否正确，该份报告终止上传********', flush=True)
            # sys.exit(0)  # Exit the program immediately
    def insertimage(self,file,xpath1,xpath2):
        driver = self.driver
        image_file_path = self.image_file_path

        with open(os.path.join(image_file_path, file), 'rb') as f:
            image_data = f.read()
            # 将图片数据编码为 base64 字符串
        image_base64 = base64.b64encode(image_data).decode('utf-8')
        # 将 base64 字符串插入到 HTML 代码中
        html_str = f'<img src="data:image/jpeg;base64,{image_base64}">'

        image_element = driver.find_element(By.XPATH, xpath1)
        ###删除
        driver.execute_script("arguments[0].remove()", image_element)

        ###添加
        start_offset = 0
        end_offset = 0
        #####图片删除后path没了，要重新看新的xpayh是什么
        element = driver.find_element(By.XPATH,xpath2)
        driver.execute_script("""
                     var range_obj = document.createRange();
                     range_obj.setStart(arguments[0], arguments[1]);
                     range_obj.setEnd(arguments[0], arguments[2]);
                     var fragment = range_obj.createContextualFragment(arguments[3]);
                     range_obj.insertNode(fragment);
                 """, element, start_offset, end_offset, html_str)
    def insertimage2(self,file,xpath):
        driver = self.driver
        image_file_path = self.image_file_path
        # image_file_path = os.path.abspath('image_file_path')
        with open(os.path.join(image_file_path, file), 'rb') as f:
            image_data = f.read()
            # 将图片数据编码为 base64 字符串
        image_base64 = base64.b64encode(image_data).decode('utf-8')
        # 将 base64 字符串插入到 HTML 代码中
        html_str = f'<img src="data:image/jpeg;base64,{image_base64}">'
        ###添加
        start_offset = 0
        end_offset = 0
        #####图片删除后path没了，要重新看新的xpayh是什么
        element = driver.find_element(By.XPATH,xpath)
        driver.execute_script("""
                                                                                                               var range_obj = document.createRange();
                                                                                                               range_obj.setStart(arguments[0], arguments[1]);
                                                                                                               range_obj.setEnd(arguments[0], arguments[2]);
                                                                                                               var fragment = range_obj.createContextualFragment(arguments[3]);
                                                                                                               range_obj.insertNode(fragment);
                                                                                                           """, element,
                              start_offset, end_offset, html_str)
    def read_pptfile(self):
        ppt_file = self.ppt_file
        image_file_path = self.image_file_path
        print("提取QCs ppt信息", flush=True)
        prs = Presentation(ppt_file)
        slide_dic = {}
        content = {}
        for index, slide in enumerate(prs.slides):
            sortedShapes = sorted(slide.shapes, key=lambda x: (x.left))
            slide_height = prs.slide_height / 914400 * 2.54
            ####遍历ppt的第二页内容,获取截图
            if index==1:
                content[index] = {}
                content[index]['QC'] = {}
                content[index]['QC']['index'] = index
                order = index
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        text = text_frame.text
                        if 'QCs' in text:
                            try:
                                top = (shape.top + shape.height) / 914400 * 2.54
                                left_position = 0
                                top_position = top / slide_height
                                bottom = 16.5
                                bottom_position = bottom / slide_height
                                right_position = 1
                                image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                img = Image.open(image_name)
                                img_size_width, img_size_height = img.size
                                crop_left_position = left_position * img_size_width
                                crop_top_position = top_position * img_size_height
                                crop_right_position = right_position * img_size_width
                                crop_bottom_position = bottom_position * img_size_height
                                name = 'QC' + '_' + str(order) + ".jpg"
                                crop_image_path = os.path.join(image_file_path, name)
                                crop_img = img.crop(
                                    (crop_left_position, crop_top_position, crop_right_position,
                                     crop_bottom_position))
                                if crop_img.mode == 'RGBA':
                                    crop_img = crop_img.convert("RGB")
                                crop_img.save(crop_image_path, quality=95, subsampling=2)
                            except:
                                print('检查QC图片是否有问题')
                                pass
            for shape in sortedShapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text = text_frame.text
                    if '项目编号' in text:
                        #####遍历含有项目编号的第一页的内容，获得相关信息
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame
                                # 读取文本框内容
                                text = text_frame.text
                                if '项目编号' in text:
                                    try:
                                        try:
                                            ProjectName = text.split(":")[1].strip()
                                            slide_dic['ProjectName'] = ProjectName
                                        except:
                                            ProjectName = text.split("：")[1].strip()
                                            slide_dic['ProjectName'] = ProjectName
                                    except:
                                        print('检查项目编号处内容', flush=True)
                                        slide_dic['ProjectName'] = ''
                                if len(text.splitlines()) == 3:
                                    try:
                                        code = text.splitlines()[0].split(",")[0]
                                        slide_dic['code'] = code
                                    except:
                                        slide_dic['code'] = ''
                                    try:
                                        条目名称 = text.splitlines()[0].split(",")[1].strip()
                                        slide_dic['条目名称'] = 条目名称
                                    except:
                                        slide_dic['条目名称'] = ''
                                    try:
                                        ProteinName = text.splitlines()[1].strip()
                                        slide_dic['ProteinName'] = ProteinName
                                    except:
                                        slide_dic['ProteinName'] = ''
                    if 'QC2: LC-MS (zoom in)' in text:
                        content[order]['QC'][index] = index
                        # print(index)
                        # if str(index - 1) != str(order):

                            # print('QC2: LC-MS (zoom in)包含多张图片', flush=True)
                        for shape in prs.slides[index].shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame
                                # 读取文本框内容
                                text = text_frame.text
                                try:
                                    result_top = ''
                                    conclusion_top = ''
                                    if "zoom in" in text:
                                        for shape in slide.shapes:
                                            if shape.has_text_frame:
                                                text_frame = shape.text_frame
                                                text = text_frame.text
                                                if "Batch No" in text:
                                                    result_top = (shape.height + shape.top) / 914400 * 2.54
                                                    # print(result_top)
                                                if "Conclusions" in text:
                                                    conclusion_top = shape.top / 914400 * 2.54

                                        image_name = os.path.join(image_file_path, "幻灯片{}.PNG".format(index + 1))
                                        left_position = 0
                                        top_position = float(result_top) / slide_height
                                        # print(conclusion_top)
                                        if conclusion_top == '':
                                            conclusion_top = 16.5
                                            bottom_position = float(conclusion_top) / slide_height
                                        else:
                                            bottom_position = float(conclusion_top) / slide_height
                                        right_position = 1
                                        img = Image.open(image_name)
                                        img_size_width, img_size_height = img.size
                                        crop_left_position = left_position * img_size_width
                                        crop_top_position = top_position * img_size_height
                                        crop_right_position = right_position * img_size_width
                                        crop_bottom_position = bottom_position * img_size_height
                                        # print(order, index)
                                        name = 'QCzoomin' + '_' + str(order) + '_' + str(index) + ".jpg"
                                        crop_image_path = os.path.join(image_file_path, name)
                                        # print(crop_image_path)
                                        crop_img = img.crop(
                                            (crop_left_position, crop_top_position, crop_right_position,
                                             crop_bottom_position))
                                        if crop_img.mode == 'RGBA':
                                            crop_img = crop_img.convert("RGB")
                                        crop_img.save(crop_image_path, quality=95, subsampling=2)
                                except:
                                    pass
                                    print('检查QCzoom格式是否正确', flush=True)
        # print(slide_dic)
        # print(content)
        #####获取步骤
        step_dic = {}
        for key, value in content.items():
            for k, v in value.items():
                try:
                    i = int(v['index'])
                    step_dic[i] = k
                except:
                    print('********注意检查' + k + ',该步骤没有抓取到相关信息')
                    pass
        step_order = sorted(step_dic.items(), key=lambda x: x[0], reverse=False)
        step_list = []
        for i in step_order:
            step = i[1]
            step_list.append(step)
        # print(step_list)

        ####获取调整过的序号--content字典
        content_list = []
        content_dic = {}
        try:
            for key, value in content.items():
                content_list.append(value)
            for index, i in enumerate(content_list):
                content_dic[index + 1] = i
        except:
            pass
        # print(content_dic)
        content_file = os.path.join(image_file_path, 'content.json')
        with open(content_file, "w", encoding="utf-8") as f:
            json.dump(content_dic, f, indent=4, ensure_ascii=False)
        return slide_dic, step_list, content_dic
    def upload_ELN(self):

        slide_dic, step_list, content_dic = self.read_pptfile()
        driver = self.driver
        driver.get(self.website)
        driver.maximize_window()
        account=self.account
        password=self.password
        image_file_path = self.image_file_path
        # image_file_path = os.path.abspath('image_file_path')
        ppt_file = self.ppt_file
        name = self.name
    #########登录
        # 等待输入框加载，输入账号

        self.sendkeys(r'//*[@id="__apppoint"]/table/tbody/tr[2]/td/table/tbody/tr[2]/td[2]/input',account)

        # 等待输入框加载，输入密码

        self.sendkeys(r'//*[@id="__apppoint"]/table/tbody/tr[2]/td/table/tbody/tr[3]/td[2]/input', password+ "\n")


        # print("登录成功",flush=True)

        ####测试用直接点进去
        #####点击My Notebooks
        '''
        self.clickoption('//*[@id="__apppoint"]/table/tbody/tr/td[1]/table/tbody/tr/td/div/table/tbody/tr[3]/td/div/div[1]/div[1]/span')
        # 点击 elntest
        self.clickoption('//*[@id="__apppoint"]/table/tbody/tr/td[1]/table/tbody/tr/td/div/table/tbody/tr[3]/td/div/div[1]/div[2]/div[1]/div')

        # self.clickoption('// *[ @ id = "__apppoint"] / table / tbody / tr / td[1] / table / tbody / tr / td / div / table / tbody / tr[3] / td / div / div[1] / div[2] / div[17] / div')

        #####创建功能可以实现,直接点进去,不重复创建
        self.clickoption('//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[1]/table/tbody/tr[2]/td/div/div[1]/div/table/tbody/tr[4]/td[6]/div/a')

        '''
        ######创建project
        ######创建项目
        ##点击Home
        try:
            self.clickoption(r'// *[ @ id = "__scil_toolbar_left"] / table / tbody / tr / td[2] / table / tbody / tr[2] / td')
            #####点击New进行创建
            self.clickoption(r'// *[ @ id = "div"] / table / tbody / tr[3] / td[3] / table / tbody / tr / td / div / a[1]')
        except  Exception:
            print('********网络无法响应，该份报告终止上传********', flush=True)
            sys.exit(0)
        ######框里添加内容
        try:
            self.sendkeys(r'/ html / body / table[4] / tbody / tr[2] / td / div / div / table / tbody / tr[2] / td[2] / input', slide_dic['ProjectName']+"\n")
        except  Exception:
            print('********1检查项目编号是否有问题********', flush=True)
            sys.exit(0)
        try:
            driver = self.driver
            for i in range(5):
                try:
                    WebDriverWait(driver, 1).until(
                        EC.presence_of_element_located(
                            (By.XPATH, r'/html/body/div[5]/div[1]'))).click()
                    time.sleep(1)
                    break

                except:
                    time.sleep(0.5)
            else:
                pass

        except  Exception:
            pass
        try:
            self.sendkeys(r'/html/body/table[4]/tbody/tr[2]/td/div/div/table/tbody/tr[3]/td[2]/input', 'Protein purification'+ "\n")
        except  Exception:
            print('********2检查项目编号是否有问题********', flush=True)
            sys.exit(0)
        try:
        #####点击创建、关闭
            self.clickoption(r'/ html / body / table[4] / tbody / tr[2] / td / div / div / table / tbody / tr[5] / td[2] / button')
        except  Exception:
            print('********网络无法响应，该份报告终止上传********', flush=True)
            sys.exit(0)

        time.sleep(3)
        for i in range(20):
            try:
                ok = driver.find_elements(By.XPATH, r"/ html / body / table[6] / tbody / tr[2] / td / div / div / div")
                time.sleep(1)
                break
            except:

                time.sleep(5)
        else:
            print('********网络无法响应，无法点击创建项目，该份报告终止上传********', flush=True)
            sys.exit(0)  # Exit the program immediately

        status=''
        if len(ok) == 0:
            status=True
            print('项目创建成功', flush=True)
            time.sleep(2)

        if len(ok) != 0:
            oktext=WebDriverWait(driver, 5).until(EC.presence_of_element_located((By.XPATH,r'/ html / body / table[6] / tbody / tr[2] / td / div / div / div[1]'))).text
            if oktext=='[ERROR]: One project is allowed only one notebook for one user':
                status = True
                try:
                    self.clickoption(r'/ html / body / table[6] / tbody / tr[2] / td / div / div / div[2] / button')

                    self.clickoption(r'/ html / body / table[4] / tbody / tr[1] / td[2] / img')

                    self.sendkeys(
                        r'// *[ @ id = "div"] / table / tbody / tr[1] / td / table / tbody / tr[1] / td[1] / input',
                        slide_dic['ProjectName'])

                    self.clickoption(
                        r'// *[ @ id = "div"] / table / tbody / tr[1] / td / table / tbody / tr[1] / td[2] / span')

                    time.sleep(5)

                    self.clickoption(
                        r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / select[6]')
                    self.clickoption(
                        r' // *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / select[6] / option[4]')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)
                table = driver.find_element(By.XPATH,
                                            r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[1] / div / table')
                #####查找项目，先找第一页的并循环刷新多次，如果没有就翻页
                try:
                    ######这个查找同时包含名字和项目编号的内容
                    ######先查找第一页是否含有含姓名和项目编号的内容，并刷新循环多次，如果没有的话就来翻页查找
                    for i in range(20):
                        # print('1')
                        row = table.find_element(By.XPATH,
                                                 r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(
                                                     slide_dic['ProjectName'], name))
                        # print('2')
                        break
                    else:
                        try:
                            for i in range(20):
                                try:
                                    self.clickoption(
                                        r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / img[1]')
                                    # print('3')
                                    row = table.find_element(By.XPATH,
                                                             r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(
                                                                 slide_dic['ProjectName'], name))
                                    break
                                except:
                                    time.sleep(2)
                            else:
                                print('********1网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                                sys.exit(0)  # Exit the program immediately
                        except:
                            print('********2网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                            sys.exit(0)  # Exit the program immediately

                except:
                    #####翻页查找内容
                    for i in range(20):
                        try:
                            self.clickoption(
                                r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[2] / div[contains(text(), "Next")]')
                            table = driver.find_element(By.XPATH,
                                                        r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[1] / div / table')
                            row = table.find_element(By.XPATH,
                                                     r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(
                                                         slide_dic['ProjectName'], name))
                            # print(i,flush=True)
                            break
                        except:
                            time.sleep(2)
                    else:
                        print('********3网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                time.sleep(3)
                for i in range(20):
                    try:
                        WebDriverWait(row, 5).until(
                            EC.presence_of_element_located(
                                (By.XPATH, r'./td[6]'))
                        ).click()
                        # print('点击项目成功',flush=True)
                        time.sleep(1)
                        break
                    except:
                        time.sleep(5)
                else:
                    print('********4网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately

            if oktext == '[ERROR]: No privilege to create experiment in the project':
                status = False
                print('********无法创建该项目，请检查项目编号是否正确，该份报告终止上传********', flush=True)
                pass

        if status==True:
            #####创建新表单
            ####新表单名字
            #
            # # 点击加号
                try:
                    self.clickoption(r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div/table/tbody/tr[1]/td/span[4]/img')

                    #####点击Biortus
                    self.clickoption(r'/html/body/table[4]/tbody/tr/td/div/div[3]')

                    #######点击纯化  xpath标签总是会变,用last函数
                    self.clickoption(r'/html/body/table[last()]/tbody/tr/td/div/div[7]')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)
                #######part1创建表单填写表格内容
                try:
                    NewEnityname=slide_dic['条目名称']+ ' ' + slide_dic['ProteinName'] +' QCs'
                except:
                    NewEnityname=' '
                    print('注意检查表单名称',flush=True)

                ######输入选项输入创建表单名称
                try:
                    self.sendkeys(r'/ html / body / table[6] / tbody / tr[2] / td / div / div / table / tbody / tr[4] / td[2] / input', NewEnityname)
                except:
                    self.sendkeys(r'/ html / body / table[6] / tbody / tr[2] / td / div / div / table / tbody / tr[4] / td[2] / input',' ')
                    print('检查表单名称',flush=True)

                #####点击创建
                try:
                    self.clickoption(r'/html/body/table[6]/tbody/tr[2]/td/div/div/table/tbody/tr[6]/td[2]/button')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)

                ######填写表格内容
                ####添加title

                try:
                    self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[1]/div',slide_dic['条目名称'])
                except:

                    print('注意检查蛋白批次号',flush=True)
                    pass

                time.sleep(3)
                # print("表格填写完毕",flush=True)
                try:
                #######part2复制 cart并改名字
                    self.clickoption(r'//*[@id="__scil_toolbar_right"]/table/tbody/tr/td[2]/table/tbody/tr[1]/td/img')
                    ######点击show 加上关键词匹配
                    self.clickoption(r'//span[contains(text(), "Show")]')
                    #####关闭

                    self.clickoption(r' / html / body / table[last()] / tbody / tr[1] / td[2] / img')
                    ######再次点击
                    self.clickoption(r'//*[@id="__scil_toolbar_right"]/table/tbody/tr/td[2]/table/tbody/tr[1]/td/img')

                    ######点击show 加上关键词匹配
                    self.clickoption(r'//span[contains(text(), "Show")]')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)

                #####从页面中选择table
                table= driver.find_element(By.XPATH,r"/html/body/table[last()]/tbody/tr[2]/td/div/div/table/tbody/tr[1]/td[2]/div/div/table")

                #####根据列表元素复制模板
                # print(step_list,flush=True)

                for step in step_list:
                    for i in range(20):
                        try:
                            row=WebDriverWait(table, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH, r'.//tr[td[7]="{}"]'.format(step)))
                            )
                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********1网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately
                    time.sleep(2)

                    for i in range(20):
                        try:
                            WebDriverWait(row, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH, r'./td[2]'))
                            ).click()
                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********2网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately
                    time.sleep(2)
                    ####复制
                    try:
                        self.clickoption(r'/html/body/table[last()]/tbody/tr[2]/td/div/div/table/tbody/tr[3]/td[2]/button[1]')
                    except  Exception:
                        print('********网络无法响应，该份报告终止上传********', flush=True)
                        sys.exit(0)
                    time.sleep(1)
                    for i in range(20):
                        try:
                            WebDriverWait(row, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH, r'./td[2]'))
                            ).click()
                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********3网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately
                    time.sleep(2)
                ######关闭
                try:
                    self.clickoption(r'/html/body/table[last()]/tbody/tr[1]/td[2]/img')

                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)
                # print("复制模板完毕",flush=True)
                ###复制过来后改名字
                # print("根据步骤修改标题",flush=True)

                if len(step_list)==1:

                    for i in range(20):
                        try:

                            WebDriverWait(driver, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[1]/div/span'))
                            ).click()

                            time.sleep(1)
                            break
                        except:
                            time.sleep(5)
                    else:
                        print('********网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                    for i in range(20):
                        try:

                            WebDriverWait(driver, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[1]/td/img[5]'))
                            ).click()
                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********网络无法响应，无法点击修改标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                    for i in range(20):
                        try:
                            WebDriverWait(driver, 10).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     '/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[2]/input'))
                            ).clear()

                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********网络无法响应，无法删除标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately
                    ######模板名改成QC
                    for i in range(20):
                        try:
                            WebDriverWait(driver, 30).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     '/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[2]/input'))
                            ).send_keys('QC')

                            time.sleep(2)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********网络无法响应，无法修改标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                    ######save
                    for i in range(20):
                        try:
                            WebDriverWait(driver, 30).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     r'/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[3]/button'))
                            ).click()
                            time.sleep(2)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********网络无法响应，无法保存标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                else:
                    pass
                    print("QCs含有多项内容，手动补充",flush=True)

                time.sleep(2)
                try:
                #####修改完刷新一下
                    self.clickoption(r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[1]/span')
                    self.clickoption(r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[1]/td/img[1]')
                    time.sleep(2)
                    ####part3  修改模板内容，qc只含有一种情况，所以直接来写就写就行，把之前的根据内容来匹配删掉
                    self.clickoption(r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[1]/div/span')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)

                for key, value in content_dic.items():

                    try:
                        for k, v in value.items():
                            ######修改QCs
                                time.sleep(2)
                                print("修改模板QC",flush=True)
                                iframe_list = driver.find_elements(By.XPATH,r'//iframe')
                                driver.switch_to.frame(iframe_list[2])
                                #####获取所有的qc zoomin
                                try:
                                    i_list = []
                                    for i, j in v.items():
                                        if type(i) == int:
                                            # print(i)
                                            i_list.append(i)
                                except:
                                    pass
                                # print(i_list)
                                ####删除原有相关内容
                                for i in range(1,12):
                                    element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/div/p[1]')
                                    driver.execute_script("arguments[0].remove()", element)
                                # print("删除原有相关内容",flush=True)
                                body = driver.find_element(By.TAG_NAME, 'body')
                                # 将光标移到文本末尾处
                                body.send_keys(Keys.END)
                                # 按下回车键换行
                                body.send_keys(Keys.ENTER)
                                ###插入QC第一张
                                try:
                                    image = 'QC_' + str(v['index']) + '.jpg'

                                    self.insertimage2(image, r'// *[ @ id = "tinymce"] / div[1]')
                                except:
                                    print('检查QC1,QC2,QC3图片是否有问题', flush=True)
                                    pass
                                ####插入QCzoomin 第一张
                                try:
                                    body = driver.find_element(By.TAG_NAME, 'body')
                                    # 将光标移到文本末尾处
                                    body.send_keys(Keys.END)
                                    # 按下回车键换行
                                    body.send_keys(Keys.ENTER)
                                    image = 'QCzoomin_' + str(v['index']) + '_' + str(i_list[0]) +  '.jpg'
                                    self.insertimage2(image, '// *[ @ id = "tinymce"] / div[2]')
                                except:
                                    print('检查QC (zoom in) 图片', flush=True)
                                    pass
                                ####如果有多张图插入QCzoomout 第三张
                                try:
                                    if len(i_list)>1:
                                        # print('需要添加图片')
                                        for i in range(1, len(i_list)):
                                            body = driver.find_element(By.TAG_NAME, 'body')
                                            # 将光标移到文本末尾处
                                            body.send_keys(Keys.END)
                                            # 按下回车键换行
                                            body.send_keys(Keys.ENTER)
                                            image = 'QCzoomin_' + str(v['index']) + '_' + str(i_list[i]) + '.jpg'
                                            self.insertimage2(image, r'// *[ @ id = "tinymce"] / div[{}]'.format(int(2+i)))
                                except:
                                    print('检查QC (zoom in) 图片', flush=True)
                                driver.switch_to.default_content()
                                self.saveoption(slide_dic['条目名称'])
                                print("模板QC内容修改成功",flush=True)
                    except:
                        pass
                print("完成修改模板内容任务",flush=True)
                try:
                    destination_path = r'完成上传'
                    shutil.copy2(ppt_file, destination_path)
                    # 删除源文件
                    os.remove(ppt_file)

                    time.sleep(1)
                    os.remove(image_file_path)
                except:
                    pass

        else:
            try:
                destination_path = r'项目编号有问题'
                shutil.copy2(ppt_file, destination_path)
                # 删除源文件
                os.remove(ppt_file)

            except:
                pass
    def convert_string_to_list(self, string):
        try:
            result = string.replace(',', '+').replace('，','+').split('+')  # 将中文逗号替换为英文逗号，并以加号为分隔符拆分字符串
        except:
            pass
            print('检查质粒编号', flush=True)
        return result
    def read_pptfile_wb(self):
        ppt_file = self.ppt_file
        # image_file_path = os.path.abspath('image_file_path')
        image_file_path = self.image_file_path
        print("提取WB ppt信息", flush=True)
        prs = Presentation(ppt_file)
        content = {}
        slide_dic = {}
        sample_list = []
        table_str = ''
        table_html = ''
        keywords_list = ["Sample", "Plasmid", "Theoretical", "A280", "PI"]
        keyword_match_num = 0
        for index, slide in enumerate(prs.slides):
            sortedShapes = sorted(slide.shapes, key=lambda x: (x.left))
            slide_height = prs.slide_height / 914400 * 2.54
            for shape in sortedShapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text = text_frame.text
                    if '项目编号' in text:
                        #####遍历含有项目编号的第一页的内容，获得相关信息
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame
                                # 读取文本框内容
                                text = text_frame.text
                                if '项目编号' in text:
                                    try:
                                        try:
                                            ProjectName = text.split(":")[1].strip()
                                            slide_dic['ProjectName'] = ProjectName
                                        except:
                                            ProjectName = text.split("：")[1].strip()
                                            slide_dic['ProjectName'] = ProjectName
                                    except:
                                        print('检查项目编号处内容', flush=True)
                                        slide_dic['ProjectName'] = ''
                                if 'Blot' in text:
                                    try:
                                        条目名称 = text.splitlines()[0].strip()
                                        slide_dic['条目名称'] = 条目名称
                                    except:
                                        slide_dic['条目名称'] = ''
                                    try:
                                        left = text.split('')[1:]
                                        left = ' '.join(left)
                                        slide_dic['left'] = left
                                        slide_dic['title'] = text.splitlines()[0].strip().split('-')[1] + '-' + \
                                                             text.splitlines()[0].strip().split('-')[0] + ' ' + left
                                    except:
                                        slide_dic['left'] = ''
                                        slide_dic['title'] = ''
                    if 'Western Blot Test' in text:
                        antibody_text = text
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text_frame = shape.text_frame
                                # 读取文本框内容
                                text = text_frame.text
                                if 'Experiment' in text:
                                    content[index] = {}
                                    content[index]['western blot'] = {}
                                    content[index]['western blot']['index'] = index

                                    try:
                                        sample = text.split("Sample:")[1].split('Exposure time')[0].strip()
                                        # print(Sample)
                                        content[index]['western blot']['sample'] = sample
                                        batch = sample.split(',')[0].split('.')[0].strip()
                                        content[index]['western blot']['batch'] = batch
                                    except:
                                        sample = ''
                                        content[index]['western blot']['sample'] = sample
                                        batch = ''
                                        content[index]['western blot']['batch'] = batch
                                    try:
                                        antibody = antibody_text.split('(')[1].split(')')[0]
                                        content[index]['western blot']['antibody'] = antibody
                                    except:
                                        antibody = ''
                                        content[index]['western blot']['antibody'] = antibody
                                    try:
                                        Exposure_time = text.split("Exposure time:")[1].split('Results:')[0].strip()
                                        # print(Sample)
                                        content[index]['western blot']['Exposure_time'] = Exposure_time
                                    except:
                                        Exposure_time = ''
                                        content[index]['western blot']['Exposure_time'] = Exposure_time


                                    try:
                                        image_name = os.path.join(image_file_path,
                                                                  "幻灯片{}.PNG".format(index + 1))
                                        left_position = 0

                                        top_position = 0
                                        conclusion_top = 16.5

                                        bottom_position = float(conclusion_top) / slide_height
                                        right_position = 1
                                        img = Image.open(image_name)
                                        img_size_width, img_size_height = img.size
                                        crop_left_position = left_position * img_size_width
                                        crop_top_position = top_position * img_size_height
                                        crop_right_position = right_position * img_size_width
                                        crop_bottom_position = bottom_position * img_size_height

                                        name = 'western blot' + '_' + str(index) + ".jpg"
                                        crop_image_path = os.path.join(image_file_path, name)

                                        crop_img = img.crop(
                                            (crop_left_position, crop_top_position, crop_right_position,
                                             crop_bottom_position))
                                        if crop_img.mode == 'RGBA':
                                            crop_img = crop_img.convert("RGB")
                                        crop_img.save(crop_image_path, quality=95, subsampling=2)
                                    except:
                                        pass
                                        print('检查western blot 图片', flush=True)
            ####获取第二页的表格信息
            if index == 1:
                for shape in slide.shapes:
                    name_mark = 0
                    if index == 1:
                        for shape in slide.shapes:
                            if isinstance(shape, pptx.shapes.graphfrm.GraphicFrame):
                                # 表格问题,出现表格问题跳过
                                try:
                                    keyword_match_num = 0
                                    for i in range(len(shape.table.columns)):
                                        for j in keywords_list:
                                            if j.upper() in shape.table.cell(0, i).text.upper():
                                                keyword_match_num += 1
                                except:
                                    break
                    if keyword_match_num >= 3:
                        table_data = []
                        if shape.has_table:
                            # 获取表格对象
                            table = shape.table
                            # 创建一个列表用于保存表格内容
                            # 遍历表格中的所有行和列
                            for i, row in enumerate(table.rows):
                                row_data = []
                                for j, cell in enumerate(row.cells):
                                    # 获取单元格中的文本
                                    cell_text = cell.text_frame.text.strip()
                                    row_data.append(cell_text)
                                # 将该行添加到表格内容列表中
                                table_data.append(row_data)
                            sample_list.append(table_data)
                            for row in shape.table.rows:
                                row_html = ''
                                for cell in row.cells:
                                    row_html += f'<td>{cell.text_frame.text}</td>'
                                table_html += f'<tr>{row_html}</tr>'
                            table_str += f'<table>{table_html}</table>'
                #####获取表格信息

        if table_str == '':
            print('检查该报告第二页是否含有所有表格信息,并手动补充', flush=True)
        #####获取步骤
        # print(content)
        # print(slide_dic)
        step_dic = {}
        for key, value in content.items():
            for k, v in value.items():
                try:
                    i = int(v['index'])
                    step_dic[i] = k
                except:
                    print('********注意检查' + k + ',该步骤没有抓取到相关信息')
                    pass
        step_order = sorted(step_dic.items(), key=lambda x: x[0], reverse=False)
        step_list = []
        step_antibody = []
        for i in step_order:
            # print(i)
            step = i[1]
            antibody = content[i[0]][step]['antibody']
            step_list.append(step)
            if antibody == '':
                step_antibody.append(step)
            else:
                step_antibody.append(step + '-' + antibody)

        ####获取带有标题名的的内容
        # print(step_antibody)
        ####获调整过的序号--content字典
        content_list = []
        content_dic = {}
        try:
            for key, value in content.items():
                content_list.append(value)
            for index, i in enumerate(content_list):
                content_dic[index + 1] = i
        except:

            pass

        # print('获取到ppt的相关信息：'+str(content_dic))
        content_file = os.path.join(image_file_path, 'content.json')
        with open(content_file, "w", encoding="utf-8") as f:
            json.dump(content_dic, f, indent=4, ensure_ascii=False)

        return slide_dic, step_list, sample_list, content_dic, step_antibody, table_str
    def saveoption(self,text):
        driver = self.driver
        for i in range(20):
            try:
                double_click_element = driver.find_element(By.XPATH,
                                                           r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[1]/div')
                time.sleep(1)
                actions = ActionChains(driver)
                # 双击需要输入文本的框
                actions.double_click(double_click_element).perform()
                # 在框中输入文本
                active_element = driver.switch_to.active_element
                active_element.send_keys(text)
                driver.switch_to.default_content()
                break

            except:
                time.sleep(5)
        else:
            print('********无法创建该项目，请检查项目编号是否正确，该份报告终止上传********', flush=True)
            sys.exit(0)  # Exit the program immediately

        ###点击save
        time.sleep(3)
        for i in range(20):
            try:
                element = WebDriverWait(driver, 10).until(
                    EC.presence_of_element_located(
                        (By.XPATH,
                         r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[1]/td/img[3]'))
                )
                actions = ActionChains(driver)
                for i in range(5):
                    actions.move_to_element(element).perform()
                    actions.click(element).perform()
                    time.sleep(0.1)
                for i in range(5):
                    element.click()
                    time.sleep(0.1)
                time.sleep(2)

                break

            except:
                time.sleep(5)
        else:
            print('********无法创建该项目，请检查项目编号是否正确，该份报告终止上传********', flush=True)
            sys.exit(0)  # Exit the program immediately
    def upload_ELN_wb(self):
        slide_dic, step_list, sample_list, content_dic,step_antibody,table_str = self.read_pptfile_wb()
        driver = self.driver
        driver.get(self.website)
        driver.maximize_window()
        account = self.account
        password = self.password
        image_file_path=self.image_file_path
        # image_file_path = os.path.abspath('image_file_path')
        ppt_file = self.ppt_file
        name = self.name
        #########登录
        # 等待输入框加载，输入账号
        try:
            self.sendkeys(r'//*[@id="__apppoint"]/table/tbody/tr[2]/td/table/tbody/tr[2]/td[2]/input', account + "\n")

            # 等待输入框加载，输入密码
            self.sendkeys(r'//*[@id="__apppoint"]/table/tbody/tr[2]/td/table/tbody/tr[3]/td[2]/input', password + "\n")
        except:
            print('********检查项目编号是否有问题********', flush=True)
            sys.exit(0)

        # print("登录成功", flush=True)

        ####测试用直接点进去
        #####点击My Notebooks
        '''
        self.clickoption(
            '//*[@id="__apppoint"]/table/tbody/tr/td[1]/table/tbody/tr/td/div/table/tbody/tr[3]/td/div/div[1]/div[1]/span')
        # 点击 elntest
        self.clickoption(
            '//*[@id="__apppoint"]/table/tbody/tr/td[1]/table/tbody/tr/td/div/table/tbody/tr[3]/td/div/div[1]/div[2]/div[1]/div')

        # self.clickoption('// *[ @ id = "__apppoint"] / table / tbody / tr / td[1] / table / tbody / tr / td / div / table / tbody / tr[3] / td / div / div[1] / div[2] / div[17] / div')

        #####创建功能可以实现,直接点进去,不重复创建
        self.clickoption(
            '//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[1]/table/tbody/tr[2]/td/div/div[1]/div/table/tbody/tr[4]/td[6]/div/a')
        '''
        ######创建project
        ######创建项目
        ##点击Home
        try:
            self.clickoption(r'// *[ @ id = "__scil_toolbar_left"] / table / tbody / tr / td[2] / table / tbody / tr[2] / td')
            #####点击New进行创建
            self.clickoption(r'// *[ @ id = "div"] / table / tbody / tr[3] / td[3] / table / tbody / tr / td / div / a[1]')
            ######框里添加内容
            self.sendkeys(r'/ html / body / table[4] / tbody / tr[2] / td / div / div / table / tbody / tr[2] / td[2] / input', slide_dic['ProjectName'])
        except  Exception:
            print('********检查项目编号是否有问题********', flush=True)
            sys.exit(0)

        try:
            driver = self.driver
            for i in range(5):
                try:
                    WebDriverWait(driver, 1).until(
                        EC.presence_of_element_located(
                            (By.XPATH, r'/html/body/div[5]/div[1]'))).click()
                    time.sleep(1)
                    break

                except:
                    time.sleep(0.5)
            else:
                pass

        except  Exception:
            pass

        try:
            self.sendkeys(r'/html/body/table[4]/tbody/tr[2]/td/div/div/table/tbody/tr[3]/td[2]/input',
                          'Protein purification' + "\n")
            #####点击创建、关闭
            self.clickoption(
                r'/ html / body / table[4] / tbody / tr[2] / td / div / div / table / tbody / tr[5] / td[2] / button')
        except:
            print('********检查项目编号是否有问题********', flush=True)
            sys.exit(0)

        time.sleep(3)
        for i in range(20):
            try:
                ok = driver.find_elements(By.XPATH, r"/ html / body / table[6] / tbody / tr[2] / td / div / div / div")
                time.sleep(1)
                break
            except:

                time.sleep(5)
        else:
            print('*********网络无法响应，无法点击创建项目，该份报告终止上传********', flush=True)
            sys.exit(0)  # Exit the program immediately
        status = ''
        if len(ok) == 0:
            status = True
            print('项目创建成功', flush=True)
            time.sleep(2)

        if len(ok) != 0:
            oktext=WebDriverWait(driver, 5).until(EC.presence_of_element_located((By.XPATH,r'/ html / body / table[6] / tbody / tr[2] / td / div / div / div[1]'))).text
            if oktext=='[ERROR]: One project is allowed only one notebook for one user':
                # print('项目已存在',flush=True)
                status = True
                try:
                    self.clickoption(r'/ html / body / table[6] / tbody / tr[2] / td / div / div / div[2] / button')

                    self.clickoption(r'/ html / body / table[4] / tbody / tr[1] / td[2] / img')

                    self.sendkeys(r'// *[ @ id = "div"] / table / tbody / tr[1] / td / table / tbody / tr[1] / td[1] / input', slide_dic['ProjectName'])

                    self.clickoption(r'// *[ @ id = "div"] / table / tbody / tr[1] / td / table / tbody / tr[1] / td[2] / span')

                    time.sleep(5)

                    self.clickoption(r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / select[6]')
                    self.clickoption(r' // *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / select[6] / option[4]')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)
                table = driver.find_element(By.XPATH,r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[1] / div / table')
                #####查找项目，先找第一页的并循环刷新多次，如果没有就翻页
                try:
                ######这个查找同时包含名字和项目编号的内容
                    ######先查找第一页是否含有含姓名和项目编号的内容，并刷新循环多次，如果没有的话就来翻页查找
                    for i in range(20):
                        # print('1')
                        row = table.find_element(By.XPATH,r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(slide_dic['ProjectName'], name))
                        # print('2')
                        break
                    else:
                        try:
                            for i in range(20):
                                try:
                                    self.clickoption(r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[1] / td / img[1]')
                                    # print('3')
                                    row = table.find_element(By.XPATH,r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(slide_dic['ProjectName'], name))
                                    break
                                except:
                                    time.sleep(2)
                            else:
                                print('********1网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                                sys.exit(0)  # Exit the program immediately

                        except:
                            print('********2网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                            sys.exit(0)  # Exit the program immediately
                except:
                    #####翻页查找内容
                    for i in range(20):
                        try:
                            self.clickoption(r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[2] / div[contains(text(), "Next")]')
                            table = driver.find_element(By.XPATH,r'// *[ @ id = "__apppoint"] / table / tbody / tr / td[3] / table / tbody / tr[2] / td / div / table / tbody / tr[2] / td / div / div[1] / div / table')
                            row = table.find_element(By.XPATH,r'.//tr[contains(td[12], "{}")][contains(td[8], "{}")]'.format(slide_dic['ProjectName'], name))
                            # print(i,flush=True)
                            break
                        except:
                            time.sleep(2)
                    else:
                        print('********3网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately

                time.sleep(3)
                for i in range(20):
                    try:
                        WebDriverWait(row, 5).until(
                            EC.presence_of_element_located(
                                (By.XPATH, r'./td[6]'))
                        ).click()
                        # print('点击项目成功',flush=True)
                        time.sleep(1)
                        break
                    except:
                        time.sleep(5)
                else:
                    print('********4网络无法相应，无法点进个人项目页，该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately
            if oktext == '[ERROR]: No privilege to create experiment in the project':
                status = False
                print('********无法创建该项目，请检查项目编号是否正确，该份报告终止上传********', flush=True)
                pass
        if status == True:
            #####创建新表单
            ####新表单名字
            #
            # # 点击加号
                try:
                    self.clickoption(r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div/table/tbody/tr[1]/td/span[4]/img')

                    #####点击Biortus
                    self.clickoption(r'/html/body/table[4]/tbody/tr/td/div/div[3]')

                    #######点击纯化  xpath标签总是会变,用last函数
                    self.clickoption(r'/html/body/table[last()]/tbody/tr/td/div/div[7]')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)

                #######part1创建表单填写表格内容
                try:
                    NewEnityname=slide_dic['title']
                except:
                    NewEnityname=' '
                    print('注意检查表单名称',flush=True)

                ######输入选项输入创建表单名称
                try:
                    self.sendkeys(r'/ html / body / table[6] / tbody / tr[2] / td / div / div / table / tbody / tr[4] / td[2] / input', NewEnityname)
                except:
                    self.sendkeys(r'/ html / body / table[6] / tbody / tr[2] / td / div / div / table / tbody / tr[4] / td[2] / input',' ')
                    print('检查表单名称',flush=True)
                try:
                #####点击创建
                    self.clickoption(r'/html/body/table[6]/tbody/tr[2]/td/div/div/table/tbody/tr[6]/td[2]/button')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)

                ######填写表格内容
                ####添加title
                try:
                    self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[1]/div',slide_dic['条目名称'].split('-')[1] + ' ' + slide_dic['条目名称'].split('-')[0])
                except:
                    print('注意检查title', flush=True)

                try:
                    #####添加Sample ID
                    sampleid = slide_dic['条目名称'].split('-')[1].replace('BP', '')
                    self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[2]/div',sampleid)
                except:
                    print('注意检查Sample ID', flush=True)

                try:
                    Reg = slide_dic['条目名称'].split("-")[1]
                    self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[3]/div',Reg)
                except:
                    print('注意检查Reg. No(Parent ID)', flush=True)

                #####添Batch No(date)
                try:
                    Batch = slide_dic['条目名称'].split("-")[0]
                    self.writetable(r'//*[@id="_eformNaN2019042926.summary"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[4]/div',Batch)
                except:
                    print('注意检查Batch No(date)', flush=True)

                try:
                    sampleid_list = self.convert_string_to_list(sampleid)
                except:
                    print('注意检查Sample ID', flush=True)
                    sampleid_list = []
                if len(sampleid_list) == 1:
                    try:
                        #####PlasmidName
                        self.writetable(
                            r'//*[@id="_eformNaN2019042926.protein1"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[1]/div',
                            sample_list[0][1][1])
                    except:
                        print('注意检查PlasmidName', flush=True)
                        pass

                    #####MW
                    try:
                        self.writetable(r'//*[@id="_eformNaN2019042926.protein1"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[2]/div',sample_list[0][1][2])
                    except:
                        print('注意检查MW', flush=True)
                        pass
                    try:
                        ######1A280
                        self.writetable(r'//*[@id="_eformNaN2019042926.protein1"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[3]/div',sample_list[0][1][3])
                    except:
                        print('注意检查1A280', flush=True)
                        pass
                    try:
                        ######P1
                        self.writetable( r'//*[@id="_eformNaN2019042926.protein1"]/table/tbody/tr[2]/td/table[2]/tbody/tr/td[4]/div',sample_list[0][1][4])
                    except:
                        print('注意检查P1', flush=True)
                        pass
                else:
                    pass

                # print("表格填写完毕", flush=True)
                #######part2复制 cart并改名字
                try:
                    self.clickoption(r'//*[@id="__scil_toolbar_right"]/table/tbody/tr/td[2]/table/tbody/tr[1]/td/img')
                    ######点击show 加上关键词匹配
                    self.clickoption(r'//span[contains(text(), "Show")]')
                    #####关闭

                    self.clickoption(r' / html / body / table[last()] / tbody / tr[1] / td[2] / img')
                    ######再次点击
                    self.clickoption(r'//*[@id="__scil_toolbar_right"]/table/tbody/tr/td[2]/table/tbody/tr[1]/td/img')

                    ######点击show 加上关键词匹配
                    self.clickoption(r'//span[contains(text(), "Show")]')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)

                #####从页面中选择table
                table = driver.find_element(By.XPATH,r"/html/body/table[last()]/tbody/tr[2]/td/div/div/table/tbody/tr[1]/td[2]/div/div/table")

                #####根据列表元素复制模板
                # print(step_list, flush=True)
                #####设置一个外层循环，如果复制成功就跳出循环
                break_outer = False
                for step in step_list:
                    execute_successfully = False
                    for i in range(20):
                        try:
                            row = WebDriverWait(table, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH, r'.//tr[td[7]="{}"]'.format(step)))
                            )
                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********1网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately
                    time.sleep(2)

                    for i in range(20):
                        try:
                            WebDriverWait(row, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH, r'./td[2]'))
                            ).click()
                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********2网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately
                    time.sleep(2)
                    ####复制
                    try:
                        self.clickoption(r'/html/body/table[last()]/tbody/tr[2]/td/div/div/table/tbody/tr[3]/td[2]/button[1]')
                    except  Exception:
                        print('********网络无法响应，该份报告终止上传********', flush=True)
                        sys.exit(0)
                    time.sleep(1)
                    for i in range(20):
                        try:
                            WebDriverWait(row, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH, r'./td[2]'))
                            ).click()
                            time.sleep(1)
                            break
                        except:
                            time.sleep(5)
                    else:
                        print('********3网络无法响应，无法复制模板该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately
                    time.sleep(2)
                    execute_successfully = True  # Set flag to indicate successful execution for this element

                    if not execute_successfully:
                        break_outer = True  # Set the flag to break the outer loop
                        break  # Break the inner loop

                    if break_outer:
                        break  # Break the outer loop

                ######关闭
                try:
                    self.clickoption(r'/html/body/table[last()]/tbody/tr[1]/td[2]/img')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)
                # print("复制模板完毕", flush=True)
                ###复制过来后改名字
                # print("根据步骤修改标题", flush=True)

                for i in range(20):
                    try:

                        NBKS = driver.find_elements(By.XPATH,
                                                    r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[last()]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div/div')

                        time.sleep(1)
                        break

                    except:
                        time.sleep(5)
                else:
                    print('********网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                    sys.exit(0)  # Exit the program immediately
                ####修改命名方式
                for j in range(1, len(NBKS) + 1):
                    try:
                        for i in range(20):
                            try:

                                WebDriverWait(driver, 5).until(
                                    EC.presence_of_element_located(
                                        (By.XPATH,
                                         r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[{}]/div/span'.format(
                                             j)))
                                ).click()

                                time.sleep(1)
                                break
                            except:
                                time.sleep(5)
                        else:
                            print('********1网络无法响应，修改标题名，该份报告终止上传********', flush=True)
                            sys.exit(0)  # Exit the program immediately

                        for i in range(20):
                            try:

                                WebDriverWait(driver, 5).until(
                                    EC.presence_of_element_located(
                                        (By.XPATH,
                                         r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[1]/td/img[5]'))
                                ).click()
                                time.sleep(1)
                                break

                            except:
                                time.sleep(5)
                        else:
                            print('********2网络无法响应，修改标题名，该份报告终止上传********', flush=True)
                            sys.exit(0)  # Exit the program immediately

                        for i in range(20):
                            try:
                                WebDriverWait(driver, 10).until(
                                    EC.presence_of_element_located(
                                        (By.XPATH,
                                         '/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[2]/input'))
                                ).clear()

                                time.sleep(1)
                                break

                            except:
                                time.sleep(5)
                        else:
                            print('********3网络无法响应，修改标题名，该份报告终止上传********', flush=True)
                            sys.exit(0)  # Exit the program immediately
                        ######文字

                        for index, i in enumerate(step_antibody):
                            # print(index, j)
                            if index+1==j:
                                name=str(j)+' '+step_antibody[index]
                                # print('name'+name, flush=True)
                                step_antibody[index] = "matched"

                                ####用last函数
                                for i in range(20):
                                    try:
                                        WebDriverWait(driver, 30).until(
                                            EC.presence_of_element_located(
                                                (By.XPATH,
                                                 '/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[2]/input'))
                                        ).send_keys(name)
                                        time.sleep(2)
                                        break

                                    except:
                                        time.sleep(5)
                                else:
                                    print('********4网络无法响应，修改标题名，该份报告终止上传********', flush=True)
                                    sys.exit(0)  # Exit the program immediately

                                    ######save
                                for i in range(20):
                                    try:
                                        WebDriverWait(driver, 30).until(
                                            EC.presence_of_element_located(
                                                (By.XPATH,
                                                 r'/html/body/table[last()]/tbody/tr[2]/td/div/table/tbody/tr/td[2]/div[3]/button'))
                                        ).click()
                                        time.sleep(2)
                                        break

                                    except:
                                        time.sleep(5)
                                else:
                                    print('********5网络无法响应，修改标题名，该份报告终止上传********', flush=True)
                                    sys.exit(0)  # Exit the program immediately
                                # 修改元素值
                                # 在下一轮匹配中跳过已匹配的元素

                                break

                            else:
                                pass

                    except:
                        pass

                time.sleep(2)
                #####修改完刷新一下
                try:
                    self.clickoption(
                        r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[1]/span')
                    self.clickoption(
                        r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[1]/td/img[1]')
                except  Exception:
                    print('********网络无法响应，该份报告终止上传********', flush=True)
                    sys.exit(0)

                time.sleep(2)

                ####part3  修改不同模板内容
                for j in range(1, len(NBKS) + 1):
                    self.clickoption(r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[{}]/div/span'.format(j))
                    time.sleep(1)
                    ######文字
                    inner2_list = []
                    for i in range(20):
                        try:
                            element1 = WebDriverWait(driver, 5).until(
                                EC.presence_of_element_located(
                                    (By.XPATH,
                                     r'//*[@id="__apppoint"]/table/tbody/tr/td[3]/table/tbody/tr[2]/td/div[2]/table/tbody/tr[2]/td/div/div[1]/div/div[2]/div[{}]/div/span'.format(j)))
                            )
                            if element1.is_displayed():
                                # 获取元素的文本信息
                                text = element1.text
                                # print(text)
                                inner2_list.append(text)
                            else:
                                # 元素不可见，获取元素内的HTML代码
                                inner2 = element1.get_attribute('innerHTML')
                                inner2_list.append(inner2)
                            time.sleep(1)
                            break

                        except:
                            time.sleep(5)
                    else:
                        print('********网络无法响应，无法获取标题名，该份报告终止上传********', flush=True)
                        sys.exit(0)  # Exit the program immediately
                    #####修改内容,根据step_antibody进行匹配
                    ####获取每一条标签的名称

                    denum = inner2_list[0].split("]")[1]
                    denum_match = inner2_list[0].split("]")[1].replace(" ", "")
                    # print(denum, flush=True)

                    for key, value in content_dic.items():
                        try:
                            for k, v in value.items():
                                # print(str(key) + ' ' + str(k) + '-' + v['antibody'])
                                ####标题中有内容和无内容放一起
                                havetext=str(key) + ' ' + str(k)
                                notext=str(key) + ' ' + str(k) + '-' + v['antibody']
                                if denum_match == havetext.replace(' ','')or denum_match == notext.replace(' ',''):
                                    print("修改模板" + denum, flush=True)
                                    sampleid = slide_dic['条目名称'].split('-')[1].replace('BP', '')
                                    sampleid_list = self.convert_string_to_list(sampleid)
                                    ######根据质粒个数来分别填写，需要加表格还是不加
                                    # print(sampleid_list, flush=True)
                                    if len(sampleid_list) == 1:
                                        iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                        driver.switch_to.frame(iframe_list[1])
                                        ####修改标题名
                                        if v['antibody'] == '':
                                            element = driver.find_element(By.XPATH,'// *[ @ id = "tinymce"] / ol / li[1] / p')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format('Western blot test'), element)
                                        else:
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / ol / li[1] / p')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format('Western blot test ('+v['antibody']+')'),element)
                                        element = driver.find_element(By.XPATH, ' // *[ @ id = "tinymce"] / ol / li[2] / p')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('Sample: ' + v['sample']),element)


                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/ol/li[3]/span/span/span')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('Exposure time: ' + v['Exposure_time']),element)
                                        driver.switch_to.default_content()
                                        ######插入图片部分Results & Discussion:

                                        driver.switch_to.frame(iframe_list[2])
                                        ####插入图片

                                        try:
                                            image = 'western blot_' + str(v['index']) + '.jpg'
                                            self.insertimage(image, '// *[ @ id = "tinymce"] / p[1] / img',
                                                             '// *[ @ id = "tinymce"] / p[1]')
                                        except:
                                            print('检查western blot图片', flush=True)
                                            pass
                                        ####修改结论
                                        try:
                                            element = driver.find_element(By.XPATH, ' // *[ @ id = "tinymce"] / p[2] / span[3]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format('The target protein of ' + v['batch'].split('BP')[1].split('-')[0]+' could be detected by western blot test.'), element)

                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / p[2] / span[5]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / p[2] / span[4]')
                                            driver.execute_script("arguments[0].remove()", element)
                                        except:
                                            print('sample中不含batch', flush=True)
                                            element = driver.find_element(By.XPATH, ' //*[@id="tinymce"]/p[2]/ span[3]')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format('The target protein could be detected by western blot test.'),element)
                                            element = driver.find_element(By.XPATH,'// *[ @ id = "tinymce"] / p[2] / span[5]')
                                            driver.execute_script("arguments[0].remove()", element)
                                            element = driver.find_element(By.XPATH,'// *[ @ id = "tinymce"] / p[2] / span[4]')
                                            driver.execute_script("arguments[0].remove()", element)
                                        driver.switch_to.default_content()
                                        self.saveoption(slide_dic['条目名称'].split('-')[1] + ' ' + slide_dic['条目名称'].split('-')[0])
                                        print(denum+'内容修改成功', flush=True)
                                    elif len(sampleid_list) > 1:
                                        iframe_list = driver.find_elements(By.XPATH, r'//iframe')
                                        driver.switch_to.frame(iframe_list[1])
                                        if v['antibody'] == '':
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / ol / li[1] / p')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format('Western blot test'),element)
                                        else:
                                            element = driver.find_element(By.XPATH, '// *[ @ id = "tinymce"] / ol / li[1] / p')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format('Western blot test (' + v['antibody'] + ')'), element)

                                        element = driver.find_element(By.XPATH, ' // *[ @ id = "tinymce"] / ol / li[2] / p')
                                        driver.execute_script("arguments[0].innerHTML = '{}';".format('Sample: ' + v['sample']),element)

                                        try:
                                            element = driver.find_element(By.XPATH,'//*[@id="tinymce"]/ol/li[3]/span/span/span')
                                            driver.execute_script("arguments[0].innerHTML = '{}';".format('Exposure time: ' + v['Exposure_time']), element)
                                        except:
                                            print('检查Exposure time', flush=True)
                                            pass
                                        driver.switch_to.default_content()
                                        ######插入图片部分Results & Discussion:
                                        iframe_list = driver.find_elements(By.XPATH, r'//iframe')

                                        driver.switch_to.frame(iframe_list[2])
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/p[1]/img')

                                        driver.execute_script("arguments[0].remove()", element)

                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/p[1]')
                                        ####换行符编码方式
                                        table_str = table_str.replace('\u000b', '\n')
                                        start_offset = 0
                                        end_offset = 0
                                        driver.execute_script("""
                                                                                                   var range_obj = document.createRange();
                                                                                                   range_obj.setStart(arguments[0], arguments[1]);
                                                                                                   range_obj.setEnd(arguments[0], arguments[2]);
                                                                                                   var fragment = range_obj.createContextualFragment(arguments[3]);
                                                                                                   range_obj.insertNode(fragment);
                                                                                               """,
                                                              element, start_offset, end_offset, table_str)

                                        # driver.execute_script(insert_script)

                                        # print('插入表格成功', flush=True)
                                        ####移动光标到最后一行
                                        body = driver.find_element(By.TAG_NAME, 'body')
                                        # 将光标移到文本末尾处
                                        body.send_keys(Keys.END)
                                        body.send_keys(Keys.ENTER)

                                        ####插入图片
                                        try:
                                            image = 'western blot_' + str(v['index']) + '.jpg'
                                            self.insertimage2(image, '// *[ @ id = "tinymce"] / p[3]')
                                            # print('插入图片成功', flush=True)
                                        except:
                                            print('检查western blot图片', flush=True)
                                            pass
                                        ####修改结论
                                        # print('修改结论', flush=True)
                                        body = driver.find_element(By.TAG_NAME, 'body')
                                        # 将光标移到文本末尾处
                                        body.send_keys(Keys.END)
                                        body.send_keys(Keys.ENTER)
                                        # print('移动光标到最后一行', flush=True)

                                        #####这一页的batch有多少个
                                        matches = re.findall(r'BP(\d+(?:\+\d+)?(?:\,\d+)?(?:\/\d+)?)', v['batch'])
                                        result = matches
                                        # print(result, flush=True)
                                        batch_list = []
                                        for i in result:
                                            i = i.replace('+', ',').replace('/', ',').replace('，', ',')
                                            parts = re.split(r',', i)
                                            batch_list.extend(parts)
                                        # print(batch_list, flush=True)
                                        #######看这一页的batch，如果一个那么直接写conclusion，否则不写
                                        if len(batch_list) == 1:
                                            # print('单个batch', flush=True)
                                            try:
                                                element = driver.find_element(By.XPATH, ' //*[@id="tinymce"]/p[4]')
                                                driver.execute_script("arguments[0].innerHTML = '{}';".format('Conclusions: The target protein of ' +v['batch'].split('BP')[1].split('-')[0] + ' could be detected by western blot test.'), element)
                                            except:
                                                element = driver.find_element(By.XPATH, ' //*[@id="tinymce"]/p[4]')
                                                driver.execute_script("arguments[0].innerHTML = '{}';".format('Conclusions: The target protein could be detected by western blot test. '),element)
                                                print('sample中不含batch', flush=True)
                                        else:
                                            try:
                                            # print('多个batch', flush=True)
                                                element = driver.find_element(By.XPATH, ' //*[@id="tinymce"]/p[4]')
                                                driver.execute_script("arguments[0].innerHTML = '{}';".format('Conclusions: The target protein could be detected by western blot test. '),element)
                                            except:
                                                pass
                                                # print('sample中不含batch', flush=True)
                                        #####删除之前的结论
                                        element = driver.find_element(By.XPATH, '//*[@id="tinymce"]/p[2]')
                                        driver.execute_script("arguments[0].remove()", element)
                                        driver.switch_to.default_content()
                                        self.saveoption(slide_dic['条目名称'].split('-')[1] + ' ' + slide_dic['条目名称'].split('-')[0])
                                        print(denum+'内容修改成功', flush=True)
                                    else:
                                        print('无法获取表单信息，请检查ppt第一页内容', flush=True)
                                        pass
                        except:
                            pass
                print("完成修改模板内容任务", flush=True)
                try:
                    destination_path = r'完成上传'
                    shutil.copy2(ppt_file, destination_path)
                    # 删除源文件
                    os.remove(ppt_file)
                    time.sleep(1)
                    os.remove(image_file_path)
                except:
                    pass

        else:
            try:
                destination_path = r'项目编号有问题'
                shutil.copy2(ppt_file, destination_path)
                # 删除源文件
                os.remove(ppt_file)

            except:
                pass


@Gooey(
    richtext_controls=True,  # 打开终端对颜色支持
    program_name="蛋白制备ELN报告上传",  # 程序名称
    clear_before_run=True,
    required_cols=2
)
def main():

    settings_msg = '蛋白制备ELN报告上传'
    parser = GooeyParser(description=settings_msg)  # 添加上方的应用信息
    subs = parser.add_subparsers(help='commands', dest='command')
    # 导入配置文件
    account,password, ppt_file, name= out.load_configuration()
    SingleSearchparser = subs.add_parser('蛋白制备ELN报告上传')
    SingleSearchparser.add_argument("script", choices=["QCs", "Western Bolt"], help="选择上传报告类型")
    SingleSearchparser.add_argument("account", metavar='ELN账号')
    SingleSearchparser.add_argument("password", metavar='ELN账号密码')
    SingleSearchparser.add_argument("name", metavar='ELN Author')
    SingleSearchparser.add_argument("ppt_files", metavar='ppt文件夹',widget='DirChooser',help="选择ppt文件夹")

    """
       程序功能实现
    """
    args = parser.parse_args()
    """界面获得的数据进行处理"""
    if args.command == "蛋白制备ELN报告上传":
        account = args.account.strip()
        password = args.password.strip()
        ppt_path = args.ppt_files.strip()
        name=args.name.strip()
        if not os.path.exists('content'):
            os.makedirs('content')
        if os.path.exists('content'):
            # print('删除content文件夹下的所有文件夹', flush=True)
            for root, dirs, files in os.walk('content', topdown=False):
                for file in files:
                    os.remove(os.path.join(root, file))
                for dir in dirs:
                    os.rmdir(os.path.join(root, dir))

    ppt_files = [file for file in os.listdir(ppt_path) if file.endswith('.pptx')]
    total_count = len(ppt_files)
    if not os.path.exists('完成上传'):
        os.makedirs('完成上传')
    if not os.path.exists('项目编号有问题'):
        os.makedirs('项目编号有问题')

    if args.script == "QCs":
    ####检查Chrome
        if out.check_broswer():
            make_print_to_file()
            current_time = datetime.now()
            print('-------------程序开始运行时间：' + str(current_time) + '-------------')
            for index, file in enumerate(ppt_files):
                ppt_file = os.path.join(ppt_path, file)
                ppt_file_name = file.split('.pptx')[0]
                # print(index, ppt_file_name,ppt_file)
                ppt_file_name=ppt_file_name.replace('.','-')
                ######获取存放文件的路径
                image_file = os.path.abspath('content')
                ######获取每一个pptx存放图片的路径
                image_file_path = os.path.join(image_file, ppt_file_name)
                if not os.path.exists(image_file_path):
                    os.makedirs(image_file_path)
                if os.path.exists(image_file_path):
                    for roots, dirs, files in os.walk(image_file_path):
                        for file in files:
                            os.remove(os.path.join(image_file_path, file))
                try:
                    delete_text_boxes_with_keyword(ppt_file, image_file_path)
                except:
                    print('将该份pptx文件保存为图片失败', flush=True)

                print('*********正在处理第{}/{}个文件*********'.format(index + 1, total_count))
                print(ppt_file)

                try_count = 0
                max_try_count = 3

                while try_count < max_try_count:
                    try:
                        eln_auto = out(account, password, ppt_file, image_file_path, name)

                        eln_auto.upload_ELN()
                        break
                    except Exception:
                        time.sleep(2)
                        try_count += 1
                        print(f'*********{ppt_file_name}+第{try_count}次上传失败*********', flush=True)
                        if try_count < max_try_count:
                            print(f'尝试进行第{try_count + 1} 次上传...')
                            time.sleep(1)
                if try_count == max_try_count:
                    print(f'*********{ppt_file_name}+上传最终失败*********', flush=True)
            print('*********所有内容上传完成,注意检查核对*********', flush=True)
            current_time = datetime.now()
            print('-------------程序结束运行时间：' + str(current_time) + '-------------', flush=True)
        else:
            # print('Chrome浏览器或驱动安装失败, 请手动安装重新运行程序, 或联系技术人员',flush=True)
            win32api.MessageBox(0, "Chrome浏览器或驱动安装失败, 请手动安装重新运行程序, 或联系技术人员", "提醒", win32con.MB_TOPMOST)
    if args.script == "Western Bolt":
        make_print_to_file()
        if out.check_broswer():
            # 下载
            current_time = datetime.now()
            print('-------------程序开始运行时间：' + str(current_time) + '-------------')
            for index, file in enumerate(ppt_files):
                ppt_file = os.path.join(ppt_path, file)
                ppt_file_name = file.split('.pptx')[0]
                # print(index, ppt_file_name,ppt_file)
                ppt_file_name = ppt_file_name.replace('.', '-')
                ######获取存放文件的路径
                image_file = os.path.abspath('content')
                ######获取每一个pptx存放图片的路径
                image_file_path = os.path.join(image_file, ppt_file_name)
                if not os.path.exists(image_file_path):
                    os.makedirs(image_file_path)
                if os.path.exists(image_file_path):
                    for roots, dirs, files in os.walk(image_file_path):
                        for file in files:
                            os.remove(os.path.join(image_file_path, file))

                # eln_auto.make_print_to_file()
                print('*********正在处理第{}/{}个文件*********'.format(index + 1, total_count))
                print(ppt_file)
                try_count = 0
                max_try_count = 3
                while try_count < max_try_count:
                    try:
                        eln_auto = out(account, password, ppt_file, image_file_path, name)
                        utils.save_pptx_as_png(image_file_path, ppt_file, overwrite_folder=True)
                        # 实例化对象
                        eln_auto.upload_ELN_wb()
                        break
                    except:
                        time.sleep(1)
                        print(f'*********{ppt_file_name}+第{try_count}次上传失败*********', flush=True)
                        try_count += 1
                        if try_count < max_try_count:
                            print(f'尝试进行第{try_count + 1} 次上传...')
                            time.sleep(1)  # Add a delay before retrying
                if try_count == max_try_count:
                    print(f'*********{ppt_file_name}+上传最终失败*********', flush=True)
            print('*********所有内容上传完成,注意检查核对*********', flush=True)
            current_time = datetime.now()
            print('-------------程序结束运行时间：' + str(current_time) + '-------------', flush=True)

        else:
            # print('Chrome浏览器或驱动安装失败, 请手动安装重新运行程序, 或联系技术人员',flush=True)
            win32api.MessageBox(0, "Chrome浏览器或驱动安装失败, 请手动安装重新运行程序, 或联系技术人员", "提醒", win32con.MB_TOPMOST)

if __name__ == '__main__':
    try:
        main()
        win32api.MessageBox(0, "任务完成", "提醒", win32con.MB_TOPMOST)
        input('\n')
    except Exception as e:
        traceback.print_exc()
        input("任务出错，联系计算结构平台")
        input("按任意键退出:")
